#!/usr/bin/env python3
"""
Create MutBench PhD Defense Presentation (.pptx)
Styled after the reference KNU presentation.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Constants ──────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Colors
DARK_BAR   = RGBColor(44, 62, 80)
WHITE      = RGBColor(255, 255, 255)
DARK_GRAY  = RGBColor(51, 51, 51)
MED_GRAY   = RGBColor(100, 100, 100)
LIGHT_BG   = RGBColor(245, 245, 245)
BLUE       = RGBColor(41, 128, 185)
ORANGE     = RGBColor(230, 126, 34)
GREEN      = RGBColor(39, 174, 96)
RED        = RGBColor(231, 76, 60)
PURPLE     = RGBColor(142, 68, 173)
LIGHT_BLUE = RGBColor(214, 234, 248)
LIGHT_GREEN= RGBColor(212, 239, 223)
LIGHT_RED  = RGBColor(250, 219, 216)
LIGHT_ORANGE = RGBColor(253, 235, 208)

# Layout
TOP_BAR_H  = Inches(1.0)
BOT_BAR_H  = Inches(0.3)
BOT_BAR_Y  = SLIDE_H - BOT_BAR_H
CONTENT_TOP = Inches(1.2)
CONTENT_LEFT = Inches(0.6)
CONTENT_W   = Inches(12.1)
CONTENT_H   = Inches(5.6)
FOOTER_TEXT = "Kyungpook National University  |  Dept. of Computer Science  |  Hwijun Kwon"

# Logo path
LOGO_PATH = "/proj/paper/docs/presentation/knu_symbol_white.png"
HAS_LOGO = os.path.exists(LOGO_PATH)

# ── Helper functions ───────────────────────────────────────────────────

def add_top_bar(slide, title_text):
    """Add the dark top bar with title text."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, TOP_BAR_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BAR
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.color.rgb = WHITE
    p.font.size = Pt(28)
    p.font.bold = True

    if HAS_LOGO:
        slide.shapes.add_picture(
            LOGO_PATH, Inches(12.1), Inches(0.12), height=Inches(0.7)
        )


def add_bottom_bar(slide, page_num):
    """Add the thin bottom bar with footer and page number."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), BOT_BAR_Y, SLIDE_W, BOT_BAR_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BAR
    bar.line.fill.background()

    tf = bar.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = FOOTER_TEXT
    p.font.color.rgb = WHITE
    p.font.size = Pt(10)

    # Page number on the right
    pn = slide.shapes.add_textbox(
        Inches(12.4), BOT_BAR_Y, Inches(0.8), BOT_BAR_H
    )
    tf2 = pn.text_frame
    tf2.margin_top = Inches(0.02)
    p2 = tf2.paragraphs[0]
    p2.text = str(page_num)
    p2.font.color.rgb = WHITE
    p2.font.size = Pt(10)
    p2.alignment = PP_ALIGN.RIGHT


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT,
                line_spacing=1.3):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.line_spacing = Pt(font_size * line_spacing)
    return tf


def add_bullets(slide, left, top, width, height, items, font_size=18,
                color=DARK_GRAY, bullet_char="\u2022"):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.line_spacing = Pt(font_size * 1.5)
        p.space_after = Pt(6)
    return tf


def add_colored_box(slide, left, top, width, height, title, body,
                    fill_color, title_color=WHITE, body_color=DARK_GRAY,
                    title_size=16, body_size=14, border_color=None):
    """Add a colored box with title and body text."""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    if border_color:
        box.line.color.rgb = border_color
        box.line.width = Pt(2)
    else:
        box.line.fill.background()

    # Title textbox
    ttl = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.1),
        width - Inches(0.3), Inches(0.4)
    )
    tf = ttl.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = title_color

    # Body textbox
    bdy = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.5),
        width - Inches(0.3), height - Inches(0.6)
    )
    tf2 = bdy.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = body
    p2.font.size = Pt(body_size)
    p2.font.color.rgb = body_color
    p2.line_spacing = Pt(body_size * 1.4)
    return box


def add_simple_table(slide, left, top, rows, cols, data, col_widths=None,
                     header_color=DARK_BAR, font_size=14):
    """Add a table with header row styling."""
    table_shape = slide.shapes.add_table(rows, cols, left, top,
                                          Inches(11), Inches(0.4 * rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                if r == 0:
                    paragraph.font.color.rgb = WHITE
                    paragraph.font.bold = True
                else:
                    paragraph.font.color.rgb = DARK_GRAY

            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else LIGHT_BG

    return table


def add_arrow_text(slide, left, top, text, color=BLUE, size=24):
    """Add arrow indicator text."""
    add_textbox(slide, left, top, Inches(12), Inches(0.5),
                f"\u2192 {text}", font_size=size, bold=True, color=color)


def add_big_number(slide, left, top, number, label, color=BLUE, num_size=48):
    """Add a large number with label below."""
    add_textbox(slide, left, top, Inches(3), Inches(1),
                str(number), font_size=num_size, bold=True, color=color,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + Inches(0.9), Inches(3), Inches(0.6),
                label, font_size=14, color=DARK_GRAY,
                alignment=PP_ALIGN.CENTER)


# ── Slide builders ─────────────────────────────────────────────────────

def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Full dark background for title slide
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BAR
    bg.line.fill.background()

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(2), Inches(2.7), Inches(9.3), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()

    # Title
    add_textbox(slide, Inches(2), Inches(1.0), Inches(9.3), Inches(1.6),
                "MutBench: Systematic Benchmarking\nfor Viral Mutation Hotspot Detection",
                font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

    # Subtitle
    add_textbox(slide, Inches(2), Inches(2.9), Inches(9.3), Inches(0.5),
                "PhD Dissertation Defense",
                font_size=22, color=RGBColor(189, 195, 199), alignment=PP_ALIGN.LEFT)

    # Author
    add_textbox(slide, Inches(2), Inches(3.7), Inches(9.3), Inches(0.5),
                "Hwijun Kwon  |  Advisor: Prof. Inuk Jung",
                font_size=20, color=WHITE, alignment=PP_ALIGN.LEFT)

    # Institution
    add_textbox(slide, Inches(2), Inches(4.3), Inches(9.3), Inches(0.5),
                "School of Computer Science, Kyungpook National University",
                font_size=16, color=RGBColor(189, 195, 199), alignment=PP_ALIGN.LEFT)

    # Date
    add_textbox(slide, Inches(2), Inches(4.9), Inches(9.3), Inches(0.5),
                "August 2026",
                font_size=16, color=RGBColor(189, 195, 199), alignment=PP_ALIGN.LEFT)

    if HAS_LOGO:
        slide.shapes.add_picture(
            LOGO_PATH, Inches(11.3), Inches(0.5), height=Inches(0.9)
        )


def slide_02_hotspots(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "What are Mutation Hotspots?")
    add_bottom_bar(slide, 2)

    add_bullets(slide, CONTENT_LEFT, Inches(1.4), CONTENT_W, Inches(3.5), [
        "Genomic positions with statistically elevated mutation rates",
        "Associated with transmissibility, immune evasion, drug resistance",
        "Critical for vaccine design, antiviral development, genomic surveillance",
    ], font_size=20)

    # Example box
    add_colored_box(slide, Inches(1.0), Inches(4.5), Inches(11.3), Inches(1.5),
                    "SARS-CoV-2 Spike Examples",
                    "N501Y (receptor binding)    E484K (immune evasion)    L452R (transmissibility)",
                    LIGHT_BLUE, title_color=BLUE, body_color=DARK_GRAY,
                    border_color=BLUE)


def slide_03_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "The Problem \u2014 No Benchmark Exists")
    add_bottom_bar(slide, 3)

    box_w = Inches(3.5)
    box_h = Inches(3.0)
    gap = Inches(0.4)
    start_x = Inches(0.8)

    boxes = [
        ("No Standardized\nEvaluation", "Existing studies use ad-hoc\nmetrics; results are\nnot comparable across\nstudies", RED),
        ("No Independent\nGround Truth", "Self-referencing validation;\nliterature-derived labels\nare circular", ORANGE),
        ("No Statistical Evidence\non Pathogen-Dependence", "Methods are assumed\nuniversal; no formal test of\ninteraction effects", BLUE),
    ]

    for i, (title, body, color) in enumerate(boxes):
        x = start_x + i * (box_w + gap)
        add_colored_box(slide, x, Inches(1.5), box_w, box_h,
                        title, body, color, WHITE, WHITE, 18, 15)

    add_arrow_text(slide, Inches(1.0), Inches(5.2),
                   "Which method should we use for which pathogen?", DARK_BAR, 22)


def slide_04_contributions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "Research Contributions")
    add_bottom_bar(slide, 4)

    # Left column
    add_colored_box(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(4.5),
                    "Contribution 1: MutBench Framework",
                    "\u2022 3-layer ground truth\n   (convergent + purifying + DMS)\n\n"
                    "\u2022 Hotspot-score metric\n   (recall + precision + stability)\n\n"
                    "\u2022 2-stage experimental design\n   (depth-first + breadth-first)",
                    LIGHT_BLUE, BLUE, DARK_GRAY, 20, 16, BLUE)

    # Right column
    add_colored_box(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(4.5),
                    "Contribution 2: Statistical Demonstration",
                    "\u2022 \u03c9\u00b2 = 0.285 interaction effect\n   (largest variance source)\n\n"
                    "\u2022 LOPO: 0/9 generalization\n   (no method transfers)\n\n"
                    "\u2022 Friedman p = 0.0015\n   (significant rank differences)",
                    LIGHT_ORANGE, ORANGE, DARK_GRAY, 20, 16, ORANGE)


def slide_05_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "MutBench Pipeline Overview")
    add_bottom_bar(slide, 5)

    # Pipeline boxes
    stages = [
        ("MSA\nSequences", BLUE),
        ("Scoring\n(9 methods)", GREEN),
        ("Detection\n(41 methods)", ORANGE),
        ("Evaluation\n(MCC)", RED),
        ("Statistical\nAnalysis", PURPLE),
    ]

    box_w = Inches(2.0)
    box_h = Inches(1.4)
    y = Inches(2.0)
    start_x = Inches(0.6)
    gap = Inches(0.45)

    for i, (label, color) in enumerate(stages):
        x = start_x + i * (box_w + gap)
        add_colored_box(slide, x, y, box_w, box_h, label, "",
                        color, WHITE, WHITE, 17, 12)

        # Arrow between boxes
        if i < len(stages) - 1:
            ax = x + box_w + Inches(0.05)
            add_textbox(slide, ax, y + Inches(0.4), Inches(0.35), Inches(0.5),
                        "\u25B6", font_size=20, color=DARK_GRAY,
                        alignment=PP_ALIGN.CENTER)

    # Ground truth box feeding into evaluation
    add_colored_box(slide, Inches(4.0), Inches(4.0), Inches(3.0), Inches(1.2),
                    "3-Layer Ground Truth", "Convergent + Purifying + DMS",
                    LIGHT_GREEN, GREEN, DARK_GRAY, 16, 13, GREEN)

    # Arrow up from GT to Evaluation
    add_textbox(slide, Inches(5.1), Inches(3.5), Inches(0.8), Inches(0.5),
                "\u25B2", font_size=20, color=GREEN, alignment=PP_ALIGN.CENTER)

    # ANOVA / Friedman / LOPO labels
    add_textbox(slide, Inches(10.0), Inches(3.8), Inches(3.0), Inches(2.0),
                "ANOVA (\u03c9\u00b2)\nFriedman (\u03c7\u00b2)\nLOPO (0/9)",
                font_size=16, color=PURPLE)


def slide_06_ground_truth(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "3-Layer Ground Truth")
    add_bottom_bar(slide, 6)

    layers = [
        ("Layer A: Convergent Evolution", "Positive labels (TP)\nMultiple independent lineages\nacquire same mutation",
         GREEN, LIGHT_GREEN),
        ("Layer B: Purifying Selection", "Negative labels (TN)\ndN/dS \u226a 1 positions\nunder strong constraint",
         RED, LIGHT_RED),
        ("Layer C: DMS Experimental", "Validation layer\nDeep mutational scanning\nfitness measurements",
         BLUE, LIGHT_BLUE),
    ]

    for i, (title, body, color, fill) in enumerate(layers):
        x = Inches(0.6) + i * Inches(4.2)
        add_colored_box(slide, x, Inches(1.5), Inches(3.8), Inches(3.5),
                        title, body, fill, color, DARK_GRAY, 17, 15, color)

    # Independence note
    add_colored_box(slide, Inches(2.5), Inches(5.5), Inches(8.3), Inches(1.0),
                    "", "Jaccard(DMS, convergent) = 0.006  \u2192  Layers are independent",
                    LIGHT_BG, DARK_GRAY, DARK_BAR, 14, 18, MED_GRAY)


def slide_07_hotspot_score(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "Hotspot-score Metric")
    add_bottom_bar(slide, 7)

    add_textbox(slide, Inches(1.5), Inches(1.4), Inches(10), Inches(0.8),
                "hotspot-score  =  ( recall  +  precision  +  stability )  /  3",
                font_size=28, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)

    # Comparison table
    data = [
        ["Method", "Recall", "Precision", "Stability", "F1", "HS"],
        ["HDBSCAN-Auto", "0.944", "0.079", "0.788", "0.146", "0.604"],
        ["Hybrid", "0.659", "0.968", "0.706", "0.784", "0.778"],
    ]
    add_simple_table(slide, Inches(1.0), Inches(2.8), 3, 6, data,
                     col_widths=[2.5, 1.7, 1.7, 1.7, 1.7, 1.7])

    add_arrow_text(slide, Inches(1.0), Inches(4.5),
                   "F1 is misleading when recall/precision are imbalanced; "
                   "hotspot-score captures all three dimensions",
                   DARK_BAR, 17)


def slide_08_two_stage(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "Two-Stage Experimental Design")
    add_bottom_bar(slide, 8)

    add_colored_box(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(3.5),
                    "Stage 1: SARS-CoV-2 Depth-First",
                    "5 detection variants + 2 baselines\n"
                    "Single pathogen, exhaustive comparison\n"
                    "Result: MutClust-Hybrid wins\n\n"
                    "HS = 0.778",
                    BLUE, WHITE, WHITE, 20, 16)

    add_colored_box(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(3.5),
                    "Stage 2: 9-Pathogen Breadth-First",
                    "9 scoring \u00d7 41 detection methods\n"
                    "= 3,321 total evaluations\n"
                    "6 RNA virus families\n\n"
                    "Result: No universal winner",
                    ORANGE, WHITE, WHITE, 20, 16)

    # Arrow between
    add_textbox(slide, Inches(6.1), Inches(2.8), Inches(1.1), Inches(0.5),
                "\u25B6", font_size=36, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    add_arrow_text(slide, Inches(1.0), Inches(5.5),
                   "Stage 1 shows Hybrid wins. But does this generalize?",
                   DARK_BAR, 20)


def slide_09_stage1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "Stage 1 \u2014 MutClust-Hybrid Wins on SARS-CoV-2")
    add_bottom_bar(slide, 9)

    data = [
        ["Method", "Precision", "Recall", "Stability", "Hotspot-score"],
        ["MutClust-Hybrid", "0.968", "0.659", "0.706", "0.778"],
        ["HDBSCAN-Auto", "0.079", "0.944", "0.788", "0.604"],
        ["Two-stage", "0.968", "0.643", "0.740", "0.784"],
    ]
    add_simple_table(slide, Inches(1.0), Inches(1.8), 4, 5, data,
                     col_widths=[3.0, 2.0, 2.0, 2.0, 2.0])

    add_bullets(slide, Inches(1.0), Inches(4.2), Inches(11), Inches(2.5), [
        "Hybrid balances all three metric components",
        "HDBSCAN achieves highest recall but worst precision (over-detection)",
        "Two-stage shows comparable but slightly lower recall",
    ], font_size=17)


def slide_10_real_gap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "Synthetic vs Real Data Gap")
    add_bottom_bar(slide, 10)

    # Two big numbers
    add_big_number(slide, Inches(1.5), Inches(1.8), "0.778",
                   "Synthetic hotspot-score", BLUE, 52)

    add_textbox(slide, Inches(5.0), Inches(2.0), Inches(1.5), Inches(0.8),
                "\u2192", font_size=52, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    add_big_number(slide, Inches(6.5), Inches(1.8), "0.123",
                   "Real GISAID F1", RED, 52)

    # Explanation
    add_colored_box(slide, Inches(1.0), Inches(4.0), Inches(11.3), Inches(2.2),
                    "Why the gap?",
                    "Enrichment ratios are preserved (6.9\u20138.0\u00d7)\n"
                    "Gap = calibration factor from simulation assumptions, not method invalidation\n"
                    "Relative rankings remain consistent",
                    LIGHT_BG, DARK_BAR, DARK_GRAY, 18, 16, MED_GRAY)


def slide_11_h3n2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "H3N2 Reverses the Rankings")
    add_bottom_bar(slide, 11)

    add_colored_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.2),
                    "SARS-CoV-2",
                    "Best: MutClust-Hybrid\nHotspot-score: 0.778\n\n"
                    "Fast-evolving, high diversity\nEntropy highly informative",
                    BLUE, WHITE, WHITE, 22, 16)

    add_colored_box(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(3.2),
                    "H3N2 Influenza",
                    "Best: Fixed threshold\nHotspot-score: 0.661\n\n"
                    "Antigenic drift, segment reassortment\nDifferent mutation landscape",
                    ORANGE, WHITE, WHITE, 22, 16)

    add_textbox(slide, Inches(6.3), Inches(2.6), Inches(0.7), Inches(0.5),
                "\u2260", font_size=48, bold=True, color=RED, alignment=PP_ALIGN.CENTER)

    add_arrow_text(slide, Inches(1.0), Inches(5.3),
                   "The best method depends on the pathogen",
                   RED, 24)


def slide_12_scale(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "Stage 2 \u2014 Scale of Evaluation")
    add_bottom_bar(slide, 12)

    # Large numbers
    nums = [
        ("9", "Pathogens", BLUE),
        ("\u00d7", "", DARK_GRAY),
        ("9", "Scoring Methods", GREEN),
        ("\u00d7", "", DARK_GRAY),
        ("41", "Detection Methods", ORANGE),
    ]

    for i, (n, label, color) in enumerate(nums):
        x = Inches(0.5) + i * Inches(2.5)
        if n in ("\u00d7",):
            add_textbox(slide, x, Inches(2.0), Inches(2.0), Inches(1.0),
                        n, font_size=48, color=color, alignment=PP_ALIGN.CENTER)
        else:
            add_big_number(slide, x, Inches(1.8), n, label, color, 60)

    # Result line
    add_textbox(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.8),
                "= 3,321 evaluations  (2,544 valid for ANOVA)",
                font_size=28, bold=True, color=DARK_BAR, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.5),
                "6 major RNA virus families",
                font_size=20, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

    # Pathogen list
    add_textbox(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(0.8),
                "SARS-CoV-2  |  H3N2  |  HIV-1  |  MERS  |  RSV  |  "
                "Dengue  |  HCV  |  Norovirus  |  Influenza B",
                font_size=15, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


def slide_13_key_finding(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "KEY FINDING \u2014 4 Lines of Evidence")
    add_bottom_bar(slide, 13)

    # Highlight bar
    highlight = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.06)
    )
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = RED
    highlight.line.fill.background()

    findings = [
        ("\u03c9\u00b2 = 0.285", "Interaction is the LARGEST variance source", BLUE),
        ("9/9 unique", "Every pathogen has a different best combination", GREEN),
        ("Friedman p = 0.0015", "Significant rank differences (\u03c7\u00b2 = 42.44)", ORANGE),
        ("LOPO: 0/9", "No method generalizes (4.9% match rate)", RED),
    ]

    for i, (num, desc, color) in enumerate(findings):
        y = Inches(1.6) + i * Inches(1.3)

        add_textbox(slide, Inches(1.0), y, Inches(4.0), Inches(0.7),
                    num, font_size=36, bold=True, color=color)

        add_textbox(slide, Inches(5.5), y + Inches(0.08), Inches(7.0), Inches(0.7),
                    desc, font_size=20, color=DARK_GRAY)

    # Conclusion
    conclusion = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(6.0),
        Inches(10.3), Inches(0.7)
    )
    conclusion.fill.solid()
    conclusion.fill.fore_color.rgb = RED
    conclusion.line.fill.background()

    add_textbox(slide, Inches(1.5), Inches(6.05), Inches(10.3), Inches(0.6),
                "No single best method exists",
                font_size=26, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)


def slide_14_esm2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "ESM-2 Cross-Paradigm Validation")
    add_bottom_bar(slide, 14)

    data = [
        ["Pathogen", "ESM-LLR Rank", "ESM-LLR Score", "Existing Best"],
        ["SARS-CoV-2", "1st", "0.168", "0.118"],
        ["MERS", "1st", "0.143", "0.088"],
        ["Norovirus", "Last", "\u2014", "0.438"],
        ["HIV-1", "Last", "\u2014", "0.182"],
    ]
    add_simple_table(slide, Inches(1.0), Inches(1.8), 5, 4, data,
                     col_widths=[3.0, 2.5, 2.5, 3.0])

    add_arrow_text(slide, Inches(1.0), Inches(5.0),
                   "No single scoring paradigm dominates across all pathogens",
                   DARK_BAR, 20)

    add_textbox(slide, Inches(1.0), Inches(5.8), Inches(11), Inches(0.8),
                "ESM-2 (protein language model) excels on Betacoronaviruses "
                "but fails on high-diversity pathogens",
                font_size=15, color=MED_GRAY)


def slide_15_window(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "Position-Exact is Conservative")
    add_bottom_bar(slide, 15)

    # Three bars showing improvement
    bars = [
        ("Exact match", 0.289, "0.289", BLUE),
        ("\u00b15 positions", 0.638, "0.638 (+121%)", GREEN),
        ("\u00b110 positions", 0.712, "0.712 (+146%)", ORANGE),
    ]

    for i, (label, frac, val_text, color) in enumerate(bars):
        y = Inches(1.8) + i * Inches(1.3)

        add_textbox(slide, Inches(0.8), y, Inches(2.5), Inches(0.5),
                    label, font_size=18, bold=True, color=DARK_GRAY)

        # Bar background
        bar_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(3.5), y + Inches(0.05),
            Inches(8.5), Inches(0.6)
        )
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = LIGHT_BG
        bar_bg.line.fill.background()

        # Bar fill
        bar_fill = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(3.5), y + Inches(0.05),
            Inches(8.5 * frac), Inches(0.6)
        )
        bar_fill.fill.solid()
        bar_fill.fill.fore_color.rgb = color
        bar_fill.line.fill.background()

        # Value text
        add_textbox(slide, Inches(3.5) + Inches(8.5 * frac) + Inches(0.1),
                    y + Inches(0.05), Inches(2.5), Inches(0.6),
                    val_text, font_size=16, bold=True, color=color)

    add_arrow_text(slide, Inches(1.0), Inches(5.8),
                   "Detected hotspots are spatially proximate to ground truth",
                   DARK_BAR, 18)


def slide_16_pahd(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "PAHD \u2014 Proof of Concept")
    add_bottom_bar(slide, 16)

    add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.6),
                "Problem: 369 combinations, only n=9 pathogens "
                "\u2192 exact prediction impossible",
                font_size=18, color=DARK_GRAY)

    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(0.5),
                "Solution: Predict at FAMILY level instead of exact combination",
                font_size=18, bold=True, color=BLUE)

    data = [
        ["Prediction Level", "Oracle Recovery", "Improvement"],
        ["Exact combination", "1.7%", "baseline"],
        ["Detection family", "48.3%", "+28\u00d7"],
        ["Scoring family", "57.7%", "+34\u00d7"],
    ]
    add_simple_table(slide, Inches(1.5), Inches(3.0), 4, 3, data,
                     col_widths=[4.0, 3.0, 3.0])

    add_arrow_text(slide, Inches(1.0), Inches(5.5),
                   "Family-level prediction makes adaptive selection feasible",
                   GREEN, 18)


def slide_17_phylo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "Phylogenetic Correction")
    add_bottom_bar(slide, 17)

    add_bullets(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.5), [
        "TreeTime on 500 real SARS-CoV-2 sequences",
        "Result: MCC = \u22120.018 (temporal mismatch)",
        "D614G: homoplasy = 0 (founder effect confirmed)",
    ], font_size=19)

    add_colored_box(slide, Inches(1.0), Inches(3.8), Inches(11.3), Inches(2.5),
                    "Key Insight",
                    "Phylogenetic correction requires temporally aligned data\n\n"
                    "Current public databases (GISAID) have sampling bias:\n"
                    "  - Over-representation of certain lineages\n"
                    "  - Temporal gaps in coverage\n"
                    "  - Geographic sampling heterogeneity",
                    LIGHT_ORANGE, ORANGE, DARK_GRAY, 18, 16, ORANGE)


def slide_18_limitations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "Limitations & Roadmap")
    add_bottom_bar(slide, 18)

    data = [
        ["Priority", "Direction", "Impact"],
        ["1", "Phylogenetic correction", "Founder-effect resolution"],
        ["2", "PAHD implementation", "Adaptive method selection"],
        ["3", "DNA virus extension", "Broader applicability"],
        ["4", "Real-time integration", "Public health utility"],
    ]
    add_simple_table(slide, Inches(0.8), Inches(1.8), 5, 3, data,
                     col_widths=[1.5, 4.5, 5.0])

    add_textbox(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.5),
                "MutBench provides the infrastructure for systematic evaluation.\n"
                "Each limitation is a concrete, addressable research direction.",
                font_size=16, color=MED_GRAY)


def slide_19_conclusions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "Conclusions")
    add_bottom_bar(slide, 19)

    # Two contribution boxes
    add_colored_box(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(3.0),
                    "1. MutBench Framework",
                    "First systematic benchmark for\nviral hotspot detection\n\n"
                    "3-layer ground truth\nHotspot-score metric\n2-stage design (3,321 evaluations)",
                    BLUE, WHITE, WHITE, 20, 15)

    add_colored_box(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(3.0),
                    "2. Statistical Demonstration",
                    "No universal best method exists\n\n"
                    "\u03c9\u00b2 = 0.285 (interaction dominates)\n"
                    "LOPO: 0/9 generalization\n"
                    "Friedman p = 0.0015",
                    GREEN, WHITE, WHITE, 20, 15)

    # Broader impact
    add_colored_box(slide, Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.5),
                    "Broader Impact",
                    "Method-pathogen interaction dominates variance \u2192 "
                    "systematic benchmarking is needed across all genomic analysis domains",
                    LIGHT_BG, DARK_BAR, DARK_GRAY, 18, 16, MED_GRAY)


def slide_20_thank_you(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Full dark background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), SLIDE_W, SLIDE_H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BAR
    bg.line.fill.background()

    add_textbox(slide, Inches(2), Inches(2.0), Inches(9.3), Inches(1.0),
                "Thank You",
                font_size=52, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(2), Inches(3.5), Inches(9.3), Inches(0.6),
                "Questions & Discussion",
                font_size=28, color=RGBColor(189, 195, 199), alignment=PP_ALIGN.CENTER)

    # Line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.5), Inches(4.3), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()

    add_textbox(slide, Inches(2), Inches(5.0), Inches(9.3), Inches(0.5),
                "github.com/hwijun/mutbench",
                font_size=18, color=BLUE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(2), Inches(5.7), Inches(9.3), Inches(0.5),
                "Hwijun Kwon  |  Kyungpook National University",
                font_size=16, color=RGBColor(189, 195, 199), alignment=PP_ALIGN.CENTER)


# ── Backup Slides ──────────────────────────────────────────────────────

def slide_b01_anova(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "BACKUP: ANOVA Diagnostics")
    add_bottom_bar(slide, "B1")

    add_bullets(slide, CONTENT_LEFT, Inches(1.5), CONTENT_W, Inches(5), [
        "3-way ANOVA: scoring \u00d7 detection \u00d7 pathogen",
        "\u03c9\u00b2 (omega-squared) used instead of \u03b7\u00b2 to avoid upward bias",
        "\u03c9\u00b2_detection = 0.350 (largest main effect)",
        "\u03c9\u00b2_scoring = 0.083",
        "\u03c9\u00b2_pathogen = 0.008",
        "\u03c9\u00b2_scoring\u00d7pathogen = 0.285 (interaction)",
        "Residuals checked: Q-Q plot, homoscedasticity, independence",
        "Levene test: p < 0.001 (heteroscedastic \u2192 Welch correction applied)",
    ], font_size=17)


def slide_b02_friedman(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "BACKUP: Friedman Power Analysis")
    add_bottom_bar(slide, "B2")

    add_bullets(slide, CONTENT_LEFT, Inches(1.5), CONTENT_W, Inches(5), [
        "Friedman test: non-parametric rank-based alternative to repeated-measures ANOVA",
        "Top 20 combinations ranked across 9 pathogens",
        "\u03c7\u00b2 = 42.44, df = 19, p = 0.0015",
        "Post-hoc Nemenyi: significant pairwise differences detected",
        "Power analysis: n=9 pathogens sufficient for W = 0.234 effect size",
        "Monte Carlo simulation: 10,000 iterations confirm stability",
    ], font_size=17)


def slide_b03_pahd(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "BACKUP: PAHD v1 vs v2 Comparison")
    add_bottom_bar(slide, "B3")

    data = [
        ["Feature", "PAHD v1", "PAHD v2 (Current)"],
        ["Prediction level", "Exact combo", "Family level"],
        ["Features", "Entropy stats", "Pathogen profile (8-dim)"],
        ["Oracle recovery", "1.7%", "48.3% (family)"],
        ["Required data", "Full MSA", "Summary statistics"],
        ["Status", "Failed", "Proof-of-concept"],
    ]
    add_simple_table(slide, Inches(0.8), Inches(1.8), 6, 3, data,
                     col_widths=[3.0, 4.0, 4.0])


def slide_b04_dms(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "BACKUP: DMS Threshold Sensitivity")
    add_bottom_bar(slide, "B4")

    add_bullets(slide, CONTENT_LEFT, Inches(1.5), CONTENT_W, Inches(5), [
        "DMS fitness threshold varied: 0.3, 0.4, 0.5, 0.6, 0.7",
        "Primary threshold: 0.5 (median fitness)",
        "Result: rankings stable across thresholds (\u03c4 > 0.85)",
        "Jaccard overlap between threshold sets: 0.72\u20130.89",
        "Convergent evolution labels unaffected (independent layer)",
        "Conclusion: results robust to DMS threshold choice",
    ], font_size=17)


def slide_b05_bug(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "BACKUP: MutClust Bug Description")
    add_bottom_bar(slide, "B5")

    add_bullets(slide, CONTENT_LEFT, Inches(1.5), CONTENT_W, Inches(5), [
        "Original MutClust: fixed p-value threshold (0.05)",
        "Bug: multiple testing correction not applied to clustering step",
        "Effect: inflated false positive rate in hotspot calls",
        "Fix: Bonferroni correction + adaptive threshold",
        "Hybrid approach: combines fixed + adaptive for robustness",
        "This bug motivated the need for systematic benchmarking",
    ], font_size=17)


def slide_b06_weights(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "BACKUP: Weight Sensitivity (171 combinations)")
    add_bottom_bar(slide, "B6")

    add_bullets(slide, CONTENT_LEFT, Inches(1.5), CONTENT_W, Inches(5), [
        "Hotspot-score weights: (w_recall, w_precision, w_stability)",
        "171 weight combinations tested (step = 0.1, sum = 1.0)",
        "Default: equal weights (1/3, 1/3, 1/3)",
        "Result: Hybrid robust across 89% of weight combinations",
        "Only extreme weights (w_stability > 0.7) change winner",
        "Equal weighting is justified: no prior reason to favor any component",
    ], font_size=17)


def slide_b07_pathogens(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "BACKUP: 9 Pathogens Coverage")
    add_bottom_bar(slide, "B7")

    data = [
        ["Pathogen", "Family", "Genome", "Sequences", "Known Hotspots"],
        ["SARS-CoV-2", "Coronaviridae", "30 kb", ">15M", "S: 501, 484, 452"],
        ["MERS-CoV", "Coronaviridae", "30 kb", "~800", "S: limited"],
        ["H3N2", "Orthomyxoviridae", "14 kb (seg)", "~50K", "HA: antigenic sites"],
        ["Influenza B", "Orthomyxoviridae", "14 kb (seg)", "~30K", "HA: limited"],
        ["HIV-1", "Retroviridae", "10 kb", ">500K", "Env: V1-V5 loops"],
        ["RSV", "Pneumoviridae", "15 kb", "~10K", "F: site \u00d8"],
        ["Dengue", "Flaviviridae", "11 kb", "~20K", "E: DIII"],
        ["HCV", "Flaviviridae", "10 kb", "~30K", "E2: HVR1"],
        ["Norovirus", "Caliciviridae", "7.5 kb", "~5K", "VP1: P2 domain"],
    ]
    add_simple_table(slide, Inches(0.3), Inches(1.5), 10, 5, data,
                     col_widths=[2.0, 2.5, 2.0, 2.0, 4.0], font_size=12)


def slide_b08_treetime(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "BACKUP: TreeTime Pipeline")
    add_bottom_bar(slide, "B8")

    add_bullets(slide, CONTENT_LEFT, Inches(1.5), CONTENT_W, Inches(5), [
        "Input: 500 SARS-CoV-2 sequences with collection dates",
        "Step 1: Multiple sequence alignment (MAFFT)",
        "Step 2: Maximum likelihood tree (IQ-TREE2)",
        "Step 3: Molecular clock fitting (TreeTime)",
        "Step 4: Ancestral state reconstruction",
        "Step 5: Homoplasy detection per site",
        "Result: D614G homoplasy = 0 (single origin, founder effect)",
        "Challenge: GISAID sampling bias distorts temporal signal",
    ], font_size=17)


def slide_b09_transfer(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "BACKUP: MOSD \u2192 MutBench Transfer")
    add_bottom_bar(slide, "B9")

    add_bullets(slide, CONTENT_LEFT, Inches(1.5), CONTENT_W, Inches(5), [
        "MOSD (Multi-Omics Subtype Discovery): benchmarking framework for multi-omics integration",
        "Key lesson: single-method evaluation insufficient across cancer types",
        "Transferred concepts to MutBench:",
        "  - Systematic method comparison framework",
        "  - Multiple evaluation metrics (not just accuracy)",
        "  - Cross-dataset generalization testing",
        "MutClust: mutation clustering algorithm, predecessor to MutBench scoring/detection",
        "MutBench = synthesis of MOSD methodology + MutClust domain knowledge",
    ], font_size=16)


def slide_b10_comparison(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, "BACKUP: Approach Comparison Table")
    add_bottom_bar(slide, "B10")

    data = [
        ["Aspect", "Previous Studies", "MutBench"],
        ["Ground truth", "Literature-derived", "3-layer (convergent+purifying+DMS)"],
        ["Evaluation", "Overlap counting", "Hotspot-score (3 components)"],
        ["Pathogens", "1 (usually SARS-CoV-2)", "9 pathogens, 6 families"],
        ["Methods", "1\u20133 methods", "369 combinations (9\u00d741)"],
        ["Statistics", "None / p-value only", "ANOVA + Friedman + LOPO + bootstrap"],
        ["Interaction", "Not tested", "\u03c9\u00b2 = 0.285 (demonstrated)"],
        ["Reproducibility", "Ad hoc scripts", "Open-source framework"],
    ]
    add_simple_table(slide, Inches(0.5), Inches(1.5), 8, 3, data,
                     col_widths=[2.5, 4.5, 5.5], font_size=13)


# ── Main ───────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Main slides
    slide_01_title(prs)
    slide_02_hotspots(prs)
    slide_03_problem(prs)
    slide_04_contributions(prs)
    slide_05_pipeline(prs)
    slide_06_ground_truth(prs)
    slide_07_hotspot_score(prs)
    slide_08_two_stage(prs)
    slide_09_stage1(prs)
    slide_10_real_gap(prs)
    slide_11_h3n2(prs)
    slide_12_scale(prs)
    slide_13_key_finding(prs)
    slide_14_esm2(prs)
    slide_15_window(prs)
    slide_16_pahd(prs)
    slide_17_phylo(prs)
    slide_18_limitations(prs)
    slide_19_conclusions(prs)
    slide_20_thank_you(prs)

    # Backup slides
    slide_b01_anova(prs)
    slide_b02_friedman(prs)
    slide_b03_pahd(prs)
    slide_b04_dms(prs)
    slide_b05_bug(prs)
    slide_b06_weights(prs)
    slide_b07_pathogens(prs)
    slide_b08_treetime(prs)
    slide_b09_transfer(prs)
    slide_b10_comparison(prs)

    out_path = "/proj/paper/paper/ppt/MutBench_Defense.pptx"
    prs.save(out_path)
    print(f"Presentation saved to {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
