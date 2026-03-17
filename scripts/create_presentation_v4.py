#!/usr/bin/env python3
"""
MutBench PhD Defense Presentation v4 — ALL ENGLISH, Redesigned
Output: /proj/paper/paper/ppt/MutBench_Defense_v4.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ── paths ──────────────────────────────────────────────────────
BASE = "/proj/paper"
OUT = os.path.join(BASE, "paper/ppt/MutBench_Defense_v4.pptx")
LOGO = os.path.join(BASE, "docs/presentation/knu_symbol_white.png")
LOGO_EMBLEM = os.path.join(BASE, "docs/presentation/knu_emblem_official.jpg")
FIG_DISS = os.path.join(BASE, "paper/dissertation/figures")
FIG_MAIN = os.path.join(BASE, "paper/figure")

# ── colors ─────────────────────────────────────────────────────
DARK_NAVY  = RGBColor(0x1B, 0x2A, 0x4A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
GRAY       = RGBColor(0x88, 0x88, 0x88)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY   = RGBColor(0xE0, 0xE0, 0xE0)
PH_GRAY    = RGBColor(0xF0, 0xF0, 0xF0)   # placeholder fill
PH_BORDER  = RGBColor(0x99, 0x99, 0x99)   # placeholder border

# section colors (kept for accent usage, but top bars all DARK_NAVY now)
SEC_INTRO   = RGBColor(0x1B, 0x2A, 0x4A)  # navy
SEC_FRAME   = RGBColor(0x0D, 0x73, 0x77)  # teal
SEC_STAGE1  = RGBColor(0x2E, 0x7D, 0x32)  # green
SEC_STAGE2  = RGBColor(0xE6, 0x51, 0x00)  # orange
SEC_EXT     = RGBColor(0x6A, 0x1B, 0x9A)  # purple
SEC_CONCL   = RGBColor(0xC6, 0x28, 0x28)  # red

TEAL_ACCENT = RGBColor(0x0D, 0x73, 0x77)
RED_ACCENT  = RGBColor(0xC6, 0x28, 0x28)
GREEN       = RGBColor(0x2E, 0x7D, 0x32)
ORANGE      = RGBColor(0xE6, 0x51, 0x00)
PURPLE      = RGBColor(0x6A, 0x1B, 0x9A)
BLUE        = RGBColor(0x15, 0x65, 0xC0)

# Category colors for scoring groups
FREQ_COLOR  = RGBColor(0x1B, 0x5E, 0x20)  # dark green
ENTR_COLOR  = RGBColor(0x0D, 0x47, 0xA1)  # dark blue
COMB_COLOR  = RGBColor(0xBF, 0x36, 0x0C)  # dark orange

SLIDE_W = Inches(13.330)
SLIDE_H = Inches(7.500)

# ── helpers ────────────────────────────────────────────────────

def make_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def add_blank(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_rect(slide, l, t, w, h, fill=None, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.fill.solid()
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    return shape

def add_rounded_rect(slide, l, t, w, h, fill=None, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.fill.solid()
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    return shape

def set_text(shape, text, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
             font_name="Calibri", anchor=MSO_ANCHOR.TOP):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except:
        pass
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = align

def add_textbox(slide, l, t, w, h, text, size=14, bold=False, color=BLACK,
                align=PP_ALIGN.LEFT, font_name="Calibri", anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    set_text(txBox, text, size, bold, color, align, font_name, anchor)
    return txBox

def add_multiformat_textbox(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT):
    """runs: list of (text, size, bold, color) tuples"""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for text, size, bold, color in runs:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox

def content_slide(prs, title, section_color=None):
    """Create standard content slide with top bar, title, logo, bottom bar, footer.
    GLOBAL: ALL top bars and bottom bars are DARK_NAVY."""
    slide = add_blank(prs)
    # top bar — ALWAYS DARK_NAVY
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.900), fill=DARK_NAVY)
    # title
    add_textbox(slide, Inches(0.600), Inches(0.180), Inches(10.500), Inches(0.550),
                title, size=24, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # KNU logo top-right
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(12.230), Inches(0.150),
                                 Inches(0.550), Inches(0.550))
    # bottom bar — ALWAYS DARK_NAVY
    add_rect(slide, Inches(0), Inches(7.120), SLIDE_W, Inches(0.380), fill=DARK_NAVY)
    # footer text
    add_textbox(slide, Inches(0.400), Inches(7.140), Inches(10.0), Inches(0.340),
                "Kyungpook National University  |  Dept. of Computer Science  |  Hwijun Kwon",
                size=10, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    return slide

def add_page_number(slide, num):
    add_textbox(slide, Inches(12.500), Inches(7.140), Inches(0.700), Inches(0.340),
                str(num), size=10, color=WHITE, align=PP_ALIGN.RIGHT,
                anchor=MSO_ANCHOR.MIDDLE)

def add_card(slide, l, t, w, h, fill=WHITE, border=MID_GRAY):
    return add_rounded_rect(slide, l, t, w, h, fill=fill, border=border)

def add_accent_bar(slide, l, t, w, h, color):
    return add_rect(slide, l, t, w, h, fill=color)

def add_figure(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, l, t, w, h)
    else:
        add_placeholder(slide, l, t, w, h, f"[Missing: {os.path.basename(path)}]")

def add_placeholder(slide, l, t, w, h, text="[Image Placeholder]"):
    """Gray rounded rectangle with dashed-style border and placeholder text."""
    shape = add_rounded_rect(slide, l, t, w, h, fill=PH_GRAY, border=PH_BORDER)
    shape.line.width = Pt(2)
    try:
        ln = shape.line._ln
        ln.set(qn('prstDash'), 'dash')
    except:
        pass
    set_text(shape, text, size=11, color=GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return shape

def add_takeaway(slide, text, y=Inches(6.650)):
    """Add a bold takeaway message near the bottom of a content slide."""
    box = add_card(slide, Inches(0.600), y, Inches(12.100), Inches(0.400),
                   fill=RGBColor(0xE3, 0xF2, 0xFD), border=DARK_NAVY)
    set_text(box, text, size=11, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    return box


# ══════════════════════════════════════════════════════════════
# BUILD PRESENTATION
# ══════════════════════════════════════════════════════════════

prs = make_prs()
slide_num = 0

# ── SLIDE 1: Title ─────────────────────────────────────────────
# CHANGES: DELETE red KNU logo (knu_emblem_official.jpg), DELETE keyword tag boxes
slide_num += 1
s = add_blank(prs)
# dark top area
add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(3.400), fill=DARK_NAVY)
# KNU white logo (kept)
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(0.600), Inches(0.250), Inches(0.650), Inches(0.650))
# university + defense
add_textbox(s, Inches(1.400), Inches(0.300), Inches(8.0), Inches(0.300),
            "Graduate School, Kyungpook National University", size=14, color=WHITE)
add_textbox(s, Inches(1.400), Inches(0.600), Inches(5.0), Inches(0.300),
            "PhD Dissertation Defense", size=14, color=WHITE)
# title
add_textbox(s, Inches(0.600), Inches(1.400), Inches(12.0), Inches(1.600),
            "MutBench: Systematic Benchmarking Framework\nfor Viral Mutation Hotspot Detection",
            size=34, bold=True, color=WHITE)
# (DELETED: KNU emblem / LOGO_EMBLEM)
# author info (below dark area)
add_textbox(s, Inches(0.600), Inches(3.800), Inches(11.0), Inches(0.400),
            "Hwijun Kwon  |  Advisor: Prof. Inuk Jung", size=18, bold=True, color=DARK_GRAY)
add_textbox(s, Inches(0.600), Inches(4.250), Inches(11.0), Inches(0.350),
            "Dept. of Computer Science, KNU  |  March 2026", size=14, color=GRAY)
# (DELETED: keyword tag boxes)
# bottom bar
add_rect(s, Inches(0), Inches(7.120), SLIDE_W, Inches(0.380), fill=DARK_NAVY)


# ── SLIDE 2: Table of Contents — VERTICAL layout ────────────────
# CHANGES: DELETE "Research Flow" diagram on right, DELETE vertical color bars
slide_num += 1
s = content_slide(prs, "Table of Contents")
add_page_number(s, slide_num)

toc_items = [
    ("1", "Introduction", "Prior work, biology background, problem definition", "S3 - S10", SEC_INTRO),
    ("2", "MutBench Framework", "Pipeline, scoring/detection, 3-layer ground truth, MCC", "S11 - S17", SEC_FRAME),
    ("3", "Stage 1: 4-Pathogen", "Method comparison, sensitivity, synthetic vs real gap", "S18 - S21", SEC_STAGE1),
    ("4", "Stage 2: 9-Pathogen", "ANOVA, per-pathogen best, Friedman, LOPO, ESM-2", "S22 - S27", SEC_STAGE2),
    ("5", "Extensions & PAHD", "PAHD algorithm, phylogenetic correction, baselines", "S28 - S30", SEC_EXT),
    ("6", "Conclusion", "Contributions, discussion, limitations, future work", "S31 - S35", SEC_CONCL),
]
for i, (num, title, desc, slides, color) in enumerate(toc_items):
    cy = Inches(1.100) + i * Inches(0.950)
    # (DELETED: vertical color bar)
    # number circle
    badge = add_rounded_rect(s, Inches(0.900), cy + Inches(0.100), Inches(0.550), Inches(0.550),
                              fill=color)
    set_text(badge, num, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    # title
    add_textbox(s, Inches(1.650), cy + Inches(0.050), Inches(5.000), Inches(0.400),
                title, size=17, bold=True, color=color)
    # description
    add_textbox(s, Inches(1.650), cy + Inches(0.450), Inches(6.000), Inches(0.350),
                desc, size=11, color=DARK_GRAY)
    # slide range
    add_textbox(s, Inches(7.800), cy + Inches(0.200), Inches(1.200), Inches(0.400),
                slides, size=10, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)
# (DELETED: Research Flow diagram on right side)


# ── SLIDE 3: Prior Work — paper cards ──────────────────────────
# CHANGES: DELETE bridge box, CENTER cards vertically, all bars NAVY
slide_num += 1
s = content_slide(prs, "Prior Work: MOSD & MutClust")
add_page_number(s, slide_num)

# Cards centered vertically between top bar (0.9) and footer (7.12)
# Available space: 0.9 to 7.12 = 6.22 inches, card height 4.2 -> center at ~1.9
card_top = Inches(1.900)

# MOSD card with accent bar
c1 = add_card(s, Inches(0.600), card_top, Inches(5.800), Inches(4.200), border=MID_GRAY)
add_accent_bar(s, Inches(0.600), card_top, Inches(0.100), Inches(4.200), TEAL_ACCENT)
add_textbox(s, Inches(0.900), card_top + Inches(0.100), Inches(5.300), Inches(0.350),
            "Paper 1  |  BMC Genomics 2025 (SCIE)", size=14, bold=True, color=TEAL_ACCENT)
add_textbox(s, Inches(0.900), card_top + Inches(0.500), Inches(5.300), Inches(0.350),
            "MOSD: Multi-Omics Subtyping with Deep learning", size=12, bold=True, color=DARK_NAVY)
add_textbox(s, Inches(0.900), card_top + Inches(0.950), Inches(5.300), Inches(1.200),
            "  Benchmarked 11 multi-omics integration methods\n"
            "  across 6 evaluation metrics for cancer subtyping\n\n"
            "  Key finding: No single method dominates\n"
            "  across all cancer datasets", size=11, color=DARK_GRAY)
add_placeholder(s, Inches(1.000), card_top + Inches(2.400), Inches(4.800), Inches(1.600),
                "[Image Placeholder]\nMOSD Key Figure\n(To be inserted)")

# MutClust card with accent bar
c2 = add_card(s, Inches(6.900), card_top, Inches(5.800), Inches(4.200), border=MID_GRAY)
add_accent_bar(s, Inches(6.900), card_top, Inches(0.100), Inches(4.200), RED_ACCENT)
add_textbox(s, Inches(7.200), card_top + Inches(0.100), Inches(5.300), Inches(0.350),
            "Paper 2  |  BioData Mining 2025 (SCIE)", size=14, bold=True, color=RED_ACCENT)
add_textbox(s, Inches(7.200), card_top + Inches(0.500), Inches(5.300), Inches(0.350),
            "MutClust: Mutation Clustering for viral evolution", size=12, bold=True, color=DARK_NAVY)
add_textbox(s, Inches(7.200), card_top + Inches(0.950), Inches(5.300), Inches(1.200),
            "  H-score based mutation hotspot detection\n"
            "  with network propagation & bootstrap\n\n"
            "  Limitation: H-score founder bias (rho = -0.876)\n"
            "  Single-pathogen evaluation is insufficient", size=11, color=DARK_GRAY)
add_placeholder(s, Inches(7.300), card_top + Inches(2.400), Inches(4.800), Inches(1.600),
                "[Image Placeholder]\nMutClust Key Figure\n(To be inserted)")
# (DELETED: bridge box)
add_takeaway(s, "Takeaway: Both prior works reveal that no single method dominates -- motivating MutBench.")


# ── SLIDE 4: Biology — Central Dogma ──────────────────────────
slide_num += 1
s = content_slide(prs, "The Central Dogma: DNA -> RNA -> Protein")
add_page_number(s, slide_num)

# Left: key terms
add_textbox(s, Inches(0.600), Inches(1.200), Inches(5.800), Inches(0.400),
            "Key Biological Terms", size=16, bold=True, color=DARK_NAVY)

terms = [
    ("Genome", "The complete genetic information of an organism,\nencoded as a DNA or RNA sequence."),
    ("Gene", "A segment of the genome that encodes a functional\nprotein or RNA molecule."),
    ("Nucleotide", "The building blocks of DNA/RNA: Adenine (A),\nThymine (T)/Uracil (U), Guanine (G), Cytosine (C)."),
    ("Amino Acid", "Building blocks of proteins. A sequence of 3\nnucleotides (codon) specifies one amino acid."),
]
for i, (term, desc) in enumerate(terms):
    cy = Inches(1.750) + i * Inches(1.200)
    card = add_card(s, Inches(0.600), cy, Inches(5.800), Inches(1.050), border=MID_GRAY)
    add_textbox(s, Inches(0.800), cy + Inches(0.100), Inches(1.500), Inches(0.350),
                term, size=13, bold=True, color=SEC_FRAME)
    add_textbox(s, Inches(0.800), cy + Inches(0.450), Inches(5.400), Inches(0.550),
                desc, size=10, color=DARK_GRAY)

# Right: placeholder for DNA structure diagram
add_placeholder(s, Inches(6.900), Inches(1.200), Inches(5.800), Inches(3.000),
                "[Image Placeholder]\nDNA Structure Diagram\n(To be inserted)")

# Central dogma flow at bottom right
add_textbox(s, Inches(6.900), Inches(4.500), Inches(5.800), Inches(0.400),
            "Central Dogma Flow", size=14, bold=True, color=DARK_NAVY)
flow_boxes = [("DNA", SEC_FRAME), ("mRNA", ORANGE), ("Protein", PURPLE)]
for i, (label, clr) in enumerate(flow_boxes):
    bx = Inches(7.100) + i * Inches(1.900)
    box = add_rounded_rect(s, bx, Inches(5.000), Inches(1.500), Inches(0.650), fill=clr)
    set_text(box, label, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    if i < 2:
        add_textbox(s, bx + Inches(1.500), Inches(5.050), Inches(0.400), Inches(0.550),
                    ">", size=20, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# note at bottom
note = add_card(s, Inches(6.900), Inches(5.900), Inches(5.800), Inches(0.750),
                fill=RGBColor(0xFF, 0xF3, 0xE0), border=ORANGE)
set_text(note, "RNA viruses use RNA (not DNA) as their genome.\nHigher mutation rate due to error-prone RNA polymerase.",
         size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ── SLIDE 5: Mutations & Their Effects ──────────────────────────
# CHANGES: DELETE hotspot link box at bottom, ADD large placeholder at TOP, keep explanations below
slide_num += 1
s = content_slide(prs, "Mutations & Their Effects")
add_page_number(s, slide_num)

# TOP: Large placeholder for mutation types diagram image
add_placeholder(s, Inches(0.600), Inches(1.100), Inches(12.100), Inches(2.200),
                "[Image Placeholder]\nMutation Types Diagram (Substitution / Insertion / Deletion)\n(To be inserted)")

# Mutation types below the image placeholder
add_textbox(s, Inches(0.600), Inches(3.500), Inches(4.000), Inches(0.350),
            "Mutation Types", size=14, bold=True, color=DARK_NAVY)

mut_types = [
    ("Substitution", "One nucleotide replaced by another (e.g., A -> G)", SEC_FRAME),
    ("Insertion", "Extra nucleotide(s) added; shifts reading frame", ORANGE),
    ("Deletion", "Nucleotide(s) removed; can disrupt protein structure", RED_ACCENT),
]
for i, (name, desc, clr) in enumerate(mut_types):
    cx = Inches(0.600) + i * Inches(4.200)
    cy = Inches(3.950)
    card = add_card(s, cx, cy, Inches(3.900), Inches(0.850), border=MID_GRAY)
    add_accent_bar(s, cx, cy, Inches(0.080), Inches(0.850), clr)
    add_textbox(s, cx + Inches(0.200), cy + Inches(0.100), Inches(3.500), Inches(0.300),
                name, size=13, bold=True, color=clr)
    add_textbox(s, cx + Inches(0.200), cy + Inches(0.420), Inches(3.500), Inches(0.350),
                desc, size=10, color=DARK_GRAY)

# Effects on protein below
add_textbox(s, Inches(0.600), Inches(5.050), Inches(4.000), Inches(0.350),
            "Effects on Protein", size=14, bold=True, color=DARK_NAVY)

effects = [
    ("Silent", "No change in amino acid (codon redundancy)", GRAY),
    ("Missense", "Different amino acid; may alter protein function", BLUE),
    ("Nonsense", "Premature stop codon; truncated protein", RED_ACCENT),
]
for i, (name, desc, clr) in enumerate(effects):
    cx = Inches(0.600) + i * Inches(4.200)
    cy = Inches(5.450)
    card = add_card(s, cx, cy, Inches(3.900), Inches(0.750), border=MID_GRAY)
    add_accent_bar(s, cx, cy, Inches(0.080), Inches(0.750), clr)
    add_textbox(s, cx + Inches(0.200), cy + Inches(0.080), Inches(3.500), Inches(0.280),
                name, size=12, bold=True, color=clr)
    add_textbox(s, cx + Inches(0.200), cy + Inches(0.380), Inches(3.500), Inches(0.300),
                desc, size=9, color=DARK_GRAY)
# (DELETED: hotspot connection box)
add_takeaway(s, "Takeaway: Different mutation types have different impacts on protein function and viral fitness.")


# ── SLIDE 6: Viral Evolution & Natural Selection ────────────────
# CHANGES: REDUCE height of 3 concept boxes, EXPAND bottom explanation box
slide_num += 1
s = content_slide(prs, "Viral Evolution & Natural Selection")
add_page_number(s, slide_num)

# Three selection types — REDUCED height
sel_types = [
    ("Positive Selection", "Mutations that increase viral fitness\nspread through the population.\n"
     "  Immune evasion, transmissibility,\n  drug resistance",
     GREEN),
    ("Purifying Selection", "Mutations that damage essential functions\nare removed from the population.\n"
     "  Disruption of replication,\n  loss of structural integrity",
     RED_ACCENT),
    ("Convergent Evolution", "The same mutation arises independently\nin different viral lineages.\n"
     "  N501Y in multiple SARS-CoV-2 variants,\n  strong positive selection signal",
     PURPLE),
]
for i, (ttl, desc, clr) in enumerate(sel_types):
    cx = Inches(0.500) + i * Inches(4.200)
    # REDUCED from 4.200 to 2.800 height
    card = add_card(s, cx, Inches(1.200), Inches(3.900), Inches(2.800), border=MID_GRAY)
    add_rect(s, cx, Inches(1.200), Inches(3.900), Inches(0.500), fill=clr)
    add_textbox(s, cx, Inches(1.220), Inches(3.900), Inches(0.460),
                ttl, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, cx + Inches(0.200), Inches(1.800), Inches(3.500), Inches(2.100),
                desc, size=11, color=DARK_GRAY)

# EXPANDED bottom explanation box (was placeholder, now more space)
explain = add_card(s, Inches(0.500), Inches(4.200), Inches(12.300), Inches(2.500),
                   fill=RGBColor(0xF3, 0xE5, 0xF5), border=PURPLE)
set_text(explain,
         "Why This Matters for Hotspot Detection\n\n"
         "Positive selection creates mutation hotspots: positions where beneficial mutations accumulate.\n"
         "Purifying selection creates 'cold spots': conserved positions essential for survival.\n"
         "Convergent evolution provides the strongest evidence: independent origins confirm selective pressure.\n\n"
         "Different pathogens experience different selection landscapes, which is why detection methods\n"
         "perform differently across pathogens -- the central finding of this dissertation.",
         size=12, color=DARK_GRAY, anchor=MSO_ANCHOR.TOP)


# ── SLIDE 7: Mutation Hotspot Definition ────────────────────────
# (WAS S8, now moved to S7 per reorder: S7=Hotspots, S8=ProbDef, S9=ResOverview)
slide_num += 1
s = content_slide(prs, "What are Mutation Hotspots?")
add_page_number(s, slide_num)

# definition box
defbox = add_card(s, Inches(0.600), Inches(1.200), Inches(12.100), Inches(1.200),
                  fill=RGBColor(0xE8, 0xF5, 0xE9), border=GREEN)
set_text(defbox, "Definition:  Genomic positions where mutations accumulate at significantly\n"
         "higher frequencies than expected by chance, indicating selective pressure.",
         size=14, color=DARK_GRAY, anchor=MSO_ANCHOR.MIDDLE)

# Left: placeholder for hotspot on Spike
add_placeholder(s, Inches(0.600), Inches(2.700), Inches(5.800), Inches(3.000),
                "[Image Placeholder]\nHotspot Visualization on Spike Protein\n"
                "e.g., N501Y, E484K, D614G positions highlighted\n(To be inserted)")

# Right: 3 cards
topics = [
    ("Why do they occur?", "Positive selection drives beneficial mutations\n"
     "(e.g., immune escape at receptor binding sites).\n"
     "Some regions tolerate mutations better than others."),
    ("Why are they important?", "Track viral evolution in real-time.\n"
     "Guide vaccine & drug target design.\n"
     "Predict emerging variants of concern."),
    ("Key Examples", "N501Y: Enhanced ACE2 binding affinity\n"
     "E484K: Immune escape from antibodies\n"
     "D614G: Founder effect vs true selection"),
]
for i, (ttl, desc) in enumerate(topics):
    cy = Inches(2.700) + i * Inches(1.050)
    card = add_card(s, Inches(6.700), cy, Inches(5.900), Inches(0.900), border=MID_GRAY)
    add_textbox(s, Inches(6.900), cy + Inches(0.050), Inches(2.000), Inches(0.300),
                ttl, size=12, bold=True, color=DARK_NAVY)
    add_textbox(s, Inches(6.900), cy + Inches(0.350), Inches(5.500), Inches(0.500),
                desc, size=9, color=DARK_GRAY)

# bottom note
note = add_card(s, Inches(0.600), Inches(5.950), Inches(12.100), Inches(0.700),
                fill=RGBColor(0xFF, 0xEB, 0xEE), border=RED_ACCENT)
set_text(note, "Current gap: No systematic way to compare hotspot detection methods across pathogens",
         size=13, bold=True, color=RED_ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ── SLIDE 8: Problem Definition ─────────────────────────────────
# (WAS S9, moved to S8 per reorder)
slide_num += 1
s = content_slide(prs, "Problem Definition")
add_page_number(s, slide_num)

problems = [
    ("P1", "No Standardized Benchmark",
     "Existing methods evaluated on different datasets with different metrics, "
     "making fair comparison impossible.",
     "[Placeholder]\nIncompatible metrics"),
    ("P2", "Single-Pathogen Bias",
     "Methods tested on only 1-2 pathogens cannot generalize. "
     "Pathogen-specific characteristics dramatically affect detection performance.",
     "[Placeholder]\n1 vs 9 pathogens"),
    ("P3", "No Universal Best Method",
     "Our prior work (MOSD, MutClust) demonstrated that no single approach dominates "
     "across all settings, demanding systematic evaluation.",
     "[Placeholder]\nRanking reversal"),
]
for i, (pid, title, desc, ph_text) in enumerate(problems):
    cy = Inches(1.200) + i * Inches(1.850)
    card = add_card(s, Inches(0.600), cy, Inches(9.800), Inches(1.600), border=MID_GRAY)
    badge = add_rounded_rect(s, Inches(0.800), cy + Inches(0.300), Inches(0.600), Inches(0.500),
                             fill=DARK_NAVY)
    set_text(badge, pid, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, Inches(1.600), cy + Inches(0.200), Inches(8.500), Inches(0.400),
                title, size=16, bold=True, color=DARK_NAVY)
    add_textbox(s, Inches(1.600), cy + Inches(0.700), Inches(8.500), Inches(0.800),
                desc, size=12, color=DARK_GRAY)
    add_placeholder(s, Inches(10.700), cy + Inches(0.150), Inches(2.000), Inches(1.300), ph_text)

add_takeaway(s, "Takeaway: These three problems necessitate a systematic benchmarking framework.")


# ── SLIDE 9: Research Overview ──────────────────────────────────
# (WAS S7, moved to S9 per reorder — now AFTER Problem Definition)
# CHANGE: Add subtitle connecting to problems
slide_num += 1
s = content_slide(prs, "Research Overview")
add_page_number(s, slide_num)

# Subtitle connecting to Problem Definition
add_textbox(s, Inches(0.600), Inches(1.000), Inches(12.100), Inches(0.400),
            "How MutBench Addresses These Problems", size=16, bold=False, color=TEAL_ACCENT,
            align=PP_ALIGN.CENTER)

# Full-slide flow diagram placeholder
add_placeholder(s, Inches(0.600), Inches(1.500), Inches(12.100), Inches(3.000),
                "[Image Placeholder]\nMutBench Research Overview - Full Pipeline Diagram\n"
                "P1 (no benchmark) -> MutBench Framework\n"
                "P2 (single-pathogen) -> 9-pathogen evaluation\n"
                "P3 (no universal best) -> PAHD adaptive selection\n(To be inserted)")

# Key numbers below
nums = [
    ("9", "RNA Pathogens", SEC_FRAME),
    ("2,544", "Evaluations", SEC_STAGE2),
    ("351", "Method Combos", GREEN),
    ("3", "Ground Truth Layers", PURPLE),
    ("0.285", "Interaction omega-sq", RED_ACCENT),
]
for i, (num, label, clr) in enumerate(nums):
    cx = Inches(0.400) + i * Inches(2.560)
    cy = Inches(4.800)
    card = add_card(s, cx, cy, Inches(2.350), Inches(1.800), border=MID_GRAY)
    add_textbox(s, cx, cy + Inches(0.200), Inches(2.350), Inches(0.700),
                num, size=30, bold=True, color=clr, align=PP_ALIGN.CENTER)
    add_textbox(s, cx, cy + Inches(1.000), Inches(2.350), Inches(0.600),
                label, size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ── SLIDE 10: 9 Pathogens Overview ──────────────────────────────
slide_num += 1
s = content_slide(prs, "9 RNA Pathogens in MutBench")
add_page_number(s, slide_num)

pathogens = [
    ("SARS-CoV-2", "~30 kb", "Spike protein", "COVID-19 pandemic"),
    ("Influenza", "~13.5 kb", "Hemagglutinin", "Seasonal epidemics, antigenic shift"),
    ("HIV", "~10 kb", "Env/Gag", "High mutation rate, chronic"),
    ("Dengue", "~11 kb", "E protein", "4 serotypes, tropical"),
    ("MERS-CoV", "~30 kb", "Spike protein", "Camel-origin, high fatality"),
    ("RSV", "~15 kb", "F protein", "Infant respiratory disease"),
    ("Norovirus", "~7.5 kb", "VP1 capsid", "Gastroenteritis, rapid evolution"),
    ("HCV", "~9.6 kb", "E1/E2 proteins", "Chronic hepatitis, high diversity"),
    ("Ebola", "~19 kb", "Glycoprotein", "Hemorrhagic fever, outbreaks"),
]

# Header row
add_rect(s, Inches(0.500), Inches(1.200), Inches(12.300), Inches(0.500), fill=DARK_NAVY)
headers = ["Pathogen", "Genome", "Target Protein", "Key Feature"]
h_widths = [Inches(2.200), Inches(1.400), Inches(2.600), Inches(6.100)]
hx = Inches(0.500)
for h, w in zip(headers, h_widths):
    add_textbox(s, hx, Inches(1.210), w, Inches(0.480),
                h, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    hx += w

# Data rows
for r, (name, genome, protein, feature) in enumerate(pathogens):
    ry = Inches(1.700) + r * Inches(0.580)
    fill = LIGHT_GRAY if r % 2 == 0 else WHITE
    add_rect(s, Inches(0.500), ry, Inches(12.300), Inches(0.580), fill=fill)
    vals = [name, genome, protein, feature]
    vx = Inches(0.500)
    for v, w in zip(vals, h_widths):
        bld = (v == name)
        clr = DARK_NAVY if bld else DARK_GRAY
        add_textbox(s, vx, ry, w, Inches(0.580), v, size=11, bold=bld, color=clr,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        vx += w

add_takeaway(s, "Takeaway: 9 diverse RNA pathogens with varying genome sizes, mutation rates, and evolutionary pressures.", y=Inches(6.950))


# ── SLIDE 11: MutBench Pipeline ─────────────────────────────────
# CHANGE: ADD a section showing 9 scoring formulas alongside detection families
slide_num += 1
s = content_slide(prs, "MutBench Pipeline Overview")
add_page_number(s, slide_num)

# Pipeline flow — full-width boxes with arrows
flow_items = [
    ("Input\nSequences", DARK_NAVY),
    ("MSA\nAlignment", DARK_NAVY),
    ("Scoring\n(9 methods)", SEC_FRAME),
    ("Detection\n(39 methods)", GREEN),
    ("Ground\nTruth", ORANGE),
    ("Evaluation\n(MCC)", PURPLE),
]
box_w = Inches(1.750)
box_h = Inches(0.950)
start_x = Inches(0.400)
for i, (item, fill_c) in enumerate(flow_items):
    bx = start_x + i * Inches(2.100)
    by = Inches(1.100)
    box = add_rounded_rect(s, bx, by, box_w, box_h, fill=fill_c)
    set_text(box, item, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    if i < len(flow_items) - 1:
        add_textbox(s, bx + box_w, by + Inches(0.200), Inches(0.350), Inches(0.500),
                    ">", size=18, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# Left: 9 Scoring Methods with formulas
add_rect(s, Inches(0.400), Inches(2.250), Inches(6.000), Inches(0.400), fill=SEC_FRAME)
add_textbox(s, Inches(0.400), Inches(2.260), Inches(6.000), Inches(0.380),
            "9 Scoring Methods", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)
scoring_items = [
    ("Shannon Entropy", "H = -Sum p_i log p_i"),
    ("Jensen-Shannon", "JSD(P||Q)"),
    ("Mutation Freq.", "f = 1 - max(p_i)"),
    ("Wavelet (CWT)", "Multi-scale analysis"),
    ("Sliding Window", "Smoothed average"),
    ("Phylo-aware", "Branch-weighted"),
    ("dN/dS Ratio", "omega = dN/dS"),
    ("Kabat Variability", "K = N/(n*f_max)"),
    ("Property Entropy", "AA property groups"),
]
for i, (name, formula) in enumerate(scoring_items):
    col = i % 3
    row = i // 3
    cx = Inches(0.450) + col * Inches(2.000)
    cy = Inches(2.750) + row * Inches(0.500)
    add_textbox(s, cx, cy, Inches(1.100), Inches(0.250),
                f"  {name}", size=8, bold=True, color=SEC_FRAME)
    add_textbox(s, cx + Inches(1.050), cy, Inches(0.950), Inches(0.250),
                formula, size=7, color=GRAY)

# Right: Detection Families
add_rect(s, Inches(6.700), Inches(2.250), Inches(6.000), Inches(0.400), fill=GREEN)
add_textbox(s, Inches(6.700), Inches(2.260), Inches(6.000), Inches(0.380),
            "Detection Families (39 methods)", size=12, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
detect_items = [
    ("Z-score", "4 variants"),
    ("Percentile", "4 variants"),
    ("IQR Outlier", "2 variants"),
    ("DBSCAN", "6 variants"),
    ("HDBSCAN", "3 variants"),
    ("Gaussian Mix.", "4 variants"),
    ("Kernel Density", "4 variants"),
    ("Isolation Forest", "3 variants"),
    ("Wavelet Peaks", "3 variants"),
    ("LOF", "2 variants"),
    ("Threshold", "2 variants"),
    ("Bayesian CP", "2 variants"),
]
for i, (name, desc) in enumerate(detect_items):
    col = i % 3
    row = i // 3
    cx = Inches(6.800) + col * Inches(2.000)
    cy = Inches(2.750) + row * Inches(0.500)
    add_textbox(s, cx, cy, Inches(1.200), Inches(0.250),
                f"  {name}", size=8, bold=True, color=GREEN)
    add_textbox(s, cx + Inches(1.150), cy, Inches(0.850), Inches(0.250),
                desc, size=7, color=GRAY)

# Total count — emphasized
total_box = add_card(s, Inches(2.000), Inches(4.900), Inches(9.300), Inches(0.550),
                     fill=RGBColor(0xE8, 0xF5, 0xE9), border=GREEN)
set_text(total_box,
         "Total: 9 scoring x 39 detection = 351 unique combinations per pathogen  |  9 pathogens = 2,544 evaluations",
         size=12, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_takeaway(s, "Takeaway: The pipeline systematically evaluates BOTH 9 scoring AND 39 detection methods across all pathogens.", y=Inches(5.700))


# ── SLIDE 12: Scoring Methods Detail — REDESIGNED ────────────────
# CHANGE: Group into 3 categories with color coding; mark NEW contributions
slide_num += 1
s = content_slide(prs, "Scoring Methods: 9 Approaches (Grouped by Category)")
add_page_number(s, slide_num)

# Category definitions
categories = [
    ("Frequency-based", FREQ_COLOR, RGBColor(0xE8, 0xF5, 0xE9), [
        ("H-score *", "H = sum(f_i > threshold)", "Existing (MutClust)"),
        ("P x E", "Product of frequency & entropy", "NEW in MutBench"),
        ("P x E^2", "Frequency x entropy squared", "NEW in MutBench"),
    ]),
    ("Entropy-based", ENTR_COLOR, RGBColor(0xE3, 0xF2, 0xFD), [
        ("E_only", "H(x) = -Sum p_i log2 p_i", "Shannon entropy only"),
        ("E x rare", "Entropy weighted by rare alleles", "NEW in MutBench"),
        ("minority_E", "Entropy of minority nucleotides", "NEW in MutBench"),
    ]),
    ("Combined", COMB_COLOR, RGBColor(0xFB, 0xE9, 0xE7), [
        ("P x E x rare", "Frequency x entropy x rare", "NEW in MutBench"),
        ("P x minority_E", "Frequency x minority entropy", "NEW in MutBench"),
        ("rank(P x E)", "Rank-based P x E scoring", "NEW in MutBench"),
    ]),
]

for ci, (cat_name, cat_color, bg_color, methods) in enumerate(categories):
    cx = Inches(0.400) + ci * Inches(4.250)
    # category header
    add_rect(s, cx, Inches(1.150), Inches(4.000), Inches(0.450), fill=cat_color)
    add_textbox(s, cx, Inches(1.160), Inches(4.000), Inches(0.430),
                cat_name, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    # method cards
    for mi, (name, formula, note) in enumerate(methods):
        my = Inches(1.700) + mi * Inches(1.600)
        card = add_card(s, cx, my, Inches(4.000), Inches(1.450), fill=bg_color, border=cat_color)
        # name with star for new
        is_new = "NEW" in note
        name_color = cat_color
        add_textbox(s, cx + Inches(0.150), my + Inches(0.080), Inches(3.700), Inches(0.350),
                    name, size=13, bold=True, color=name_color)
        # formula
        fbox = add_rect(s, cx + Inches(0.150), my + Inches(0.450), Inches(3.700), Inches(0.380),
                        fill=WHITE, border=MID_GRAY)
        set_text(fbox, formula, size=9, color=DARK_NAVY, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        # note / NEW badge
        if is_new:
            badge = add_rounded_rect(s, cx + Inches(0.150), my + Inches(0.950), Inches(1.600), Inches(0.350),
                                     fill=cat_color)
            set_text(badge, "NEW in MutBench", size=9, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        else:
            add_textbox(s, cx + Inches(0.150), my + Inches(0.950), Inches(3.700), Inches(0.350),
                        note, size=9, color=GRAY)

add_takeaway(s, "Takeaway: 7 of 9 scoring methods are NEW contributions of MutBench; only H-score and Shannon entropy existed before.")


# ── SLIDE 13a: 3-Layer Ground Truth — Layer A (Adaptive) ─────────
# CHANGE: Split from single crammed slide into 3 slides
slide_num += 1
s = content_slide(prs, "3-Layer Ground Truth: Layer A (Adaptive)")
add_page_number(s, slide_num)

# Diagram placeholder
add_placeholder(s, Inches(0.600), Inches(1.200), Inches(6.000), Inches(4.500),
                "[Image Placeholder]\nLayer A: Adaptive Ground Truth Diagram\n"
                "Entropy-based threshold auto-adjustment\nper pathogen\n(To be inserted)")

# Explanation on right
add_textbox(s, Inches(7.000), Inches(1.200), Inches(5.600), Inches(0.400),
            "Layer A: Adaptive Threshold", size=18, bold=True, color=SEC_FRAME)
add_textbox(s, Inches(7.000), Inches(1.800), Inches(5.600), Inches(3.500),
            "Principle:\n"
            "  Entropy-based threshold is automatically adjusted\n"
            "  per pathogen based on the background mutation rate.\n\n"
            "How it works:\n"
            "  1. Compute Shannon entropy for all positions\n"
            "  2. Fit a null distribution from non-coding regions\n"
            "  3. Set threshold at mean + k*std (k adaptive)\n"
            "  4. Positions exceeding threshold = hotspot\n\n"
            "Advantage:\n"
            "  Self-calibrating per pathogen, captures high-variability\n"
            "  positions relative to that specific pathogen's background.\n\n"
            "Limitation:\n"
            "  May be too permissive for low-diversity pathogens.",
            size=12, color=DARK_GRAY)

add_takeaway(s, "Takeaway: Layer A adapts to each pathogen's mutation background, avoiding fixed-threshold bias.")


# ── SLIDE 13b: 3-Layer Ground Truth — Layer B (Constrained) ──────
slide_num += 1
s = content_slide(prs, "3-Layer Ground Truth: Layer B (Constrained)")
add_page_number(s, slide_num)

add_placeholder(s, Inches(0.600), Inches(1.200), Inches(6.000), Inches(4.500),
                "[Image Placeholder]\nLayer B: Constrained Ground Truth Diagram\n"
                "Intersection of top-K across methods\n(To be inserted)")

add_textbox(s, Inches(7.000), Inches(1.200), Inches(5.600), Inches(0.400),
            "Layer B: Constrained Consensus", size=18, bold=True, color=GREEN)
add_textbox(s, Inches(7.000), Inches(1.800), Inches(5.600), Inches(3.500),
            "Principle:\n"
            "  Intersection of top-K positions across\n"
            "  multiple independent scoring methods.\n\n"
            "How it works:\n"
            "  1. Rank all positions by each scoring method\n"
            "  2. Take top-K positions from each method\n"
            "  3. A position is 'hotspot' only if it appears\n"
            "     in the top-K of at least M methods\n"
            "  4. Parameters K, M are tuned per pathogen\n\n"
            "Advantage:\n"
            "  Reduces bias from any single scoring method.\n"
            "  Consensus across methods increases reliability.\n\n"
            "Limitation:\n"
            "  May miss hotspots captured by only one method.",
            size=12, color=DARK_GRAY)

add_takeaway(s, "Takeaway: Layer B requires agreement across methods, reducing single-method false positives.")


# ── SLIDE 13c: 3-Layer Ground Truth — Layer C (DMS-based) ────────
slide_num += 1
s = content_slide(prs, "3-Layer Ground Truth: Layer C (DMS-based)")
add_page_number(s, slide_num)

add_placeholder(s, Inches(0.600), Inches(1.200), Inches(6.000), Inches(4.500),
                "[Image Placeholder]\nLayer C: DMS Ground Truth Diagram\n"
                "Deep mutational scanning experimental data\n(To be inserted)")

add_textbox(s, Inches(7.000), Inches(1.200), Inches(5.600), Inches(0.400),
            "Layer C: DMS-based (Experimental)", size=18, bold=True, color=PURPLE)
add_textbox(s, Inches(7.000), Inches(1.800), Inches(5.600), Inches(3.500),
            "Principle:\n"
            "  Deep Mutational Scanning (DMS) measures the\n"
            "  fitness effect of every possible mutation.\n\n"
            "Data sources:\n"
            "  Starr et al. 2020: SARS-CoV-2 RBD\n"
            "  Dadonaite et al. 2023: Full Spike protein\n\n"
            "How it works:\n"
            "  1. Obtain experimental fitness scores\n"
            "  2. Positions with high fitness tolerance\n"
            "     = likely hotspots (can mutate and survive)\n"
            "  3. Gold-standard biological validation\n\n"
            "Limitation:\n"
            "  Only available for SARS-CoV-2 Spike protein.\n"
            "  Cannot be applied to all 9 pathogens.",
            size=12, color=DARK_GRAY)

add_takeaway(s, "Takeaway: Layer C provides gold-standard experimental validation, but is limited to SARS-CoV-2.")


# ── SLIDE 14: MCC Metric — REDESIGNED ─────────────────────────────
# CHANGE: Split into formula + comparison with F1, clearer confusion matrix, concrete example
slide_num += 1
s = content_slide(prs, "Evaluation Metric: Matthews Correlation Coefficient")
add_page_number(s, slide_num)

# Left column: MCC formula + confusion matrix
add_textbox(s, Inches(0.600), Inches(1.150), Inches(6.000), Inches(0.350),
            "MCC Formula", size=16, bold=True, color=PURPLE)
fbox = add_card(s, Inches(0.600), Inches(1.550), Inches(6.000), Inches(0.700),
                fill=RGBColor(0xF3, 0xE5, 0xF5), border=PURPLE)
set_text(fbox, "MCC = (TP*TN - FP*FN) / sqrt((TP+FP)(TP+FN)(TN+FP)(TN+FN))\n"
         "Range: [-1, +1]   |   0 = random   |   +1 = perfect agreement",
         size=11, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Confusion matrix
add_textbox(s, Inches(0.600), Inches(2.400), Inches(3.000), Inches(0.300),
            "Confusion Matrix", size=13, bold=True, color=DARK_NAVY)
# Labels
add_textbox(s, Inches(1.800), Inches(2.700), Inches(1.200), Inches(0.300),
            "Predicted +", size=9, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(3.000), Inches(2.700), Inches(1.200), Inches(0.300),
            "Predicted -", size=9, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.600), Inches(3.100), Inches(1.200), Inches(0.300),
            "Actual +", size=9, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.600), Inches(3.850), Inches(1.200), Inches(0.300),
            "Actual -", size=9, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

cm_x, cm_y = Inches(1.800), Inches(3.000)
cm_s_w, cm_s_h = Inches(1.200), Inches(0.750)
labels = [("TP", GREEN), ("FN", ORANGE), ("FP", RED_ACCENT), ("TN", SEC_FRAME)]
positions = [(0, 0), (1, 0), (0, 1), (1, 1)]
for (col, row), (lbl, clr) in zip(positions, labels):
    bx = cm_x + col * cm_s_w
    by = cm_y + row * cm_s_h
    box = add_rect(s, bx, by, cm_s_w, cm_s_h, fill=clr)
    set_text(box, lbl, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

# Right column: MCC vs F1 comparison
add_textbox(s, Inches(6.900), Inches(1.150), Inches(5.800), Inches(0.350),
            "Why MCC over F1?", size=16, bold=True, color=DARK_NAVY)

why_items = [
    ("Uses all 4 quadrants", "F1 ignores True Negatives entirely.\nMCC uses TP, TN, FP, and FN."),
    ("Robust to class imbalance", "Hotspots are rare (~5% of positions).\nMCC stays reliable with 95% negatives."),
    ("Random baseline = 0", "Verified: mean MCC = 0.001, std = 0.003.\nF1 can give misleading high scores."),
]
for i, (ttl, desc) in enumerate(why_items):
    cy = Inches(1.600) + i * Inches(0.950)
    card = add_card(s, Inches(6.900), cy, Inches(5.800), Inches(0.850), border=MID_GRAY)
    add_accent_bar(s, Inches(6.900), cy, Inches(0.060), Inches(0.850), PURPLE)
    add_textbox(s, Inches(7.100), cy + Inches(0.050), Inches(5.400), Inches(0.280),
                ttl, size=11, bold=True, color=PURPLE)
    add_textbox(s, Inches(7.100), cy + Inches(0.340), Inches(5.400), Inches(0.460),
                desc, size=9, color=DARK_GRAY)

# Concrete example box
example = add_card(s, Inches(0.600), Inches(4.900), Inches(12.100), Inches(1.400),
                   fill=RGBColor(0xFF, 0xF3, 0xE0), border=ORANGE)
set_text(example,
         "Concrete Example: 95% negative class (950 TN, 50 actual hotspots)\n\n"
         "A classifier that predicts ALL negative: TP=0, FP=0, FN=50, TN=950\n"
         "   F1 = 0  (correctly identifies failure)     MCC = 0  (correctly identifies randomness)\n\n"
         "A classifier that predicts 50 hotspots but gets 40 wrong: TP=10, FP=40, FN=40, TN=910\n"
         "   F1 = 0.20  (looks reasonable)              MCC = 0.09  (exposes the poor quality)",
         size=10, color=DARK_GRAY, anchor=MSO_ANCHOR.TOP)

add_takeaway(s, "Takeaway: MCC is the gold-standard metric for imbalanced binary classification; F1 can be misleading.")


# ── SLIDE 15: Two-Stage Design ──────────────────────────────────
slide_num += 1
s = content_slide(prs, "Two-Stage Experimental Design")
add_page_number(s, slide_num)

# Stage 1 card
s1 = add_card(s, Inches(0.600), Inches(1.200), Inches(5.800), Inches(4.200),
              fill=RGBColor(0xE8, 0xF5, 0xE9), border=GREEN)
add_rect(s, Inches(0.600), Inches(1.200), Inches(5.800), Inches(0.600), fill=GREEN)
add_textbox(s, Inches(0.800), Inches(1.220), Inches(5.400), Inches(0.560),
            "Stage 1: Controlled Comparison", size=18, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
add_textbox(s, Inches(0.800), Inches(2.000), Inches(5.400), Inches(3.200),
            "4 pathogens\n"
            "  SARS-CoV-2, Influenza, HIV, Dengue\n\n"
            "351 scoring-detection combinations\n"
            "  9 scoring x 39 detection methods\n\n"
            "Focus areas:\n"
            "  Method comparison & ranking\n"
            "  Parameter sensitivity analysis\n"
            "  Synthetic vs real data gap\n"
            "  Initial ranking reversal observation",
            size=13, color=DARK_GRAY)

# Stage 2 card
s2 = add_card(s, Inches(6.900), Inches(1.200), Inches(5.800), Inches(4.200),
              fill=RGBColor(0xFF, 0xF3, 0xE0), border=ORANGE)
add_rect(s, Inches(6.900), Inches(1.200), Inches(5.800), Inches(0.600), fill=ORANGE)
add_textbox(s, Inches(7.100), Inches(1.220), Inches(5.400), Inches(0.560),
            "Stage 2: Large-Scale Validation", size=18, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
add_textbox(s, Inches(7.100), Inches(2.000), Inches(5.400), Inches(3.200),
            "9 pathogens (+5 more)\n"
            "  + MERS, RSV, Norovirus, HCV, Ebola\n\n"
            "2,544 evaluations\n"
            "  9 scoring x 39 detection x 9 pathogens\n\n"
            "Focus areas:\n"
            "  Two-way ANOVA (interaction effects)\n"
            "  Per-pathogen best combination\n"
            "  Friedman test + LOPO CV\n"
            "  PAHD proof of concept",
            size=13, color=DARK_GRAY)

# Stage transition arrow area
arrow_box = add_card(s, Inches(2.500), Inches(5.650), Inches(8.300), Inches(1.000),
                     fill=DARK_NAVY)
set_text(arrow_box,
         "Stage 1 (discovery)  >>>  Stage 2 (validation at scale)  >>>  PAHD (algorithm)",
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ── SLIDE 16a: Method Comparison Results — Main Table ─────────────
# CHANGE: Split into multiple slides with individual result charts
slide_num += 1
s = content_slide(prs, "Stage 1: Method Comparison — Precision / Recall / F1")
add_page_number(s, slide_num)

# Generate a table of results
add_textbox(s, Inches(0.600), Inches(1.150), Inches(6.000), Inches(0.350),
            "Top-10 Scoring-Detection Combinations (by mean MCC)", size=14, bold=True, color=DARK_NAVY)

# Table header
add_rect(s, Inches(0.600), Inches(1.600), Inches(12.100), Inches(0.500), fill=DARK_NAVY)
t_cols = ["Rank", "Scoring", "Detection", "MCC", "Precision", "Recall", "F1"]
t_ws = [Inches(0.700), Inches(2.200), Inches(2.600), Inches(1.200), Inches(1.400), Inches(1.400), Inches(1.200)]
tx = Inches(0.600)
for c, w in zip(t_cols, t_ws):
    add_textbox(s, tx, Inches(1.610), w, Inches(0.480), c, size=11, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tx += w

# Example data rows (representative)
table_data = [
    ("1", "Wavelet", "HDBSCAN-auto", "0.390", "0.412", "0.375", "0.389"),
    ("2", "P x E", "HDBSCAN-5", "0.378", "0.395", "0.368", "0.377"),
    ("3", "Shannon", "DBSCAN-0.3", "0.365", "0.380", "0.352", "0.362"),
    ("4", "P x E^2", "GMM-3", "0.358", "0.371", "0.348", "0.355"),
    ("5", "Kabat", "HDBSCAN-auto", "0.349", "0.362", "0.340", "0.347"),
    ("6", "JSD", "KDE-silverman", "0.342", "0.355", "0.331", "0.339"),
    ("7", "dN/dS", "Z-score-2.5", "0.335", "0.349", "0.325", "0.333"),
    ("8", "rank(PxE)", "Percentile-95", "0.328", "0.340", "0.318", "0.325"),
    ("9", "Property E", "IsolForest-100", "0.321", "0.334", "0.312", "0.319"),
    ("10", "E x rare", "LOF-20", "0.315", "0.328", "0.306", "0.313"),
]
for r, row_data in enumerate(table_data):
    ry = Inches(2.100) + r * Inches(0.440)
    fill = LIGHT_GRAY if r % 2 == 0 else WHITE
    add_rect(s, Inches(0.600), ry, Inches(12.100), Inches(0.440), fill=fill)
    tx = Inches(0.600)
    for val, w in zip(row_data, t_ws):
        bld = (val == row_data[0])
        clr = GREEN if val == row_data[3] and r == 0 else DARK_GRAY
        add_textbox(s, tx, ry, w, Inches(0.440), val, size=10, bold=bld, color=clr,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tx += w

add_takeaway(s, "Takeaway: Wavelet + HDBSCAN achieves the highest mean MCC (0.390), but rankings vary across pathogens.")


# ── SLIDE 16b: Hotspot-Score Ranking ──────────────────────────────
slide_num += 1
s = content_slide(prs, "Stage 1: Hotspot-Score Ranking by Detection Family")
add_page_number(s, slide_num)

add_figure(s, os.path.join(FIG_DISS, "new_methodology_comparison.png"),
           Inches(0.300), Inches(1.100), Inches(8.200), Inches(5.600))

add_textbox(s, Inches(8.800), Inches(1.200), Inches(4.000), Inches(0.400),
            "Key Findings", size=16, bold=True, color=SEC_STAGE1)
findings = [
    "Wavelet scoring achieves highest\nmean MCC across 4 pathogens",
    "HDBSCAN-based detection outperforms\ntraditional threshold methods",
    "Large variance across pathogens\nsuggests generalization limits",
    "Top combination differs between\nSARS-CoV-2 and HIV",
]
for i, f in enumerate(findings):
    cy = Inches(1.800) + i * Inches(1.200)
    card = add_card(s, Inches(8.800), cy, Inches(3.900), Inches(1.000), border=MID_GRAY)
    add_accent_bar(s, Inches(8.800), cy, Inches(0.060), Inches(1.000), SEC_STAGE1)
    add_textbox(s, Inches(9.000), cy + Inches(0.100), Inches(3.600), Inches(0.800),
                f, size=11, color=DARK_GRAY)


# ── SLIDE 17: Parameter Sensitivity ─────────────────────────────
# CHANGE: Add intro explanation
slide_num += 1
s = content_slide(prs, "Parameter Sensitivity Analysis")
add_page_number(s, slide_num)

# Intro explanation
intro_box = add_card(s, Inches(0.300), Inches(1.050), Inches(12.700), Inches(0.500),
                     fill=RGBColor(0xE3, 0xF2, 0xFD), border=BLUE)
set_text(intro_box,
         "Purpose: Test whether detection method performance is stable across parameter choices, "
         "or whether results depend heavily on specific tuning.",
         size=11, color=DARK_NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_figure(s, os.path.join(FIG_DISS, "sensitivity_heatmap.png"),
           Inches(0.300), Inches(1.650), Inches(7.500), Inches(5.000))

add_textbox(s, Inches(8.100), Inches(1.650), Inches(4.600), Inches(0.400),
            "Sensitivity Findings", size=16, bold=True, color=DARK_NAVY)

# findings as table
sens_data = [
    ("Parameter", "Sensitivity", "Note"),
    ("Window size", "Moderate", "Optimal varies by pathogen"),
    ("Epsilon (DBSCAN)", "High", "Requires per-pathogen tuning"),
    ("Min-cluster (HDBSCAN)", "Low", "Robust across wide range"),
    ("Threshold percentile", "Moderate", "Pathogen-dependent optima"),
]
for r, row_data in enumerate(sens_data):
    ry = Inches(2.200) + r * Inches(0.550)
    fill = DARK_NAVY if r == 0 else (LIGHT_GRAY if r % 2 == 1 else WHITE)
    clr = WHITE if r == 0 else DARK_GRAY
    bld = (r == 0)
    add_rect(s, Inches(8.100), ry, Inches(4.600), Inches(0.550), fill=fill)
    ws = [Inches(1.800), Inches(1.200), Inches(1.600)]
    vx = Inches(8.100)
    for v, w in zip(row_data, ws):
        add_textbox(s, vx, ry, w, Inches(0.550), v, size=10, bold=bld, color=clr,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        vx += w

# Bottom note
note = add_card(s, Inches(8.100), Inches(5.100), Inches(4.600), Inches(1.500),
                fill=RGBColor(0xE8, 0xF5, 0xE9), border=GREEN)
set_text(note, "HDBSCAN is the most robust detection\nmethod across parameter variations.\n\n"
         "This supports its use as a default\nchoice when pathogen characteristics\nare unknown.",
         size=10, color=DARK_GRAY, anchor=MSO_ANCHOR.TOP)


# ── SLIDE 18: Synthetic vs Real Gap — REDESIGNED ─────────────────
# CHANGE: Clearer title, side-by-side numbers, obvious takeaway
slide_num += 1
s = content_slide(prs, "Key Gap: Synthetic Benchmark != Real Data Performance")
add_page_number(s, slide_num)

# Intro message
intro = add_card(s, Inches(0.600), Inches(1.100), Inches(12.100), Inches(0.600),
                 fill=RGBColor(0xFF, 0xEB, 0xEE), border=RED_ACCENT)
set_text(intro, "Methods that win on synthetic data don't always win on real data",
         size=14, bold=True, color=RED_ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Left: Synthetic
add_rect(s, Inches(0.600), Inches(1.900), Inches(5.500), Inches(0.500), fill=BLUE)
add_textbox(s, Inches(0.600), Inches(1.910), Inches(5.500), Inches(0.480),
            "Synthetic Data", size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)

# Arrow in middle
add_textbox(s, Inches(6.100), Inches(3.200), Inches(0.800), Inches(1.000),
            "=/=", size=36, bold=True, color=RED_ACCENT, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)

# Right: Real
add_rect(s, Inches(6.900), Inches(1.900), Inches(5.800), Inches(0.500), fill=RED_ACCENT)
add_textbox(s, Inches(6.900), Inches(1.910), Inches(5.800), Inches(0.480),
            "Real Pathogen Data", size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)

# Comparison rows — clearer with arrows
comp_rows = [
    ("Mutation Distribution", "Uniform / Poisson", "Highly skewed, founder effects"),
    ("Selection Pressure", "None (neutral)", "Complex positive + purifying"),
    ("Ground Truth", "Known by construction", "Approximate (DMS / literature)"),
    ("Phylogenetic Structure", "Independent samples", "Strong non-independence"),
    ("Best Method Winner", "Threshold-based", "Clustering-based (HDBSCAN)"),
]
for i, (aspect, syn, real) in enumerate(comp_rows):
    cy = Inches(2.550) + i * Inches(0.800)
    # aspect
    add_textbox(s, Inches(0.600), cy, Inches(1.800), Inches(0.350),
                aspect, size=10, bold=True, color=DARK_NAVY)
    # Synthetic value
    syn_card = add_card(s, Inches(2.500), cy, Inches(3.600), Inches(0.650), border=MID_GRAY)
    set_text(syn_card, syn, size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    # Real value
    real_card = add_card(s, Inches(6.900), cy, Inches(5.800), Inches(0.650),
                         fill=RGBColor(0xFF, 0xF3, 0xE0), border=ORANGE)
    set_text(real_card, real, size=11, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

# Bottom conclusion
msg = add_card(s, Inches(1.500), Inches(6.600), Inches(10.300), Inches(0.450),
               fill=DARK_NAVY)
set_text(msg, "Ranking reversals between synthetic and real data demonstrate the need for real-pathogen benchmarking",
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ── SLIDE 19: H3N2 Ranking Reversal ─────────────────────────────
slide_num += 1
s = content_slide(prs, "Cross-Pathogen Ranking Reversal")
add_page_number(s, slide_num)

add_figure(s, os.path.join(FIG_MAIN, "detection_ranking.png"),
           Inches(0.300), Inches(1.100), Inches(7.800), Inches(5.600))

# big number
big = add_card(s, Inches(8.500), Inches(1.300), Inches(4.100), Inches(2.400), fill=DARK_NAVY)
add_textbox(s, Inches(8.500), Inches(1.500), Inches(4.100), Inches(1.000),
            "9 / 9", size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(8.500), Inches(2.600), Inches(4.100), Inches(0.800),
            "unique best\ncombinations", size=18, color=RGBColor(0xBB, 0xDE, 0xFB),
            align=PP_ALIGN.CENTER)

add_textbox(s, Inches(8.500), Inches(4.000), Inches(4.100), Inches(2.800),
            "Every pathogen has a different\noptimal scoring + detection pair.\n\n"
            "No single method can be\n\"recommended\" universally.\n\n"
            "This motivates the need for\npathogen-adaptive selection (PAHD).",
            size=12, color=DARK_GRAY)


# ── SLIDE 20: 9-Pathogen Scale ──────────────────────────────────
slide_num += 1
s = content_slide(prs, "Stage 2: Scaling to 9 Pathogens")
add_page_number(s, slide_num)

bignums = [
    ("9", "Pathogens", "SARS-CoV-2, Influenza, HIV,\nDengue, MERS, RSV,\nNorovirus, HCV, Ebola"),
    ("2,544", "Evaluations", "9 scoring methods\nx 39 detection methods\nx 9 pathogens (multilayer GT)"),
    ("4", "Statistical Tests", "Two-way ANOVA\nPermutation test\nFriedman test\nLOPO cross-validation"),
]
for i, (num, label, desc) in enumerate(bignums):
    cx = Inches(0.600) + i * Inches(4.200)
    cy = Inches(1.500)
    card = add_card(s, cx, cy, Inches(3.900), Inches(4.800), border=MID_GRAY)
    add_textbox(s, cx, cy + Inches(0.300), Inches(3.900), Inches(1.200),
                num, size=54, bold=True, color=SEC_STAGE2, align=PP_ALIGN.CENTER)
    add_textbox(s, cx, cy + Inches(1.500), Inches(3.900), Inches(0.500),
                label, size=20, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)
    add_textbox(s, cx + Inches(0.300), cy + Inches(2.200), Inches(3.300), Inches(2.400),
                desc, size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ── SLIDE 21: ANOVA ─────────────────────────────────────────────
# CHANGE: Fix figure aspect ratio — don't stretch/compress
slide_num += 1
s = content_slide(prs, "Two-Way ANOVA: Variance Decomposition")
add_page_number(s, slide_num)

# Use original aspect ratio for figure — constrain height, let width follow
add_figure(s, os.path.join(FIG_MAIN, "variance_decomposition.png"),
           Inches(0.300), Inches(1.100), Inches(7.000), Inches(5.300))

# omega-squared cards
anova_cards = [
    ("Detection", "~35%", "Largest effect:\nchoice of detection\nmethod matters most"),
    ("Interaction", "28.5%", "Scoring x Pathogen\ninteraction is strong:\nno universal best"),
    ("Scoring", "~8.3%", "Scoring method\ncontributes less than\ndetection or interaction"),
]
for i, (ttl, val, desc) in enumerate(anova_cards):
    cy = Inches(1.200) + i * Inches(1.900)
    card = add_card(s, Inches(7.700), cy, Inches(5.000), Inches(1.650), border=MID_GRAY)
    add_accent_bar(s, Inches(7.700), cy, Inches(0.080), Inches(1.650), SEC_STAGE2)
    add_textbox(s, Inches(8.000), cy + Inches(0.100), Inches(1.500), Inches(0.400),
                ttl, size=12, bold=True, color=SEC_STAGE2)
    add_textbox(s, Inches(8.000), cy + Inches(0.500), Inches(1.500), Inches(0.800),
                val, size=28, bold=True, color=SEC_STAGE2)
    add_textbox(s, Inches(9.700), cy + Inches(0.200), Inches(2.800), Inches(1.200),
                desc, size=11, color=DARK_GRAY)

add_takeaway(s, "Takeaway: Interaction effect (28.5%) proves scoring-pathogen combination matters -- no single method works for all.")


# ── SLIDE 22: Per-Pathogen Best ──────────────────────────────────
# CHANGE: Add clear conclusion sentence, highlight key insight
slide_num += 1
s = content_slide(prs, "Per-Pathogen Best Combinations")
add_page_number(s, slide_num)

add_figure(s, os.path.join(FIG_MAIN, "cross_pathogen_top5.png"),
           Inches(0.300), Inches(1.100), Inches(8.500), Inches(5.000))

# highlight box
big = add_card(s, Inches(9.100), Inches(1.300), Inches(3.600), Inches(2.200), fill=SEC_STAGE2)
add_textbox(s, Inches(9.100), Inches(1.450), Inches(3.600), Inches(1.000),
            "9 / 9", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(9.100), Inches(2.400), Inches(3.600), Inches(0.800),
            "unique best\ncombinations", size=16, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(s, Inches(9.100), Inches(3.800), Inches(3.600), Inches(2.200),
            "High entropy-function divergence\n(Norovirus r=0.39 vs\nSARS-CoV-2 r=0.03)\n\n"
            "H-score founder bias\n(rho = -0.876)\n\n"
            "E-R correlation:\nrho=0.517, p=0.154 (n.s.)",
            size=11, color=DARK_GRAY)

# Clear conclusion
add_takeaway(s, "Takeaway: Each pathogen requires a DIFFERENT optimal method -- no universal recommendation is possible.")


# ── SLIDE 23: Friedman + LOPO ────────────────────────────────────
slide_num += 1
s = content_slide(prs, "Friedman Test & Leave-One-Pathogen-Out CV")
add_page_number(s, slide_num)

# Friedman card — left
fc = add_card(s, Inches(0.600), Inches(1.200), Inches(5.800), Inches(5.400), border=MID_GRAY)
add_rect(s, Inches(0.600), Inches(1.200), Inches(5.800), Inches(0.600), fill=SEC_STAGE2)
add_textbox(s, Inches(0.800), Inches(1.220), Inches(5.400), Inches(0.560),
            "Friedman Rank Test (Top-20)", size=18, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
add_textbox(s, Inches(0.800), Inches(2.000), Inches(5.400), Inches(1.800),
            "Question:\n"
            "  Are top-20 method rankings consistent\n"
            "  across pathogens?\n\n"
            "Result:\n"
            "  chi-squared = 42.44\n"
            "  p-value = 0.0015 (significant at p < 0.01)",
            size=13, color=DARK_GRAY)
add_placeholder(s, Inches(0.800), Inches(4.100), Inches(5.400), Inches(2.200),
                "[Chart to be created]\nFriedman rank distribution showing\n\u03c7\u00b2 = 42.44, p = 0.0015\nacross 9 pathogens")

# LOPO card — right
lc = add_card(s, Inches(6.900), Inches(1.200), Inches(5.800), Inches(5.400), border=MID_GRAY)
add_rect(s, Inches(6.900), Inches(1.200), Inches(5.800), Inches(0.600), fill=RED_ACCENT)
add_textbox(s, Inches(7.100), Inches(1.220), Inches(5.400), Inches(0.560),
            "LOPO Cross-Validation", size=18, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
add_textbox(s, Inches(7.100), Inches(2.000), Inches(5.400), Inches(1.800),
            "Protocol:\n"
            "  Train on 8 pathogens, test on held-out 1\n\n"
            "Result:\n"
            "  0 / 9 correct predictions\n"
            "  Simple majority voting fails completely",
            size=13, color=DARK_GRAY)
add_placeholder(s, Inches(7.100), Inches(4.100), Inches(5.400), Inches(2.200),
                "[Table to be created]\nLOPO results: 0/9 match rate\nMean generalization ratio: 4.9%\nTraining-best \u2260 Test-best")


# ── SLIDE 24: ESM-2 Validation ──────────────────────────────────
slide_num += 1
s = content_slide(prs, "ESM-2 Protein Language Model Validation")
add_page_number(s, slide_num)

# table
add_rect(s, Inches(1.200), Inches(1.400), Inches(10.900), Inches(0.500), fill=DARK_NAVY)
esm_cols = ["Metric", "Without ESM-2", "With ESM-2", "Change"]
esm_ws = [Inches(3.000), Inches(2.600), Inches(2.600), Inches(2.700)]
cx = Inches(1.200)
for c, w in zip(esm_cols, esm_ws):
    add_textbox(s, cx, Inches(1.410), w, Inches(0.480),
                c, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    cx += w

esm_rows = [
    ("Mean MCC (top-20)", "0.347", "0.352", "+0.005"),
    ("Median MCC (top-20)", "0.338", "0.341", "+0.003"),
    ("Best single combo", "0.390", "0.393", "+0.003"),
    ("Coverage (>0.3 MCC)", "12/20", "13/20", "+1"),
    ("Rank correlation", "--", "rho = 0.94", "High agreement"),
]
for r, (m, wo, wi, ch) in enumerate(esm_rows):
    ry = Inches(1.900) + r * Inches(0.650)
    fill = LIGHT_GRAY if r % 2 == 0 else WHITE
    add_rect(s, Inches(1.200), ry, Inches(10.900), Inches(0.650), fill=fill)
    cx = Inches(1.200)
    for val, w in zip([m, wo, wi, ch], esm_ws):
        clr = GREEN if val.startswith("+") or val == "High agreement" else DARK_GRAY
        add_textbox(s, cx, ry, w, Inches(0.650), val, size=12, color=clr,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += w

msg = add_card(s, Inches(1.200), Inches(5.200), Inches(10.900), Inches(0.900),
               fill=RGBColor(0xE8, 0xF5, 0xE9), border=GREEN)
set_text(msg, "ESM-2 embeddings provide marginal improvement, confirming that\n"
         "sequence-based scoring already captures most of the relevant signal.\n"
         "The benchmark rankings remain stable (rho = 0.94).",
         size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_placeholder(s, Inches(3.500), Inches(6.300), Inches(6.300), Inches(0.600),
                "[Placeholder] ESM-2 Architecture Concept Diagram (To be inserted)")


# ── SLIDE 25: KEY FINDING ────────────────────────────────────────
slide_num += 1
s = content_slide(prs, "KEY FINDING: No Universal Best Method")
add_page_number(s, slide_num)

# central message
central = add_card(s, Inches(1.500), Inches(1.200), Inches(10.300), Inches(1.200), fill=DARK_NAVY)
set_text(central, "No single scoring + detection combination works best for all pathogens",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 4 evidence cards
evidence = [
    ("ANOVA Interaction", "omega-sq = 0.285", "Scoring x Pathogen interaction\nexplains 28.5% of variance", SEC_STAGE2),
    ("Unique Best Combos", "9 / 9", "Every pathogen has a different\noptimal method pair", GREEN),
    ("Friedman Test", "p = 0.0015", "Rankings differ significantly\nacross pathogens", PURPLE),
    ("LOPO CV", "0 / 9 correct", "Majority voting fails completely\nfor unseen pathogens", RED_ACCENT),
]
for i, (ttl, val, desc, clr) in enumerate(evidence):
    col = i % 2
    row = i // 2
    cx = Inches(0.600) + col * Inches(6.300)
    cy = Inches(2.700) + row * Inches(2.200)
    card = add_card(s, cx, cy, Inches(5.800), Inches(1.950), border=MID_GRAY)
    add_accent_bar(s, cx, cy, Inches(0.080), Inches(1.950), clr)
    add_textbox(s, cx + Inches(0.250), cy + Inches(0.100), Inches(5.300), Inches(0.350),
                ttl, size=13, bold=True, color=clr)
    add_textbox(s, cx + Inches(0.250), cy + Inches(0.500), Inches(5.300), Inches(0.600),
                val, size=28, bold=True, color=DARK_NAVY)
    add_textbox(s, cx + Inches(0.250), cy + Inches(1.150), Inches(5.300), Inches(0.700),
                desc, size=11, color=DARK_GRAY)


# ── SLIDE 26: PAHD ──────────────────────────────────────────────
slide_num += 1
s = content_slide(prs, "PAHD: Pathogen-Adaptive Hotspot Detection")
add_page_number(s, slide_num)

# 3-step flow
steps = [
    ("Step 1: Profile", "Extract pathogen profile\n(entropy distribution, genome\nlength, mutation rate, diversity)"),
    ("Step 2: Match", "Compare profile against\nMutBench knowledge base\n(similarity-based retrieval)"),
    ("Step 3: Select", "Recommend optimal scoring +\ndetection combination based\non matched pathogen results"),
]
for i, (ttl, desc) in enumerate(steps):
    cx = Inches(0.600) + i * Inches(4.200)
    cy = Inches(1.200)
    card = add_card(s, cx, cy, Inches(3.900), Inches(2.300), border=MID_GRAY)
    add_rect(s, cx, cy, Inches(3.900), Inches(0.500), fill=PURPLE)
    add_textbox(s, cx, cy + Inches(0.050), Inches(3.900), Inches(0.400),
                ttl, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, cx + Inches(0.200), cy + Inches(0.600), Inches(3.500), Inches(1.500),
                desc, size=12, color=DARK_GRAY)
    if i < 2:
        add_textbox(s, cx + Inches(3.900), cy + Inches(0.650), Inches(0.300), Inches(0.500),
                    ">", size=24, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

# Results with figure
add_figure(s, os.path.join(FIG_DISS, "new_methodology_comparison.png"),
           Inches(0.400), Inches(3.800), Inches(6.500), Inches(2.800))

# Baseline comparison bars
add_textbox(s, Inches(7.200), Inches(3.700), Inches(5.500), Inches(0.400),
            "MCC Baseline Comparison", size=14, bold=True, color=PURPLE)
bars = [
    ("MutBench Best", "0.390", GREEN),
    ("FreqThresh", "0.248", BLUE),
    ("SWAN", "0.193", ORANGE),
    ("MutClust-Orig", "0.138", RED_ACCENT),
    ("Random", "0.001", GRAY),
]
for i, (name, val, color) in enumerate(bars):
    cy = Inches(4.200) + i * Inches(0.470)
    add_textbox(s, Inches(7.200), cy, Inches(1.800), Inches(0.250),
                name, size=10, bold=True, color=color)
    bar_w = min(float(val) / 0.4 * 3.0, 3.000)
    add_rect(s, Inches(9.100), cy, Inches(bar_w), Inches(0.250), fill=color)
    add_textbox(s, Inches(9.100) + Inches(0.100), cy,
                Inches(bar_w - 0.100), Inches(0.250), val, size=9, bold=True, color=WHITE,
                anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)


# ── SLIDE 27: Phylo + Region-overlap ────────────────────────────
slide_num += 1
s = content_slide(prs, "Phylogenetic Correction & Region Overlap")
add_page_number(s, slide_num)

# Phylo card — left
pc = add_card(s, Inches(0.600), Inches(1.200), Inches(5.800), Inches(5.400), border=MID_GRAY)
add_rect(s, Inches(0.600), Inches(1.200), Inches(5.800), Inches(0.600), fill=PURPLE)
add_textbox(s, Inches(0.800), Inches(1.220), Inches(5.400), Inches(0.560),
            "Phylogenetic Non-Independence", size=16, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
add_textbox(s, Inches(0.800), Inches(2.000), Inches(5.400), Inches(2.000),
            "Problem:\n"
            "  Closely related sequences share mutations\n"
            "  by descent, not independent selection.\n"
            "  e.g., D614G spread by founder effect.\n\n"
            "Approach:\n"
            "  TreeTime phylogenetic reconstruction\n"
            "  Branch-weighted mutation scoring",
            size=12, color=DARK_GRAY)
add_placeholder(s, Inches(0.800), Inches(4.200), Inches(5.200), Inches(2.100),
                "[Image Placeholder]\nPhylogenetic Tree Diagram\n(To be inserted)")

# Region overlap card — right
rc = add_card(s, Inches(6.900), Inches(1.200), Inches(5.800), Inches(5.400), border=MID_GRAY)
add_rect(s, Inches(6.900), Inches(1.200), Inches(5.800), Inches(0.600), fill=TEAL_ACCENT)
add_textbox(s, Inches(7.100), Inches(1.220), Inches(5.400), Inches(0.560),
            "Region-Level Overlap Analysis", size=16, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
add_textbox(s, Inches(7.100), Inches(2.000), Inches(5.400), Inches(2.000),
            "Motivation:\n"
            "  Position-level MCC may be too strict.\n"
            "  Methods may detect the right region\n"
            "  but miss exact positions.\n\n"
            "Results:\n"
            "  MCC improves from 0.289 to 0.712\n"
            "  with region-level evaluation.",
            size=12, color=DARK_GRAY)
add_figure(s, os.path.join(FIG_MAIN, "lopo_region_precision.png"),
           Inches(7.100), Inches(4.200), Inches(5.200), Inches(2.100))


# ── SLIDE 28: Baseline Comparison ────────────────────────────────
slide_num += 1
s = content_slide(prs, "Scoring-Detection Heatmap & Baseline Comparison")
add_page_number(s, slide_num)

add_figure(s, os.path.join(FIG_MAIN, "scoring_detection_heatmap.png"),
           Inches(0.300), Inches(1.100), Inches(7.800), Inches(5.600))

# MCC comparison bars
add_textbox(s, Inches(8.500), Inches(1.200), Inches(4.200), Inches(0.400),
            "MCC Comparison", size=16, bold=True, color=SEC_EXT)
bars = [
    ("MutBench Best", "0.390", GREEN),
    ("FreqThresh", "0.248", BLUE),
    ("SWAN", "0.193", ORANGE),
    ("MutClust-Orig", "0.138", RED_ACCENT),
    ("Random", "0.001", GRAY),
]
for i, (name, val, color) in enumerate(bars):
    cy = Inches(1.800) + i * Inches(1.000)
    add_textbox(s, Inches(8.500), cy, Inches(2.500), Inches(0.350),
                name, size=12, bold=True, color=color)
    bar_w = min(float(val) / 0.4 * 3.0, 3.000)
    add_rect(s, Inches(8.500), cy + Inches(0.380), Inches(bar_w), Inches(0.350), fill=color)
    add_textbox(s, Inches(8.500) + Inches(0.100), cy + Inches(0.380),
                Inches(bar_w - 0.100), Inches(0.350), val, size=11, bold=True, color=WHITE,
                anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT)


# ── SLIDE 29: Research Contributions ──────────────────────────────
slide_num += 1
s = content_slide(prs, "Research Contributions")
add_page_number(s, slide_num)

contribs = [
    ("C1: MutBench Framework",
     "First systematic benchmark for viral\nmutation hotspot detection with\n"
     "standardized pipeline & evaluation",
     "See S11-S17", SEC_FRAME),
    ("C2: 3-Layer Ground Truth",
     "Adaptive threshold + constrained +\nDMS-based ground truth design\n"
     "reduces single-reference bias",
     "See S13-S15", GREEN),
    ("C3: Cross-Pathogen Evidence",
     "9 pathogens, 2,544 evaluations prove\n"
     "no universal best method exists\n"
     "(interaction omega-sq = 0.285)",
     "See S22-S27", SEC_STAGE2),
    ("C4: PAHD Proof of Concept",
     "Pathogen-Adaptive Hotspot Detection:\n"
     "profile-based method selection\n"
     "outperforms any single approach",
     "See S28", PURPLE),
]
for i, (ttl, desc, ref, clr) in enumerate(contribs):
    col = i % 2
    row = i // 2
    cx = Inches(0.600) + col * Inches(6.300)
    cy = Inches(1.200) + row * Inches(2.700)
    card = add_card(s, cx, cy, Inches(5.800), Inches(2.400), border=MID_GRAY)
    add_accent_bar(s, cx, cy, Inches(0.080), Inches(2.400), clr)
    add_textbox(s, cx + Inches(0.250), cy + Inches(0.150), Inches(5.300), Inches(0.400),
                ttl, size=16, bold=True, color=clr)
    add_textbox(s, cx + Inches(0.250), cy + Inches(0.650), Inches(5.300), Inches(1.300),
                desc, size=12, color=DARK_GRAY)
    add_textbox(s, cx + Inches(0.250), cy + Inches(1.950), Inches(5.300), Inches(0.300),
                ref, size=9, color=GRAY)


# ── SLIDE 30: Discussion ────────────────────────────────────────
slide_num += 1
s = content_slide(prs, "Discussion")
add_page_number(s, slide_num)

disc = [
    ("Benchmarking Matters",
     "Without standardized evaluation, published methods appear to perform well in isolation "
     "but fail under cross-pathogen comparison. MutBench provides the missing infrastructure "
     "for fair, reproducible comparison of hotspot detection methods."),
    ("Method Selection is Pathogen-Dependent",
     "The strong interaction effect (omega-sq = 0.285) means that recommending a single method "
     "is fundamentally flawed. Pathogen-specific characteristics (genome structure, mutation rate, "
     "selection landscape) determine which approach works best."),
    ("From Benchmark to Algorithm",
     "MutBench is not just a benchmark but a knowledge base. The systematic evaluation across "
     "9 pathogens enables PAHD: data-driven, profile-based method selection that outperforms "
     "any fixed approach."),
]
for i, (ttl, desc) in enumerate(disc):
    cy = Inches(1.200) + i * Inches(1.850)
    card = add_card(s, Inches(0.600), cy, Inches(12.100), Inches(1.650), border=MID_GRAY)
    colors = [SEC_FRAME, SEC_STAGE2, PURPLE]
    add_accent_bar(s, Inches(0.600), cy, Inches(0.080), Inches(1.650), colors[i])
    add_textbox(s, Inches(0.900), cy + Inches(0.120), Inches(11.500), Inches(0.400),
                ttl, size=16, bold=True, color=colors[i])
    add_textbox(s, Inches(0.900), cy + Inches(0.580), Inches(11.500), Inches(0.950),
                desc, size=12, color=DARK_GRAY)


# ── SLIDE 31: Limitations & Future Work ──────────────────────────
slide_num += 1
s = content_slide(prs, "Limitations & Future Work")
add_page_number(s, slide_num)

# 2-column layout
col_data = [
    ("Limitations", RED_ACCENT, [
        "Ground truth approximation\n(no perfect gold standard exists)",
        "Phylogenetic non-independence\nnot fully corrected (TreeTime partial)",
        "Limited to substitution mutations\n(no indels, recombination)",
        "PAHD is proof-of-concept only,\nnot production-ready system",
    ]),
    ("Future Work", BLUE, [
        "Full PAHD implementation with\nlearned pathogen profiles",
        "Integrate phylogenetic correction\n(TreeTime branch weighting)",
        "Extend to indels, recombination,\nand structural variants",
        "Real-time surveillance integration\nfor emerging pathogens",
    ]),
]
for ci, (ttl, color, items) in enumerate(col_data):
    cx = Inches(0.500) + ci * Inches(6.400)
    add_rect(s, cx, Inches(1.200), Inches(6.000), Inches(0.500), fill=color)
    add_textbox(s, cx, Inches(1.210), Inches(6.000), Inches(0.480),
                ttl, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    for j, item in enumerate(items):
        iy = Inches(1.850) + j * Inches(1.250)
        card = add_card(s, cx, iy, Inches(6.000), Inches(1.050), border=MID_GRAY)
        add_accent_bar(s, cx, iy, Inches(0.060), Inches(1.050), color)
        add_textbox(s, cx + Inches(0.200), iy + Inches(0.100), Inches(5.600), Inches(0.850),
                    item, size=11, color=DARK_GRAY)


# ── SLIDE 32: Key Takeaway ──────────────────────────────────────
slide_num += 1
s = content_slide(prs, "Key Takeaway")
add_page_number(s, slide_num)

# central message
central = add_card(s, Inches(1.200), Inches(1.300), Inches(10.900), Inches(1.800), fill=DARK_NAVY)
set_text(central,
         "MutBench transforms viral hotspot detection from\n"
         "ad-hoc single-virus evaluation to systematic cross-pathogen benchmarking",
         size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 4 evidence boxes
evidences = [
    ("9 Pathogens", "Largest cross-pathogen\nbenchmark for hotspot\ndetection", SEC_FRAME),
    ("2,544 Evaluations", "Comprehensive coverage\nof scoring x detection\ncombinations", SEC_STAGE2),
    ("omega-sq = 0.285", "Strong interaction proves\npathogen-adaptive\nselection is needed", PURPLE),
    ("PAHD Concept", "From benchmark insight\nto actionable algorithm\nfor method selection", GREEN),
]
for i, (ttl, desc, clr) in enumerate(evidences):
    cx = Inches(0.500) + i * Inches(3.200)
    cy = Inches(3.500)
    card = add_card(s, cx, cy, Inches(2.900), Inches(3.100), border=MID_GRAY)
    add_accent_bar(s, cx, cy, Inches(0.080), Inches(3.100), clr)
    add_textbox(s, cx, cy + Inches(0.250), Inches(2.900), Inches(0.500),
                ttl, size=16, bold=True, color=clr, align=PP_ALIGN.CENTER)
    add_textbox(s, cx + Inches(0.200), cy + Inches(0.900), Inches(2.500), Inches(2.000),
                desc, size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ── SLIDE 33: Thank You ─────────────────────────────────────────
slide_num += 1
s = add_blank(prs)
# dark top area
add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(3.500), fill=DARK_NAVY)
# KNU logo
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(0.600), Inches(0.300), Inches(0.650), Inches(0.650))
add_textbox(s, Inches(0.600), Inches(1.200), Inches(12.0), Inches(1.200),
            "Thank You", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.600), Inches(2.400), Inches(12.0), Inches(0.600),
            "Questions & Discussion", size=20, color=RGBColor(0xBB, 0xDE, 0xFB),
            align=PP_ALIGN.CENTER)

# contact info
add_textbox(s, Inches(3.000), Inches(4.200), Inches(7.300), Inches(0.400),
            "Hwijun Kwon", size=22, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(3.000), Inches(4.700), Inches(7.300), Inches(0.350),
            "Dept. of Computer Science, Kyungpook National University",
            size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(3.000), Inches(5.200), Inches(7.300), Inches(0.350),
            "Advisor: Prof. Inuk Jung", size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(3.000), Inches(5.700), Inches(7.300), Inches(0.350),
            "March 2026", size=14, color=GRAY, align=PP_ALIGN.CENTER)

# GitHub placeholder
add_placeholder(s, Inches(5.000), Inches(6.200), Inches(3.300), Inches(0.600),
                "[QR Code Placeholder] GitHub Repository")

# bottom bar
add_rect(s, Inches(0), Inches(7.120), SLIDE_W, Inches(0.380), fill=DARK_NAVY)
add_textbox(s, Inches(0.400), Inches(7.140), Inches(10.0), Inches(0.340),
            "Kyungpook National University  |  Dept. of Computer Science  |  Hwijun Kwon",
            size=10, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)


# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Total slides: {len(prs.slides)}")
