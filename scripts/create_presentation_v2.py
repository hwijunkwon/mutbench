#!/usr/bin/env python3
"""Create MutBench PhD Defense PPTX presentation (v2)."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Constants ──────────────────────────────────────────────────────────────
SLIDE_W = Inches(13.330)
SLIDE_H = Inches(7.500)

SECTION_COLORS = {
    1: RGBColor(0x1B, 0x2A, 0x4A),
    2: RGBColor(0x2E, 0x50, 0x90),
    3: RGBColor(0x1A, 0x7A, 0x4C),
    4: RGBColor(0x8B, 0x45, 0x13),
    5: RGBColor(0x6B, 0x3F, 0xA0),
    6: RGBColor(0x1B, 0x2A, 0x4A),
}

BOTTOM_BAR_COLOR = RGBColor(0x1B, 0x2A, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x88, 0x88, 0x88)
LIGHT_GRAY_BG = RGBColor(0xF5, 0xF7, 0xFA)

FIG_DIR_D = "/proj/paper/paper/dissertation/figures"
FIG_DIR_F = "/proj/paper/paper/figure"

FIGURES = {
    "multi_ground_truth": f"{FIG_DIR_D}/multi_ground_truth_figure.png",
    "new_methodology": f"{FIG_DIR_D}/new_methodology_comparison.png",
    "sensitivity_heatmap": f"{FIG_DIR_D}/sensitivity_heatmap.png",
    "mcc_vs_f1": f"{FIG_DIR_F}/mcc_vs_f1_fairness.png",
    "detection_ranking": f"{FIG_DIR_F}/detection_ranking.png",
    "variance_decomposition": f"{FIG_DIR_F}/variance_decomposition.png",
    "cross_pathogen_top5": f"{FIG_DIR_F}/cross_pathogen_top5.png",
    "scoring_detection_heatmap": f"{FIG_DIR_F}/scoring_detection_heatmap.png",
}

FOOTER_TEXT = "MutBench: Systematic Benchmarking for Viral Mutation Hotspot Detection | 경북대학교 대학원"


# ── Helper functions ───────────────────────────────────────────────────────
def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_shape(slide, left, top, width, height, fill=None, border_color=None,
              border_width=None, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border_color:
        shape.line.fill.solid()
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    return shape


def set_text(shape, text, size=14, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font_name="맑은 고딕"):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = alignment
    except:
        pass
    # Set anchor
    try:
        txBody = shape._element.txBody
        bodyPr = txBody.find(qn("a:bodyPr"))
        if bodyPr is not None:
            bodyPr.set("anchor", {
                MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"
            }.get(anchor, "t"))
    except:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_text_box(slide, left, top, width, height, text, size=14, bold=False,
                 color=BLACK, alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    set_text(txBox, text, size, bold, color, alignment, anchor, font_name)
    return txBox


def add_multiline_text(shape, lines, size=14, color=BLACK, bold=False,
                       alignment=PP_ALIGN.LEFT, line_spacing=1.2, font_name="맑은 고딕"):
    """Set multiple lines in a shape's text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(size * 0.3)
    return tf


def add_bullet_text(shape, lines, size=13, color=BLACK, font_name="맑은 고딕"):
    """Add bulleted lines to a shape."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        p.level = 0
    return tf


def add_top_bar(slide, section):
    color = SECTION_COLORS[section]
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.900), fill=color)


def add_title(slide, text):
    add_text_box(slide, Inches(0.600), Inches(0.180), Inches(10.500), Inches(0.550),
                 text, size=28, bold=True, color=WHITE)


def add_knu_logo(slide):
    circle = add_shape(slide, Inches(12.230), Inches(0.150), Inches(0.550), Inches(0.550),
                       fill=WHITE, shape_type=MSO_SHAPE.OVAL)
    set_text(circle, "KNU", size=9, bold=True, color=SECTION_COLORS[1],
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_bottom_bar(slide, page_num):
    add_shape(slide, Inches(0), Inches(7.120), SLIDE_W, Inches(0.380), fill=BOTTOM_BAR_COLOR)
    add_text_box(slide, Inches(0.600), Inches(7.150), Inches(10.000), Inches(0.300),
                 FOOTER_TEXT, size=10, color=WHITE)
    add_text_box(slide, Inches(12.130), Inches(7.140), Inches(0.900), Inches(0.300),
                 str(page_num), size=10, color=WHITE, alignment=PP_ALIGN.RIGHT)


def add_common(slide, section, title_text, page_num):
    add_top_bar(slide, section)
    add_title(slide, title_text)
    add_knu_logo(slide)
    add_bottom_bar(slide, page_num)


def add_card(slide, left, top, width, height, title, bullets, header_color,
             body_fill=WHITE, title_size=16, bullet_size=13):
    """Add a card with colored header and bullet body."""
    # Header
    hdr_h = Inches(0.450)
    hdr = add_shape(slide, left, top, width, hdr_h, fill=header_color)
    set_text(hdr, title, size=title_size, bold=True, color=WHITE,
             alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    # Indent text
    tf = hdr.text_frame
    tf.margin_left = Inches(0.200)

    # Body
    body_h = height - hdr_h
    body = add_shape(slide, left, top + hdr_h, width, body_h, fill=body_fill,
                     border_color=RGBColor(0xDD, 0xDD, 0xDD), border_width=1)
    tf = body.text_frame
    tf.margin_left = Inches(0.200)
    tf.margin_right = Inches(0.200)
    tf.margin_top = Inches(0.100)
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(bullet_size)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "맑은 고딕"
        p.space_after = Pt(4)
    return body


def add_simple_card(slide, left, top, width, height, fill, border_left_color=None,
                    border_width=4):
    """Add a simple rectangular card, optionally with left border accent."""
    card = add_shape(slide, left, top, width, height, fill=fill,
                     border_color=RGBColor(0xE0, 0xE0, 0xE0), border_width=1)
    if border_left_color:
        accent = add_shape(slide, left, top, Pt(border_width), height,
                           fill=border_left_color)
    return card


def add_figure(slide, key, left, top, width, height):
    path = FIGURES.get(key)
    if path and os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)
    else:
        box = add_shape(slide, left, top, width, height, fill=LIGHT_GRAY_BG,
                        border_color=MED_GRAY, border_width=1)
        set_text(box, f"[Figure: {key}]", size=14, color=MED_GRAY,
                 alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_big_number(slide, left, top, width, height, number, label, color,
                   bg_fill=None):
    """Add a card with a big number and label below."""
    card = add_shape(slide, left, top, width, height,
                     fill=bg_fill or LIGHT_GRAY_BG,
                     border_color=color, border_width=2)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    # Number
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER
    # Label
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(12)
    p2.font.color.rgb = DARK_GRAY
    p2.font.name = "맑은 고딕"
    p2.alignment = PP_ALIGN.CENTER


def new_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


# ── Slide builders ─────────────────────────────────────────────────────────
def slide_01_title(prs):
    slide = new_slide(prs)
    # Dark top area
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(3.000),
              fill=SECTION_COLORS[1])
    # Title
    add_text_box(slide, Inches(0.600), Inches(0.700), Inches(11.000), Inches(1.200),
                 "MutBench: 체계적 벤치마킹을 통한\n바이러스 변이 핫스팟 탐지 방법론 비교 연구",
                 size=30, bold=True, color=WHITE)
    # Subtitle
    add_text_box(slide, Inches(0.600), Inches(1.950), Inches(11.000), Inches(0.500),
                 "Systematic Benchmarking Framework for Viral Mutation Hotspot Detection",
                 size=18, color=hex_to_rgb("B0C4DE"))
    # KNU logo on title
    circle = add_shape(slide, Inches(12.230), Inches(0.150), Inches(0.550), Inches(0.550),
                       fill=WHITE, shape_type=MSO_SHAPE.OVAL)
    set_text(circle, "KNU", size=9, bold=True, color=SECTION_COLORS[1],
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Author
    add_text_box(slide, Inches(0.600), Inches(3.400), Inches(10.000), Inches(0.350),
                 "발표자: 권휘준  |  지도교수: 정인욱", size=16, color=DARK_GRAY)
    # Dept
    add_text_box(slide, Inches(0.600), Inches(3.800), Inches(10.000), Inches(0.350),
                 "경북대학교 대학원 컴퓨터학부  |  2026년 3월", size=14, color=MED_GRAY)
    # Bottom bar
    add_shape(slide, Inches(0), Inches(7.120), SLIDE_W, Inches(0.380), fill=BOTTOM_BAR_COLOR)


def slide_02_toc(prs):
    slide = new_slide(prs)
    add_common(slide, 1, "목차 (Table of Contents)", 2)
    sections = [
        ("1. 서론", "연구 배경 및 문제 정의", "#1B2A4A"),
        ("2. MutBench 프레임워크", "파이프라인, GT, 지표", "#2E5090"),
        ("3. Stage 1 결과", "합성 데이터 벤치마크", "#1A7A4C"),
        ("4. Stage 2 결과", "9-병원체 실제 데이터", "#8B4513"),
        ("5. 확장 분석", "PAHD, 계통 보정, 비교", "#6B3FA0"),
        ("6. 결론", "기여, 한계, 향후 과제", "#1B2A4A"),
    ]
    positions = [
        (0.600, 1.200), (6.880, 1.200),
        (0.600, 2.800), (6.880, 2.800),
        (0.600, 4.400), (6.880, 4.400),
    ]
    for (title, desc, col_hex), (lx, ty) in zip(sections, positions):
        col = hex_to_rgb(col_hex)
        card = add_simple_card(slide, Inches(lx), Inches(ty), Inches(5.850), Inches(1.400),
                               fill=LIGHT_GRAY_BG, border_left_color=col)
        tf = card.text_frame
        tf.margin_left = Inches(0.350)
        tf.margin_top = Inches(0.200)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = col
        p.font.name = "맑은 고딕"
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = MED_GRAY
        p2.font.name = "맑은 고딕"


def slide_03_prior(prs):
    slide = new_slide(prs)
    add_common(slide, 1, "선행 연구: MOSD & MutClust", 3)
    # MOSD card
    add_card(slide, Inches(0.600), Inches(1.200), Inches(5.850), Inches(4.000),
             "MOSD", [
                 "Multi-omics 벤치마크 프레임워크",
                 "10개 방법 × 10개 암종 비교",
                 "Cluster-score 복합 지표 제안",
                 "핵심: 데이터 특성별 최적 방법 상이",
             ], SECTION_COLORS[1])
    # MutClust card
    add_card(slide, Inches(6.880), Inches(1.200), Inches(5.850), Inches(4.000),
             "MutClust", [
                 "H-score 기반 핫스팟 탐지",
                 "적응적 DBSCAN 클러스터링",
                 "SARS-CoV-2 477개 핫스팟 탐지",
                 "한계: 독립 GT 부재, 파라미터 고정",
             ], SECTION_COLORS[2])
    # Bridge box
    bridge = add_shape(slide, Inches(0.600), Inches(5.600), Inches(12.130), Inches(0.800),
                       fill=hex_to_rgb("FFF3CD"), border_color=hex_to_rgb("F0C674"), border_width=1)
    set_text(bridge, "→ 두 연구의 공통 교훈: 체계적 벤치마크 없이는 방법론 선택 불가",
             size=16, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)


def slide_04_hotspots(prs):
    slide = new_slide(prs)
    add_common(slide, 1, "변이 핫스팟(Mutation Hotspot)이란?", 4)
    # Definition box
    defbox = add_shape(slide, Inches(0.600), Inches(1.100), Inches(12.130), Inches(0.900),
                       fill=hex_to_rgb("E8F4FD"), border_color=hex_to_rgb("2E5090"), border_width=2)
    set_text(defbox,
             "정의: 바이러스 게놈에서 통계적으로 유의하게 높은 빈도로 변이가 발생하는 위치",
             size=18, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    tf = defbox.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)

    # 3 cards
    cards = [
        ("왜 발생?", hex_to_rgb("2E5090"),
         ["양성 선택 (positive selection)", "기능적 제약 (functional constraint)", "재조합 (recombination)"]),
        ("왜 중요?", hex_to_rgb("1A7A4C"),
         ["백신 표적 설계", "항바이러스제 개발", "변이 감시 및 조기 경보"]),
        ("현재 문제", hex_to_rgb("C0392B"),
         ["표준 벤치마크 부재", "병원체 의존성 미검증", "Ground Truth 불일치"]),
    ]
    col_lefts = [0.600, 4.740, 8.880]
    for (title, col, bullets), lx in zip(cards, col_lefts):
        add_card(slide, Inches(lx), Inches(2.300), Inches(3.850), Inches(3.500),
                 title, bullets, col)


def slide_05_problem(prs):
    slide = new_slide(prs)
    add_common(slide, 1, "연구 문제 정의", 5)
    problems = [
        ("P1", "표준 벤치마크의 부재",
         "기존 핫스팟 탐지 방법들은 각자의 데이터셋에서만 평가 → 공정한 비교 불가",
         hex_to_rgb("E74C3C")),
        ("P2", "병원체 특이성 미검증",
         "SARS-CoV-2에서 잘 작동하는 방법이 다른 병원체에도 유효한지 검증되지 않음",
         hex_to_rgb("F39C12")),
        ("P3", "Ground Truth 기준 불일치",
         "연구마다 다른 GT 정의 사용 → 성능 비교의 신뢰성 저하",
         hex_to_rgb("27AE60")),
    ]
    for i, (label, title, desc, col) in enumerate(problems):
        ty = 1.200 + i * 1.700
        card = add_simple_card(slide, Inches(0.600), Inches(ty), Inches(12.130), Inches(1.500),
                               fill=WHITE, border_left_color=col, border_width=6)
        tf = card.text_frame
        tf.margin_left = Inches(0.500)
        tf.margin_top = Inches(0.200)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{label}: {title}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = col
        p.font.name = "맑은 고딕"
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = DARK_GRAY
        p2.font.name = "맑은 고딕"
        p2.space_before = Pt(6)

    # Bottom arrow
    add_text_box(slide, Inches(0.600), Inches(6.400), Inches(12.130), Inches(0.500),
                 "→ MutBench로 해결", size=18, bold=True, color=hex_to_rgb("2E5090"),
                 alignment=PP_ALIGN.CENTER)


def slide_06_contributions(prs):
    slide = new_slide(prs)
    add_common(slide, 2, "연구 기여 (Contributions)", 6)
    contribs = [
        ("C1: MutBench 프레임워크", hex_to_rgb("2E5090"), hex_to_rgb("E8F4FD"),
         ["9개 병원체 × 9 scoring × 39 detection", "2,544 체계적 평가 실행", "재현 가능한 파이프라인"]),
        ("C2: 3-Layer Ground Truth", hex_to_rgb("1A7A4C"), hex_to_rgb("E8F9EF"),
         ["Adaptive + Constrained + DMS", "다각적 검증으로 편향 최소화", "기존 단일 GT의 한계 극복"]),
        ("C3: 병원체 특이성 실증", hex_to_rgb("E67E22"), hex_to_rgb("FFF5E6"),
         ["interaction ω² = 0.285", "9/9 고유 최적 조합", "단일 방법 불가 실증"]),
        ("C4: PAHD 적응형 탐지", hex_to_rgb("6B3FA0"), hex_to_rgb("F5E6FF"),
         ["병원체 프로파일 기반 선택", "Top-3 평균 MCC 0.390", "기존 최고 대비 +57% 향상"]),
    ]
    positions = [(0.600, 1.200), (6.880, 1.200), (0.600, 3.800), (6.880, 3.800)]
    for (title, col, bg, bullets), (lx, ty) in zip(contribs, positions):
        add_card(slide, Inches(lx), Inches(ty), Inches(5.850), Inches(2.400),
                 title, bullets, col, body_fill=bg)


def slide_07_pipeline(prs):
    slide = new_slide(prs)
    add_common(slide, 2, "MutBench 파이프라인", 7)
    # 5 flow boxes
    steps = [
        ("입력\n(MSA)", hex_to_rgb("3498DB")),
        ("Scoring\n(9종)", hex_to_rgb("2E5090")),
        ("Detection\n(39종)", hex_to_rgb("1A7A4C")),
        ("GT\n(3-Layer)", hex_to_rgb("E67E22")),
        ("평가\n(MCC)", hex_to_rgb("8B4513")),
    ]
    box_w = Inches(2.100)
    box_h = Inches(1.500)
    gap = Inches(0.300)
    start_x = Inches(0.600)
    ty = Inches(1.400)
    for i, (label, col) in enumerate(steps):
        lx = start_x + i * (box_w + gap)
        box = add_shape(slide, lx, ty, box_w, box_h, fill=col,
                        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        set_text(box, label, size=16, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Arrow between boxes
        if i < 4:
            arrow_x = lx + box_w
            add_text_box(slide, arrow_x, Inches(1.900), gap, Inches(0.400),
                         "→", size=24, bold=True, color=DARK_GRAY,
                         alignment=PP_ALIGN.CENTER)

    # Scoring list
    scoring_box = add_shape(slide, Inches(0.600), Inches(3.200), Inches(5.850), Inches(3.400),
                            fill=LIGHT_GRAY_BG, border_color=hex_to_rgb("2E5090"), border_width=1)
    tf = scoring_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.200)
    tf.margin_top = Inches(0.150)
    lines = [
        "Scoring Methods (9종):",
        "  Entropy, Wavelet, H-score,",
        "  KL-divergence, Jensen-Shannon,",
        "  Mutual Information, Chi-squared,",
        "  Conservation Score, Rate4Site"
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.bold = (i == 0)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "맑은 고딕"

    # Detection list
    det_box = add_shape(slide, Inches(6.880), Inches(3.200), Inches(5.850), Inches(3.400),
                        fill=LIGHT_GRAY_BG, border_color=hex_to_rgb("1A7A4C"), border_width=1)
    tf = det_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.200)
    tf.margin_top = Inches(0.150)
    lines = [
        "Detection Families (39종):",
        "  Threshold: Fixed, Percentile, MAD, IQR",
        "  Clustering: DBSCAN, HDBSCAN, OPTICS",
        "  Statistical: Z-score, Grubbs, Peaks",
        "  Sliding Window: Mean, Median, Gaussian",
        "  ML: Isolation Forest, LOF, One-Class SVM"
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.bold = (i == 0)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "맑은 고딕"


def slide_08_gt(prs):
    slide = new_slide(prs)
    add_common(slide, 2, "3-Layer Ground Truth 체계", 8)
    # Figure
    add_figure(slide, "multi_ground_truth", Inches(0.600), Inches(1.100),
               Inches(7.000), Inches(5.000))
    # Right panel - 3 layer cards
    layers = [
        ("Layer 1: Adaptive", hex_to_rgb("2E5090"), hex_to_rgb("E8F4FD"),
         "병원체별 최적 entropy 임계값 자동 결정"),
        ("Layer 2: Constrained", hex_to_rgb("1A7A4C"), hex_to_rgb("E8F9EF"),
         "기존 문헌 기반 보수적 핫스팟 정의"),
        ("Layer 3: DMS", hex_to_rgb("E67E22"), hex_to_rgb("FFF5E6"),
         "Deep Mutational Scanning 실험 데이터"),
    ]
    for i, (title, col, bg, desc) in enumerate(layers):
        ty = 1.100 + i * 1.700
        card = add_shape(slide, Inches(7.900), Inches(ty), Inches(4.830), Inches(1.500),
                         fill=bg, border_color=col, border_width=2)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.200)
        tf.margin_top = Inches(0.150)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = col
        p.font.name = "맑은 고딕"
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = DARK_GRAY
        p2.font.name = "맑은 고딕"
        p2.space_before = Pt(6)


def slide_09_mcc(prs):
    slide = new_slide(prs)
    add_common(slide, 2, "평가 지표: MCC", 9)
    # Left: formula + explanation
    left_box = add_shape(slide, Inches(0.600), Inches(1.200), Inches(5.850), Inches(5.000),
                         fill=LIGHT_GRAY_BG, border_color=hex_to_rgb("2E5090"), border_width=1)
    tf = left_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.300)
    tf.margin_top = Inches(0.300)
    tf.margin_right = Inches(0.300)
    lines = [
        ("MCC (Matthews Correlation Coefficient)", 18, True, hex_to_rgb("2E5090")),
        ("", 8, False, DARK_GRAY),
        ("MCC = (TP×TN - FP×FN) /", 16, False, DARK_GRAY),
        ("  √((TP+FP)(TP+FN)(TN+FP)(TN+FN))", 16, False, DARK_GRAY),
        ("", 8, False, DARK_GRAY),
        ("범위: [-1, +1]", 14, False, DARK_GRAY),
        ("  +1 = 완벽한 예측", 14, False, DARK_GRAY),
        ("   0 = 무작위 수준", 14, False, DARK_GRAY),
        ("  -1 = 완전한 역예측", 14, False, DARK_GRAY),
        ("", 8, False, DARK_GRAY),
        ("선택 이유:", 16, True, hex_to_rgb("2E5090")),
        ("• 클래스 불균형에 강건", 14, False, DARK_GRAY),
        ("• TP/TN/FP/FN 모두 반영", 14, False, DARK_GRAY),
        ("• F1 대비 낙관적 편향 없음", 14, False, DARK_GRAY),
        ("• Random baseline MCC ≈ 0.001", 14, False, DARK_GRAY),
    ]
    for i, (text, sz, bld, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bld
        p.font.color.rgb = col
        p.font.name = "맑은 고딕"

    # Right: figure
    add_figure(slide, "mcc_vs_f1", Inches(6.880), Inches(1.200),
               Inches(5.850), Inches(4.000))


def slide_10_twostage(prs):
    slide = new_slide(prs)
    add_common(slide, 2, "2단계 실험 설계", 10)
    # Stage 1
    add_card(slide, Inches(0.600), Inches(1.200), Inches(5.850), Inches(5.200),
             "Stage 1: 합성 데이터", [
                 "Controlled synthetic mutations",
                 "4 병원체 (SARS-CoV-2, HIV, Influenza, Dengue)",
                 "Known ground truth positions",
                 "방법론 기본 성능 검증",
                 "파라미터 민감도 분석",
             ], hex_to_rgb("1A7A4C"))
    # Stage 2
    add_card(slide, Inches(6.880), Inches(1.200), Inches(5.850), Inches(5.200),
             "Stage 2: 실제 데이터", [
                 "Real MSA from NCBI/GISAID",
                 "9 병원체 확장",
                 "3-Layer Ground Truth 적용",
                 "9 scoring × 39 detection = 351 조합",
                 "2,544 evaluations (9 pathogens)",
                 "4종 통계 검정 (ANOVA, Friedman, LOPO, Permutation)",
             ], hex_to_rgb("8B4513"))


def slide_11_stage1(prs):
    slide = new_slide(prs)
    add_common(slide, 3, "Stage 1: 합성 데이터 벤치마크", 11)
    # Figure
    add_figure(slide, "new_methodology", Inches(0.600), Inches(1.100),
               Inches(7.500), Inches(5.200))
    # Right summary
    summary = add_shape(slide, Inches(8.400), Inches(1.100), Inches(4.330), Inches(5.200),
                        fill=LIGHT_GRAY_BG, border_color=hex_to_rgb("1A7A4C"), border_width=2)
    tf = summary.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.200)
    tf.margin_top = Inches(0.200)
    lines = [
        ("주요 발견", 18, True, hex_to_rgb("1A7A4C")),
        ("", 6, False, DARK_GRAY),
        ("• Entropy + Percentile 조합이", 13, False, DARK_GRAY),
        ("  합성 데이터에서 최고 성능", 13, False, DARK_GRAY),
        ("", 6, False, DARK_GRAY),
        ("• Wavelet scoring의 높은", 13, False, DARK_GRAY),
        ("  노이즈 강건성 확인", 13, False, DARK_GRAY),
        ("", 6, False, DARK_GRAY),
        ("• HDBSCAN이 DBSCAN 대비", 13, False, DARK_GRAY),
        ("  파라미터 민감도 낮음", 13, False, DARK_GRAY),
        ("", 6, False, DARK_GRAY),
        ("• 합성 데이터 MCC 범위:", 13, False, DARK_GRAY),
        ("  0.15 ~ 0.85", 16, True, hex_to_rgb("1A7A4C")),
    ]
    for i, (text, sz, bld, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bld
        p.font.color.rgb = col
        p.font.name = "맑은 고딕"


def slide_12_sensitivity(prs):
    slide = new_slide(prs)
    add_common(slide, 3, "파라미터 민감도 분석", 12)
    add_figure(slide, "sensitivity_heatmap", Inches(0.600), Inches(1.100),
               Inches(7.000), Inches(5.000))
    # Right findings
    findings = add_shape(slide, Inches(7.900), Inches(1.100), Inches(4.830), Inches(5.000),
                         fill=LIGHT_GRAY_BG, border_color=hex_to_rgb("1A7A4C"), border_width=2)
    tf = findings.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.200)
    tf.margin_top = Inches(0.200)
    lines = [
        ("민감도 분석 결과", 18, True, hex_to_rgb("1A7A4C")),
        ("", 6, False, DARK_GRAY),
        ("• Threshold 기반 방법:", 14, True, DARK_GRAY),
        ("  임계값에 매우 민감", 13, False, DARK_GRAY),
        ("  (MCC 변동 ±0.3)", 13, False, DARK_GRAY),
        ("", 6, False, DARK_GRAY),
        ("• Clustering 기반 방법:", 14, True, DARK_GRAY),
        ("  eps 파라미터에 민감", 13, False, DARK_GRAY),
        ("  HDBSCAN은 상대적 안정", 13, False, DARK_GRAY),
        ("", 6, False, DARK_GRAY),
        ("• Statistical 방법:", 14, True, DARK_GRAY),
        ("  Z-score 기준 안정적", 13, False, DARK_GRAY),
        ("", 6, False, DARK_GRAY),
        ("시사점: 파라미터 선택이", 14, True, hex_to_rgb("E74C3C")),
        ("방법론 자체보다 중요할 수 있음", 13, False, hex_to_rgb("E74C3C")),
    ]
    for i, (text, sz, bld, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bld
        p.font.color.rgb = col
        p.font.name = "맑은 고딕"


def slide_13_gap(prs):
    slide = new_slide(prs)
    add_common(slide, 3, "Synthetic-Real 성능 격차", 13)
    # Comparison table header
    table_data = [
        ["", "Synthetic MCC", "Real MCC", "격차"],
        ["Entropy+Percentile", "0.85", "0.42", "-0.43"],
        ["Wavelet+HDBSCAN", "0.72", "0.38", "-0.34"],
        ["H-score+DBSCAN", "0.68", "0.25", "-0.43"],
        ["Conservation+Z-score", "0.61", "0.35", "-0.26"],
    ]
    rows, cols = len(table_data), len(table_data[0])
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.600), Inches(1.300),
                                         Inches(12.130), Inches(3.000))
    table = table_shape.table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = table_data[r][c]
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.name = "맑은 고딕"
            p.alignment = PP_ALIGN.CENTER
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = hex_to_rgb("1A7A4C")
            else:
                p.font.color.rgb = DARK_GRAY
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else LIGHT_GRAY_BG

    # Key message
    msg_box = add_shape(slide, Inches(0.600), Inches(4.800), Inches(12.130), Inches(1.500),
                        fill=hex_to_rgb("FFF3CD"), border_color=hex_to_rgb("F0C674"), border_width=2)
    tf = msg_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.300)
    tf.margin_top = Inches(0.200)
    lines = [
        ("핵심 시사점", 18, True, hex_to_rgb("8B4513")),
        ("합성 데이터에서의 우수한 성능이 실제 데이터로 전이되지 않음", 14, False, DARK_GRAY),
        ("→ 실제 MSA 기반 벤치마크(Stage 2)의 필요성 확인", 14, True, hex_to_rgb("2E5090")),
    ]
    for i, (text, sz, bld, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bld
        p.font.color.rgb = col
        p.font.name = "맑은 고딕"


def slide_14_ranking(prs):
    slide = new_slide(prs)
    add_common(slide, 3, "병원체 간 순위 역전", 14)
    add_figure(slide, "detection_ranking", Inches(0.600), Inches(1.100),
               Inches(7.500), Inches(5.000))
    # Right panel
    right = add_shape(slide, Inches(8.400), Inches(1.100), Inches(4.330), Inches(2.500),
                      fill=hex_to_rgb("E8F4FD"), border_color=hex_to_rgb("1A7A4C"), border_width=2)
    tf = right.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.200)
    tf.margin_top = Inches(0.150)
    p = tf.paragraphs[0]
    p.text = "9/9"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb("E74C3C")
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "모든 병원체에서\n서로 다른 최적 조합"
    p2.font.size = Pt(16)
    p2.font.color.rgb = DARK_GRAY
    p2.font.name = "맑은 고딕"
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)

    # Interpretation box
    interp = add_shape(slide, Inches(8.400), Inches(3.900), Inches(4.330), Inches(2.200),
                       fill=LIGHT_GRAY_BG, border_color=hex_to_rgb("1A7A4C"), border_width=1)
    tf2 = interp.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Inches(0.200)
    tf2.margin_top = Inches(0.150)
    lines = [
        "해석:",
        "• 한 병원체에서 1위인 방법이",
        "  다른 병원체에서 하위권",
        "• 순위 역전 = 병원체 특이성",
        "• 범용 최적 방법 부재 실증",
    ]
    for i, text in enumerate(lines):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = text
        p.font.size = Pt(13)
        p.font.bold = (i == 0)
        p.font.color.rgb = DARK_GRAY
        p.font.name = "맑은 고딕"


def slide_15_9pathogen(prs):
    slide = new_slide(prs)
    add_common(slide, 4, "Stage 2: 9-병원체 벤치마크", 15)
    # 3 big number cards
    numbers = [
        ("9", "병원체", "SARS-CoV-2, HIV, Influenza,\nDengue, HCV, RSV,\nNorovirus, MERS, H3N2",
         hex_to_rgb("2E5090")),
        ("2,544", "평가 조합", "9 scoring × 39 detection\n× 9 pathogens\n(multilayer GT 기반)",
         hex_to_rgb("1A7A4C")),
        ("4", "통계 검정", "ANOVA (ω²)\nFriedman (χ²)\nLOPO CV\nPermutation test",
         hex_to_rgb("8B4513")),
    ]
    col_lefts = [0.600, 4.740, 8.880]
    for (num, label, desc, col), lx in zip(numbers, col_lefts):
        card = add_shape(slide, Inches(lx), Inches(1.400), Inches(3.850), Inches(5.000),
                         fill=LIGHT_GRAY_BG, border_color=col, border_width=3)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.200)
        tf.margin_right = Inches(0.200)
        tf.margin_top = Inches(0.500)
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = col
        p.font.name = "맑은 고딕"
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = DARK_GRAY
        p2.font.name = "맑은 고딕"
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(12)
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(13)
        p3.font.color.rgb = MED_GRAY
        p3.font.name = "맑은 고딕"
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(16)


def slide_16_anova(prs):
    slide = new_slide(prs)
    add_common(slide, 4, "ANOVA: 분산 분해", 16)
    add_figure(slide, "variance_decomposition", Inches(0.600), Inches(1.100),
               Inches(7.500), Inches(5.000))
    # Right: 3 omega cards
    omegas = [
        ("ω² = 35%", "Detection", "탐지 방법 선택이\n가장 큰 영향", hex_to_rgb("2E5090")),
        ("ω² = 28.5%", "Pathogen×Method", "병원체-방법 상호작용\n(핵심 발견)", hex_to_rgb("E74C3C")),
        ("ω² = 8.3%", "Scoring", "점수화 방법의\n상대적 낮은 영향", hex_to_rgb("1A7A4C")),
    ]
    for i, (val, label, desc, col) in enumerate(omegas):
        ty = 1.100 + i * 1.800
        card = add_shape(slide, Inches(8.400), Inches(ty), Inches(4.330), Inches(1.600),
                         fill=LIGHT_GRAY_BG, border_color=col, border_width=2)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.200)
        tf.margin_top = Inches(0.150)
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = col
        p.font.name = "맑은 고딕"
        p2 = tf.add_paragraph()
        p2.text = f"{label}: {desc}"
        p2.font.size = Pt(12)
        p2.font.color.rgb = DARK_GRAY
        p2.font.name = "맑은 고딕"
        p2.space_before = Pt(4)


def slide_17_perpathogen(prs):
    slide = new_slide(prs)
    add_common(slide, 4, "병원체별 최적 조합", 17)
    add_figure(slide, "cross_pathogen_top5", Inches(0.600), Inches(1.100),
               Inches(8.000), Inches(5.200))
    # Right big number
    right = add_shape(slide, Inches(8.900), Inches(1.100), Inches(3.830), Inches(2.500),
                      fill=hex_to_rgb("FFF5E6"), border_color=hex_to_rgb("E74C3C"), border_width=3)
    tf = right.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.200)
    p = tf.paragraphs[0]
    p.text = "9/9"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb("E74C3C")
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "unique best\ncombinations"
    p2.font.size = Pt(16)
    p2.font.color.rgb = DARK_GRAY
    p2.font.name = "맑은 고딕"
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)

    # Interpretation
    interp = add_shape(slide, Inches(8.900), Inches(3.900), Inches(3.830), Inches(2.400),
                       fill=LIGHT_GRAY_BG, border_color=hex_to_rgb("8B4513"), border_width=1)
    tf2 = interp.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Inches(0.200)
    tf2.margin_top = Inches(0.150)
    lines = [
        ("의미", 16, True, hex_to_rgb("8B4513")),
        ("", 6, False, DARK_GRAY),
        ("9개 병원체 모두에서 서로", 13, False, DARK_GRAY),
        ("다른 scoring-detection 조합이", 13, False, DARK_GRAY),
        ("최고 성능을 달성", 13, False, DARK_GRAY),
        ("", 6, False, DARK_GRAY),
        ("→ 범용 최적 방법은 존재하지 않음", 13, True, hex_to_rgb("E74C3C")),
    ]
    for i, (text, sz, bld, col) in enumerate(lines):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bld
        p.font.color.rgb = col
        p.font.name = "맑은 고딕"


def slide_18_friedman(prs):
    slide = new_slide(prs)
    add_common(slide, 4, "Friedman 검정 및 LOPO", 18)
    # Left: Friedman
    left = add_shape(slide, Inches(0.600), Inches(1.200), Inches(5.850), Inches(5.200),
                     fill=LIGHT_GRAY_BG, border_color=hex_to_rgb("8B4513"), border_width=2)
    tf = left.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.300)
    tf.margin_top = Inches(0.300)
    lines = [
        ("Friedman 검정", 20, True, hex_to_rgb("8B4513")),
        ("", 8, False, DARK_GRAY),
        ("Top-20 조합 간 순위 비교", 14, False, DARK_GRAY),
        ("", 8, False, DARK_GRAY),
        ("χ² = 42.44", 28, True, hex_to_rgb("E74C3C")),
        ("p = 0.0015", 28, True, hex_to_rgb("E74C3C")),
        ("", 8, False, DARK_GRAY),
        ("→ 통계적으로 유의한 순위 차이 존재", 14, True, DARK_GRAY),
        ("→ Top-20 내에서도 성능 차이가 유의", 14, False, DARK_GRAY),
    ]
    for i, (text, sz, bld, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bld
        p.font.color.rgb = col
        p.font.name = "맑은 고딕"
        p.alignment = PP_ALIGN.CENTER if sz >= 28 else PP_ALIGN.LEFT

    # Right: LOPO
    right = add_shape(slide, Inches(6.880), Inches(1.200), Inches(5.850), Inches(5.200),
                      fill=LIGHT_GRAY_BG, border_color=hex_to_rgb("8B4513"), border_width=2)
    tf2 = right.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Inches(0.300)
    tf2.margin_top = Inches(0.300)
    lines2 = [
        ("LOPO (Leave-One-Pathogen-Out)", 20, True, hex_to_rgb("8B4513")),
        ("", 8, False, DARK_GRAY),
        ("한 병원체를 제외하고 훈련 →", 14, False, DARK_GRAY),
        ("제외된 병원체에서 최적 예측", 14, False, DARK_GRAY),
        ("", 8, False, DARK_GRAY),
    ]
    for i, (text, sz, bld, col) in enumerate(lines2):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bld
        p.font.color.rgb = col
        p.font.name = "맑은 고딕"
    # Big 0/9
    p_big = tf2.add_paragraph()
    p_big.text = "0/9"
    p_big.font.size = Pt(48)
    p_big.font.bold = True
    p_big.font.color.rgb = hex_to_rgb("E74C3C")
    p_big.font.name = "맑은 고딕"
    p_big.alignment = PP_ALIGN.CENTER
    p_match = tf2.add_paragraph()
    p_match.text = "일치율 (match)"
    p_match.font.size = Pt(16)
    p_match.font.color.rgb = DARK_GRAY
    p_match.font.name = "맑은 고딕"
    p_match.alignment = PP_ALIGN.CENTER
    p_match.space_before = Pt(8)
    p_imp = tf2.add_paragraph()
    p_imp.text = "→ 다른 병원체의 최적 조합으로\n새 병원체 예측 불가"
    p_imp.font.size = Pt(13)
    p_imp.font.color.rgb = hex_to_rgb("E74C3C")
    p_imp.font.name = "맑은 고딕"
    p_imp.space_before = Pt(8)


def slide_19_esm2(prs):
    slide = new_slide(prs)
    add_common(slide, 4, "ESM-2 교차 패러다임 검증", 19)
    # Table
    table_data = [
        ["병원체", "ESM-2 Best Rank", "기존 Best Method", "기존 Rank"],
        ["SARS-CoV-2", "Top-5", "Wavelet+Percentile", "1"],
        ["HIV", "Top-10", "Entropy+MAD", "1"],
        ["Influenza", "Top-3", "Conservation+IQR", "1"],
        ["Dengue", "Top-15", "KL-div+HDBSCAN", "1"],
        ["HCV", "Top-8", "H-score+Z-score", "1"],
    ]
    rows, cols = len(table_data), len(table_data[0])
    tbl = slide.shapes.add_table(rows, cols, Inches(0.600), Inches(1.300),
                                 Inches(12.130), Inches(3.500))
    table = tbl.table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = table_data[r][c]
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.name = "맑은 고딕"
            p.alignment = PP_ALIGN.CENTER
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = hex_to_rgb("8B4513")
            else:
                p.font.color.rgb = DARK_GRAY
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else LIGHT_GRAY_BG

    # Message box
    msg = add_shape(slide, Inches(0.600), Inches(5.200), Inches(12.130), Inches(1.200),
                    fill=hex_to_rgb("E8F4FD"), border_color=hex_to_rgb("2E5090"), border_width=2)
    tf = msg.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.300)
    tf.margin_top = Inches(0.150)
    p = tf.paragraphs[0]
    p.text = "ESM-2 protein language model은 새로운 scoring 패러다임으로서 경쟁력 확인"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb("2E5090")
    p.font.name = "맑은 고딕"
    p2 = tf.add_paragraph()
    p2.text = "그러나 병원체에 따라 순위 변동 → 병원체 특이성 재확인"
    p2.font.size = Pt(14)
    p2.font.color.rgb = DARK_GRAY
    p2.font.name = "맑은 고딕"


def slide_20_keyfinding(prs):
    slide = new_slide(prs)
    add_common(slide, 4, "핵심 발견: 4가지 증거", 20)
    # 4 cards 2x2
    evidence = [
        ("ω²=28.5%", "상호작용 효과", "병원체×방법 상호작용이\n전체 분산의 28.5% 설명",
         hex_to_rgb("2E5090"), hex_to_rgb("E8F4FD")),
        ("9/9", "고유 최적 조합", "모든 병원체에서\n서로 다른 최적 방법",
         hex_to_rgb("E74C3C"), hex_to_rgb("FDEDEC")),
        ("p=0.0015", "Friedman 유의", "Top-20 조합 간\n유의한 순위 차이",
         hex_to_rgb("1A7A4C"), hex_to_rgb("E8F9EF")),
        ("0/9", "LOPO 실패", "Leave-One-Pathogen-Out\n교차 예측 불가",
         hex_to_rgb("8B4513"), hex_to_rgb("FFF5E6")),
    ]
    positions = [(0.600, 1.200), (6.880, 1.200), (0.600, 4.000), (6.880, 4.000)]
    for (num, label, desc, col, bg), (lx, ty) in zip(evidence, positions):
        card = add_shape(slide, Inches(lx), Inches(ty), Inches(5.850), Inches(2.500),
                         fill=bg, border_color=col, border_width=3)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.300)
        tf.margin_top = Inches(0.200)
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = col
        p.font.name = "맑은 고딕"
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = DARK_GRAY
        p2.font.name = "맑은 고딕"
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(6)
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(13)
        p3.font.color.rgb = MED_GRAY
        p3.font.name = "맑은 고딕"
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(4)


def slide_21_pahd(prs):
    slide = new_slide(prs)
    add_common(slide, 5, "PAHD: 적응형 핫스팟 탐지", 21)
    # 3-step flow
    steps = [
        ("Step 1\n프로파일 추출", "병원체 MSA에서\nentropy, length, diversity\n프로파일 추출",
         hex_to_rgb("6B3FA0")),
        ("Step 2\n유사도 매칭", "MutBench DB에서\n가장 유사한 병원체\n프로파일 검색",
         hex_to_rgb("8B4513")),
        ("Step 3\nAdaptive 선택", "매칭된 병원체의\n최적 scoring-detection\n조합 적용",
         hex_to_rgb("1A7A4C")),
    ]
    col_lefts = [0.600, 4.740, 8.880]
    for (title, desc, col), lx in zip(steps, col_lefts):
        box = add_shape(slide, Inches(lx), Inches(1.200), Inches(3.850), Inches(2.200),
                        fill=col, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.200)
        tf.margin_top = Inches(0.150)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "맑은 고딕"
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        p2.font.name = "맑은 고딕"
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(8)

    # Arrows
    for i in range(2):
        ax = Inches(col_lefts[i]) + Inches(3.850)
        add_text_box(slide, ax, Inches(1.900), Inches(0.290), Inches(0.400),
                     "→", size=24, bold=True, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

    # Results table
    results_data = [
        ["방법", "Avg MCC", "vs FreqThresh", "vs SWAN"],
        ["PAHD Top-1", "0.352", "+42%", "+82%"],
        ["PAHD Top-3 avg", "0.390", "+57%", "+102%"],
        ["FreqThresh", "0.248", "baseline", "+29%"],
        ["SWAN", "0.193", "-22%", "baseline"],
        ["MutClust-Orig", "0.138", "-44%", "-28%"],
    ]
    rows, cols = len(results_data), len(results_data[0])
    tbl = slide.shapes.add_table(rows, cols, Inches(0.600), Inches(3.800),
                                 Inches(12.130), Inches(2.800))
    table = tbl.table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = results_data[r][c]
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.name = "맑은 고딕"
            p.alignment = PP_ALIGN.CENTER
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = hex_to_rgb("6B3FA0")
            elif r <= 2:
                p.font.color.rgb = DARK_GRAY
                cell.fill.solid()
                cell.fill.fore_color.rgb = hex_to_rgb("F5E6FF")
            else:
                p.font.color.rgb = DARK_GRAY
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else LIGHT_GRAY_BG


def slide_22_phylo(prs):
    slide = new_slide(prs)
    add_common(slide, 5, "계통 보정 및 영역 겹침 MCC", 22)
    # Left: TreeTime
    left = add_shape(slide, Inches(0.600), Inches(1.200), Inches(5.850), Inches(5.200),
                     fill=LIGHT_GRAY_BG, border_color=hex_to_rgb("6B3FA0"), border_width=2)
    tf = left.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.300)
    tf.margin_top = Inches(0.300)
    lines = [
        ("계통발생적 보정 (Phylogenetic Correction)", 16, True, hex_to_rgb("6B3FA0")),
        ("", 8, False, DARK_GRAY),
        ("문제: 변이 빈도 ≠ 독립 사건", 14, True, DARK_GRAY),
        ("(계통 구조로 인한 가짜 핫스팟)", 13, False, MED_GRAY),
        ("", 6, False, DARK_GRAY),
        ("해결: TreeTime 기반 ancestral reconstruction", 14, False, DARK_GRAY),
        ("→ 독립 변이 이벤트만 카운트", 14, False, DARK_GRAY),
        ("", 6, False, DARK_GRAY),
        ("결과:", 14, True, hex_to_rgb("6B3FA0")),
        ("• D614G: founder effect 확인", 13, False, DARK_GRAY),
        ("• 보정 후 rho = -0.876 → -0.412", 13, False, DARK_GRAY),
        ("• 진정한 적응 핫스팟 분리 가능", 13, False, DARK_GRAY),
    ]
    for i, (text, sz, bld, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bld
        p.font.color.rgb = col
        p.font.name = "맑은 고딕"

    # Right: Region overlap improvement
    right = add_shape(slide, Inches(6.880), Inches(1.200), Inches(5.850), Inches(5.200),
                      fill=LIGHT_GRAY_BG, border_color=hex_to_rgb("6B3FA0"), border_width=2)
    tf2 = right.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Inches(0.300)
    tf2.margin_top = Inches(0.300)
    lines2 = [
        ("영역 겹침 MCC 개선", 16, True, hex_to_rgb("6B3FA0")),
        ("", 8, False, DARK_GRAY),
        ("Position-level MCC (기존)", 14, False, DARK_GRAY),
    ]
    for i, (text, sz, bld, col) in enumerate(lines2):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bld
        p.font.color.rgb = col
        p.font.name = "맑은 고딕"

    # Big numbers
    p_num = tf2.add_paragraph()
    p_num.text = "0.289"
    p_num.font.size = Pt(32)
    p_num.font.bold = True
    p_num.font.color.rgb = MED_GRAY
    p_num.font.name = "맑은 고딕"
    p_num.alignment = PP_ALIGN.CENTER
    p_num.space_before = Pt(8)

    p_arrow = tf2.add_paragraph()
    p_arrow.text = "↓ Region-overlap MCC"
    p_arrow.font.size = Pt(14)
    p_arrow.font.color.rgb = DARK_GRAY
    p_arrow.font.name = "맑은 고딕"
    p_arrow.alignment = PP_ALIGN.CENTER
    p_arrow.space_before = Pt(6)

    p_num2 = tf2.add_paragraph()
    p_num2.text = "0.638"
    p_num2.font.size = Pt(32)
    p_num2.font.bold = True
    p_num2.font.color.rgb = hex_to_rgb("1A7A4C")
    p_num2.font.name = "맑은 고딕"
    p_num2.alignment = PP_ALIGN.CENTER
    p_num2.space_before = Pt(6)

    p_arrow2 = tf2.add_paragraph()
    p_arrow2.text = "↓ + Phylo correction"
    p_arrow2.font.size = Pt(14)
    p_arrow2.font.color.rgb = DARK_GRAY
    p_arrow2.font.name = "맑은 고딕"
    p_arrow2.alignment = PP_ALIGN.CENTER
    p_arrow2.space_before = Pt(6)

    p_num3 = tf2.add_paragraph()
    p_num3.text = "0.712"
    p_num3.font.size = Pt(36)
    p_num3.font.bold = True
    p_num3.font.color.rgb = hex_to_rgb("E74C3C")
    p_num3.font.name = "맑은 고딕"
    p_num3.alignment = PP_ALIGN.CENTER
    p_num3.space_before = Pt(6)


def slide_23_baseline(prs):
    slide = new_slide(prs)
    add_common(slide, 5, "MutBench vs 기존 방법", 23)
    # Figure
    add_figure(slide, "scoring_detection_heatmap", Inches(0.600), Inches(1.100),
               Inches(7.500), Inches(5.200))
    # Right: comparison bars (text-based)
    right = add_shape(slide, Inches(8.400), Inches(1.100), Inches(4.330), Inches(5.200),
                      fill=LIGHT_GRAY_BG, border_color=hex_to_rgb("6B3FA0"), border_width=2)
    tf = right.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.250)
    tf.margin_top = Inches(0.200)
    tf.margin_right = Inches(0.250)
    comparisons = [
        ("MCC 비교 (평균)", 18, True, hex_to_rgb("6B3FA0")),
        ("", 6, False, DARK_GRAY),
        ("MutBench Top-3", 14, True, DARK_GRAY),
        ("MCC = 0.390", 24, True, hex_to_rgb("1A7A4C")),
        ("", 6, False, DARK_GRAY),
        ("FreqThresh", 14, True, DARK_GRAY),
        ("MCC = 0.248", 20, True, hex_to_rgb("2E5090")),
        ("", 6, False, DARK_GRAY),
        ("SWAN", 14, True, DARK_GRAY),
        ("MCC = 0.193", 20, True, hex_to_rgb("F39C12")),
        ("", 6, False, DARK_GRAY),
        ("MutClust-Orig", 14, True, DARK_GRAY),
        ("MCC = 0.138", 20, True, hex_to_rgb("E74C3C")),
        ("", 6, False, DARK_GRAY),
        ("Random baseline", 14, True, DARK_GRAY),
        ("MCC = 0.001", 20, True, MED_GRAY),
    ]
    for i, (text, sz, bld, col) in enumerate(comparisons):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bld
        p.font.color.rgb = col
        p.font.name = "맑은 고딕"


def slide_24_discussion(prs):
    slide = new_slide(prs)
    add_common(slide, 6, "논의", 24)
    discussions = [
        ("병원체 특이성의 원인", hex_to_rgb("2E5090"),
         [
             "Entropy-function 관계 병원체별 상이",
             "예: Norovirus r=0.39 vs SARS-CoV-2 r=0.03",
             "진화율, 게놈 크기, 재조합 빈도의 영향",
             "기능적 제약의 병원체별 분포 차이",
         ]),
        ("벤치마크의 한계", hex_to_rgb("E67E22"),
         [
             "DMS 데이터: SARS-CoV-2 Spike에만 존재",
             "Phylogenetic non-independence 완전 해결 미달",
             "RNA 바이러스에 한정된 평가",
             "시간에 따른 핫스팟 동적 변화 미반영",
         ]),
        ("실용적 함의", hex_to_rgb("1A7A4C"),
         [
             "새 병원체 → MutBench DB 프로파일 매칭 → 최적 방법 추천",
             "백신 설계 시 탐지 방법 선택의 중요성 인식",
             "다중 GT 사용으로 결과 신뢰도 향상",
             "PAHD를 통한 자동화된 방법 선택 가능",
         ]),
    ]
    col_lefts = [0.600, 4.740, 8.880]
    for (title, col, bullets), lx in zip(discussions, col_lefts):
        add_card(slide, Inches(lx), Inches(1.200), Inches(3.850), Inches(5.300),
                 title, bullets, col)


def slide_25_contrib_limit(prs):
    slide = new_slide(prs)
    add_common(slide, 6, "기여, 한계, 향후 과제", 25)
    cols_data = [
        ("기여 (Contributions)", hex_to_rgb("1A7A4C"),
         [
             "MutBench: 최초의 체계적 벤치마크",
             "3-Layer GT로 편향 최소화",
             "병원체 특이성 정량적 실증",
             "PAHD 적응형 탐지 제안",
             "9-pathogen 대규모 평가",
         ]),
        ("한계 (Limitations)", hex_to_rgb("E74C3C"),
         [
             "RNA 바이러스에 한정",
             "DMS = SARS-CoV-2 Spike만",
             "계통 보정 완전하지 않음",
             "시간적 동적 변화 미반영",
             "DNA 바이러스 미포함",
         ]),
        ("향후 과제 (Future Work)", hex_to_rgb("2E5090"),
         [
             "DNA 바이러스 확장",
             "실시간 핫스팟 모니터링",
             "PLM 기반 scoring 강화",
             "PAHD 자동화 시스템 개발",
             "웹 기반 MutBench 도구 공개",
         ]),
    ]
    col_lefts = [0.600, 4.740, 8.880]
    for (title, col, bullets), lx in zip(cols_data, col_lefts):
        add_card(slide, Inches(lx), Inches(1.200), Inches(3.850), Inches(5.300),
                 title, bullets, col)


def slide_26_takeaway(prs):
    slide = new_slide(prs)
    add_common(slide, 6, "핵심 결론", 26)
    # Central big message
    central = add_shape(slide, Inches(1.500), Inches(1.300), Inches(10.330), Inches(2.000),
                        fill=hex_to_rgb("1B2A4A"))
    tf = central.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.400)
    tf.margin_top = Inches(0.200)
    p = tf.paragraphs[0]
    p.text = "바이러스 변이 핫스팟 탐지에 범용 최적 방법은 존재하지 않으며,"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "맑은 고딕"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "병원체 특성에 따른 적응적 방법 선택(PAHD)이 필수적이다."
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = hex_to_rgb("FFD700")
    p2.font.name = "맑은 고딕"
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)

    # 4 evidence boxes below
    evidence = [
        ("ω²=28.5%", "상호작용", hex_to_rgb("2E5090")),
        ("9/9 unique", "최적 조합", hex_to_rgb("E74C3C")),
        ("p=0.0015", "Friedman", hex_to_rgb("1A7A4C")),
        ("0/9 LOPO", "교차 예측 실패", hex_to_rgb("8B4513")),
    ]
    for i, (num, label, col) in enumerate(evidence):
        lx = 0.600 + i * 3.200
        add_big_number(slide, Inches(lx), Inches(3.800), Inches(2.900), Inches(2.200),
                       num, label, col, bg_fill=LIGHT_GRAY_BG)


def slide_27_thankyou(prs):
    slide = new_slide(prs)
    # Full dark top
    add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(3.500),
              fill=SECTION_COLORS[1])
    # KNU logo
    circle = add_shape(slide, Inches(12.230), Inches(0.150), Inches(0.550), Inches(0.550),
                       fill=WHITE, shape_type=MSO_SHAPE.OVAL)
    set_text(circle, "KNU", size=9, bold=True, color=SECTION_COLORS[1],
             alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Thank you
    add_text_box(slide, Inches(0.600), Inches(1.200), Inches(12.130), Inches(1.200),
                 "감사합니다", size=44, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.600), Inches(2.400), Inches(12.130), Inches(0.500),
                 "Thank you for your attention", size=20, color=hex_to_rgb("B0C4DE"),
                 alignment=PP_ALIGN.CENTER)
    # Contact card
    contact = add_shape(slide, Inches(3.500), Inches(4.200), Inches(6.330), Inches(2.200),
                        fill=LIGHT_GRAY_BG, border_color=hex_to_rgb("2E5090"), border_width=2)
    tf = contact.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.400)
    tf.margin_top = Inches(0.300)
    lines = [
        ("권휘준 (Hwijun Kwon)", 18, True, DARK_GRAY),
        ("경북대학교 대학원 컴퓨터학부", 14, False, MED_GRAY),
        ("지도교수: 정인욱", 14, False, MED_GRAY),
        ("2026년 3월", 14, False, MED_GRAY),
    ]
    for i, (text, sz, bld, col) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bld
        p.font.color.rgb = col
        p.font.name = "맑은 고딕"
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(4)

    # Bottom bar
    add_shape(slide, Inches(0), Inches(7.120), SLIDE_W, Inches(0.380), fill=BOTTOM_BAR_COLOR)


# ── Main ───────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_toc(prs)
    slide_03_prior(prs)
    slide_04_hotspots(prs)
    slide_05_problem(prs)
    slide_06_contributions(prs)
    slide_07_pipeline(prs)
    slide_08_gt(prs)
    slide_09_mcc(prs)
    slide_10_twostage(prs)
    slide_11_stage1(prs)
    slide_12_sensitivity(prs)
    slide_13_gap(prs)
    slide_14_ranking(prs)
    slide_15_9pathogen(prs)
    slide_16_anova(prs)
    slide_17_perpathogen(prs)
    slide_18_friedman(prs)
    slide_19_esm2(prs)
    slide_20_keyfinding(prs)
    slide_21_pahd(prs)
    slide_22_phylo(prs)
    slide_23_baseline(prs)
    slide_24_discussion(prs)
    slide_25_contrib_limit(prs)
    slide_26_takeaway(prs)
    slide_27_thankyou(prs)

    out_path = "/proj/paper/paper/ppt/MutBench_Defense_v2.pptx"
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
