#!/usr/bin/env python3
"""
MutBench PhD Defense Presentation v3 — ALL ENGLISH
Output: /proj/paper/paper/ppt/MutBench_Defense_v3.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── paths ──────────────────────────────────────────────────────
BASE = "/proj/paper"
OUT = os.path.join(BASE, "paper/ppt/MutBench_Defense_v3.pptx")
LOGO = os.path.join(BASE, "docs/presentation/knu_symbol_white.png")
LOGO_EMBLEM = os.path.join(BASE, "docs/presentation/knu_emblem_official.jpg")
FIG_DISS = os.path.join(BASE, "paper/dissertation/figures")
FIG_MAIN = os.path.join(BASE, "paper/figure")

# ── colors ─────────────────────────────────────────────────────
DARK_NAVY  = RGBColor(0x1B, 0x2A, 0x4A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
GRAY       = RGBColor(0x88, 0x88, 0x88)
DARK_GRAY  = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY   = RGBColor(0xE0, 0xE0, 0xE0)

# section colors
SEC_INTRO   = RGBColor(0x1B, 0x2A, 0x4A)  # navy
SEC_FRAME   = RGBColor(0x0D, 0x73, 0x77)  # teal
SEC_STAGE1  = RGBColor(0x2E, 0x7D, 0x32)  # green
SEC_STAGE2  = RGBColor(0xE6, 0x51, 0x00)  # orange
SEC_EXT     = RGBColor(0x6A, 0x1B, 0x9A)  # purple
SEC_CONCL   = RGBColor(0xC6, 0x28, 0x28)  # red

TEAL_ACCENT = RGBColor(0x0D, 0x73, 0x77)
RED_ACCENT  = RGBColor(0xC6, 0x28, 0x28)
GREEN       = RGBColor(0x2E, 0x7D, 0x32)
ORANGE      = RGBColor(0xE6, 0x51, 0x00)
PURPLE      = RGBColor(0x6A, 0x1B, 0x9A)
BLUE        = RGBColor(0x15, 0x65, 0xC0)

SLIDE_W = Inches(13.330)
SLIDE_H = Inches(7.500)

# ── helpers ────────────────────────────────────────────────────

def make_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def add_blank(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_rect(slide, l, t, w, h, fill=None, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.fill.solid()
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    return shape

def add_rounded_rect(slide, l, t, w, h, fill=None, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.fill.solid()
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    return shape

def set_text(shape, text, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
             font_name="Calibri", anchor=MSO_ANCHOR.TOP):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except:
        pass
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = align

def add_textbox(slide, l, t, w, h, text, size=14, bold=False, color=BLACK,
                align=PP_ALIGN.LEFT, font_name="Calibri", anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    set_text(txBox, text, size, bold, color, align, font_name, anchor)
    return txBox

def add_multiformat_textbox(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT):
    """runs: list of (text, size, bold, color) tuples"""
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    for text, size, bold, color in runs:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox

def content_slide(prs, title, section_color):
    """Create standard content slide with top bar, title, logo, bottom bar, footer."""
    slide = add_blank(prs)
    # top bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.900), fill=section_color)
    # title
    add_textbox(slide, Inches(0.600), Inches(0.180), Inches(10.500), Inches(0.550),
                title, size=24, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # KNU logo top-right
    slide.shapes.add_picture(LOGO, Inches(12.230), Inches(0.150),
                             Inches(0.550), Inches(0.550))
    # bottom bar
    add_rect(slide, Inches(0), Inches(7.120), SLIDE_W, Inches(0.380), fill=DARK_NAVY)
    # footer text
    add_textbox(slide, Inches(0.400), Inches(7.140), Inches(10.0), Inches(0.340),
                "Kyungpook National University  |  Dept. of Computer Science  |  Hwijun Kwon",
                size=10, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    return slide

def add_page_number(slide, num):
    add_textbox(slide, Inches(12.500), Inches(7.140), Inches(0.700), Inches(0.340),
                str(num), size=10, color=WHITE, align=PP_ALIGN.RIGHT,
                anchor=MSO_ANCHOR.MIDDLE)

def add_card(slide, l, t, w, h, fill=WHITE, border=MID_GRAY):
    return add_rounded_rect(slide, l, t, w, h, fill=fill, border=border)

def add_accent_bar(slide, l, t, w, h, color):
    return add_rect(slide, l, t, w, h, fill=color)

def add_figure(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, l, t, w, h)
    else:
        box = add_card(slide, l, t, w, h, fill=LIGHT_GRAY, border=MID_GRAY)
        set_text(box, f"[Missing: {os.path.basename(path)}]", size=10, color=GRAY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ══════════════════════════════════════════════════════════════
# BUILD PRESENTATION
# ══════════════════════════════════════════════════════════════

prs = make_prs()

# ── SLIDE 1: Title ─────────────────────────────────────────────
s = add_blank(prs)
# dark top area
add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(3.000), fill=DARK_NAVY)
# KNU logo
s.shapes.add_picture(LOGO, Inches(0.600), Inches(0.250), Inches(0.650), Inches(0.650))
# university + defense
add_textbox(s, Inches(1.400), Inches(0.300), Inches(8.0), Inches(0.300),
            "Graduate School, Kyungpook National University", size=14, color=WHITE)
add_textbox(s, Inches(1.400), Inches(0.600), Inches(5.0), Inches(0.300),
            "PhD Dissertation Defense", size=14, color=WHITE)
# title
add_textbox(s, Inches(0.600), Inches(1.250), Inches(12.0), Inches(1.500),
            "MutBench: Systematic Benchmarking Framework\nfor Viral Mutation Hotspot Detection",
            size=32, bold=True, color=WHITE)
# KNU emblem (below dark area)
if os.path.exists(LOGO_EMBLEM):
    s.shapes.add_picture(LOGO_EMBLEM, Inches(0.600), Inches(3.200), Inches(0.650), Inches(0.650))
# author info (below dark area)
add_textbox(s, Inches(1.400), Inches(3.300), Inches(11.0), Inches(0.400),
            "Hwijun Kwon  |  Advisor: Prof. Inuk Jung", size=16, color=DARK_GRAY)
add_textbox(s, Inches(1.400), Inches(3.700), Inches(11.0), Inches(0.350),
            "Dept. of Computer Science, KNU  |  March 2026", size=14, color=GRAY)

# paper cards
card_w = Inches(5.800)
card_h = Inches(2.200)
card_y = Inches(4.600)
gap = Inches(0.400)
card1_x = Inches(0.600)
card2_x = Inches(6.900)

# Card 1: BMC Genomics
c1 = add_card(s, card1_x, card_y, card_w, card_h, fill=WHITE, border=MID_GRAY)
add_accent_bar(s, card1_x, card_y, Inches(0.080), card_h, TEAL_ACCENT)
add_textbox(s, Inches(0.900), Inches(4.700), Inches(5.300), Inches(0.350),
            "Paper 1  |  BMC Genomics 2025 (SCIE)", size=13, bold=True, color=TEAL_ACCENT)
add_textbox(s, Inches(0.900), Inches(5.100), Inches(5.300), Inches(1.400),
            "MOSD: Multi-Omics Subtyping with Deep learning\n\n"
            "Multi-omics integration benchmark for cancer\n"
            "subtype discovery (11 methods, 6 metrics)",
            size=11, color=DARK_GRAY)

# Card 2: BioData Mining
c2 = add_card(s, card2_x, card_y, card_w, card_h, fill=WHITE, border=MID_GRAY)
add_accent_bar(s, card2_x, card_y, Inches(0.080), card_h, RED_ACCENT)
add_textbox(s, Inches(7.200), Inches(4.700), Inches(5.300), Inches(0.350),
            "Paper 2  |  BioData Mining 2025 (SCIE)", size=13, bold=True, color=RED_ACCENT)
add_textbox(s, Inches(7.200), Inches(5.100), Inches(5.300), Inches(1.400),
            "MutClust: Mutation Clustering for viral evolution\n\n"
            "H-score based mutation hotspot detection\n"
            "with network propagation & bootstrap",
            size=11, color=DARK_GRAY)


# ── SLIDE 2: Table of Contents ─────────────────────────────────
s = content_slide(prs, "Table of Contents", SEC_INTRO)
add_page_number(s, 2)
cards_data = [
    ("1. Introduction", "Mutation hotspots, problem definition,\nresearch contributions", SEC_INTRO),
    ("2. MutBench Framework", "Pipeline, ground truth design,\nevaluation metrics", SEC_FRAME),
    ("3. Stage 1: 4-Pathogen", "Method comparison, parameter\nsensitivity, synthetic vs real", SEC_STAGE1),
    ("4. Stage 2: 9-Pathogen", "ANOVA, per-pathogen best,\nFriedman test, LOPO CV", SEC_STAGE2),
    ("5. Extensions", "PAHD concept, phylogenetic\ncorrection, baseline comparison", SEC_EXT),
    ("6. Conclusion", "Key takeaways, contributions,\nlimitations, future work", SEC_CONCL),
]
cw, ch = Inches(5.800), Inches(1.600)
for i, (title, desc, color) in enumerate(cards_data):
    col = i % 2
    row = i // 2
    cx = Inches(0.600) + col * Inches(6.300)
    cy = Inches(1.200) + row * Inches(1.850)
    card = add_card(s, cx, cy, cw, ch, fill=WHITE, border=MID_GRAY)
    add_accent_bar(s, cx, cy, Inches(0.080), ch, color)
    add_textbox(s, cx + Inches(0.250), cy + Inches(0.150), Inches(5.300), Inches(0.400),
                title, size=16, bold=True, color=color)
    add_textbox(s, cx + Inches(0.250), cy + Inches(0.600), Inches(5.300), Inches(0.900),
                desc, size=12, color=DARK_GRAY)


# ── SLIDE 3: Prior Work ───────────────────────────────────────
s = content_slide(prs, "Prior Work: MOSD & MutClust", SEC_INTRO)
add_page_number(s, 3)
# MOSD card
c = add_card(s, Inches(0.600), Inches(1.200), Inches(5.800), Inches(3.800), border=MID_GRAY)
add_rect(s, Inches(0.600), Inches(1.200), Inches(5.800), Inches(0.550), fill=TEAL_ACCENT)
add_textbox(s, Inches(0.800), Inches(1.220), Inches(5.400), Inches(0.500),
            "MOSD (BMC Genomics 2025)", size=16, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
add_textbox(s, Inches(0.800), Inches(1.900), Inches(5.400), Inches(2.900),
            "Multi-Omics Subtyping with Deep learning\n\n"
            "  Benchmarked 11 multi-omics integration methods\n"
            "  across 6 evaluation metrics\n\n"
            "  Key finding: No single method dominates\n"
            "  across all cancer datasets\n\n"
            "  Lesson: Method selection depends on\n"
            "  data characteristics",
            size=12, color=DARK_GRAY)

# MutClust card
c = add_card(s, Inches(6.900), Inches(1.200), Inches(5.800), Inches(3.800), border=MID_GRAY)
add_rect(s, Inches(6.900), Inches(1.200), Inches(5.800), Inches(0.550), fill=RED_ACCENT)
add_textbox(s, Inches(7.100), Inches(1.220), Inches(5.400), Inches(0.500),
            "MutClust (BioData Mining 2025)", size=16, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
add_textbox(s, Inches(7.100), Inches(1.900), Inches(5.400), Inches(2.900),
            "Mutation Clustering for viral evolution\n\n"
            "  H-score based mutation hotspot detection\n"
            "  with network propagation & bootstrap\n\n"
            "  Limitation: H-score shows founder bias\n"
            "  (rho = -0.876 with sequence count)\n\n"
            "  Lesson: Single-pathogen evaluation\n"
            "  is insufficient",
            size=12, color=DARK_GRAY)

# bridge
bridge = add_card(s, Inches(2.500), Inches(5.300), Inches(8.300), Inches(0.700),
                  fill=DARK_NAVY, border=None)
set_text(bridge, "Common lesson: systematic benchmarking is essential",
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ── SLIDE 4: What are Mutation Hotspots? ──────────────────────
s = content_slide(prs, "What are Mutation Hotspots?", SEC_INTRO)
add_page_number(s, 4)

# definition box
defbox = add_card(s, Inches(0.600), Inches(1.200), Inches(12.100), Inches(1.200),
                  fill=RGBColor(0xE8, 0xF5, 0xE9), border=GREEN)
set_text(defbox, "Definition:  Genomic positions where mutations accumulate at significantly\n"
         "higher frequencies than expected by chance, indicating selective pressure.",
         size=14, color=DARK_GRAY, anchor=MSO_ANCHOR.MIDDLE)

# 3 cards
topics = [
    ("Why do they occur?", "Positive selection drives beneficial\nmutations (e.g., immune escape).\n"
     "Some regions tolerate mutations\nbetter than others."),
    ("Why are they important?", "Track viral evolution in real-time.\n"
     "Guide vaccine & drug target design.\n"
     "Predict emerging variants of concern."),
    ("Current problems", "No standard benchmark exists.\n"
     "Methods tested on 1-2 viruses only.\n"
     "No ground truth consensus."),
]
for i, (ttl, desc) in enumerate(topics):
    cx = Inches(0.600) + i * Inches(4.100)
    cy = Inches(2.700)
    card = add_card(s, cx, cy, Inches(3.800), Inches(3.800), border=MID_GRAY)
    add_textbox(s, cx + Inches(0.200), cy + Inches(0.200), Inches(3.400), Inches(0.400),
                ttl, size=14, bold=True, color=DARK_NAVY)
    add_textbox(s, cx + Inches(0.200), cy + Inches(0.700), Inches(3.400), Inches(2.900),
                desc, size=12, color=DARK_GRAY)


# ── SLIDE 4B: Key Terminology for CS Audience ────────────────
s = content_slide(prs, "Key Terminology", SEC_INTRO)
add_page_number(s, 5)

# Left column header
add_rect(s, Inches(0.600), Inches(1.200), Inches(5.800), Inches(0.500), fill=SEC_FRAME)
add_textbox(s, Inches(0.600), Inches(1.210), Inches(5.800), Inches(0.480),
            "Genomics Basics", size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)

left_terms = [
    ("Genome", "Complete genetic blueprint of an organism (DNA/RNA sequence)"),
    ("Mutation", "Change in the nucleotide sequence (A, T, G, C/U)"),
    ("Protein", "Functional molecule encoded by genes; amino acid sequence determines 3D structure"),
    ("MSA", "Multiple Sequence Alignment: aligning multiple sequences to identify conserved/variable regions"),
]
for i, (term, desc) in enumerate(left_terms):
    cy = Inches(1.850) + i * Inches(1.200)
    card = add_card(s, Inches(0.600), cy, Inches(5.800), Inches(1.050), border=MID_GRAY)
    add_multiformat_textbox(s, Inches(0.800), cy + Inches(0.150), Inches(5.400), Inches(0.350),
                            [(term, 13, True, DARK_NAVY)], align=PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.800), cy + Inches(0.500), Inches(5.400), Inches(0.500),
                desc, size=11, color=DARK_GRAY)

# Right column header
add_rect(s, Inches(6.900), Inches(1.200), Inches(5.800), Inches(0.500), fill=ORANGE)
add_textbox(s, Inches(6.900), Inches(1.210), Inches(5.800), Inches(0.480),
            "Viral Evolution", size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)

right_terms = [
    ("Variant", "Virus with distinct mutations (e.g., Alpha, Delta, Omicron)"),
    ("Positive selection", "Mutations that increase viral fitness (immune evasion, transmissibility)"),
    ("Purifying selection", "Mutations removed because they harm essential functions"),
    ("Convergent evolution", "Same mutation arising independently in different lineages"),
]
for i, (term, desc) in enumerate(right_terms):
    cy = Inches(1.850) + i * Inches(1.200)
    card = add_card(s, Inches(6.900), cy, Inches(5.800), Inches(1.050), border=MID_GRAY)
    add_multiformat_textbox(s, Inches(7.100), cy + Inches(0.150), Inches(5.400), Inches(0.350),
                            [(term, 13, True, DARK_NAVY)], align=PP_ALIGN.LEFT)
    add_textbox(s, Inches(7.100), cy + Inches(0.500), Inches(5.400), Inches(0.500),
                desc, size=11, color=DARK_GRAY)


# ── SLIDE 4C: Why Hotspot Detection Matters ──────────────────
s = content_slide(prs, "Why Hotspot Detection Matters", SEC_INTRO)
add_page_number(s, 6)

# 3 cards
why_cards = [
    ("Vaccine Design", SEC_FRAME,
     "Hotspots = regions likely to mutate\n\n"
     "Vaccines should target conserved\nregions or predict variable ones"),
    ("Variant Surveillance", ORANGE,
     "Early detection of new hotspots\n\n"
     "Early warning system for\nemerging variants"),
    ("Drug Development", PURPLE,
     "Hotspots in drug-binding sites\n\n"
     "Predict drug resistance\nmutations before they spread"),
]
for i, (ttl, color, desc) in enumerate(why_cards):
    cx = Inches(0.600) + i * Inches(4.100)
    cy = Inches(1.400)
    card = add_card(s, cx, cy, Inches(3.800), Inches(3.800), border=MID_GRAY)
    add_rect(s, cx, cy, Inches(3.800), Inches(0.550), fill=color)
    add_textbox(s, cx, cy + Inches(0.070), Inches(3.800), Inches(0.450),
                ttl, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, cx + Inches(0.250), cy + Inches(0.800), Inches(3.300), Inches(2.800),
                desc, size=13, color=DARK_GRAY)

# bottom gap message
gap_msg = add_card(s, Inches(1.500), Inches(5.600), Inches(10.300), Inches(0.900),
                   fill=RGBColor(0xFF, 0xEB, 0xEE), border=RED_ACCENT)
set_text(gap_msg,
         "Current gap: No systematic way to compare hotspot detection methods across pathogens",
         size=14, bold=True, color=RED_ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ── SLIDE 5: Problem Definition ──────────────────────────────
s = content_slide(prs, "Problem Definition", SEC_INTRO)
add_page_number(s, 7)
problems = [
    ("P1", "No Standardized Benchmark",
     "Existing methods evaluated on different datasets with different metrics, "
     "making fair comparison impossible."),
    ("P2", "Single-Pathogen Bias",
     "Methods tested on only 1-2 pathogens cannot generalize. "
     "Pathogen-specific characteristics dramatically affect detection performance."),
    ("P3", "No Universal Best Method",
     "Our prior work (MOSD, MutClust) demonstrated that no single approach dominates "
     "across all settings, demanding systematic evaluation."),
]
for i, (pid, title, desc) in enumerate(problems):
    cy = Inches(1.200) + i * Inches(1.800)
    card = add_card(s, Inches(0.600), cy, Inches(12.100), Inches(1.550), border=MID_GRAY)
    # problem id badge
    badge = add_rounded_rect(s, Inches(0.800), cy + Inches(0.250), Inches(0.600), Inches(0.450),
                             fill=DARK_NAVY)
    set_text(badge, pid, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, Inches(1.600), cy + Inches(0.200), Inches(10.500), Inches(0.400),
                title, size=16, bold=True, color=DARK_NAVY)
    add_textbox(s, Inches(1.600), cy + Inches(0.650), Inches(10.500), Inches(0.800),
                desc, size=12, color=DARK_GRAY)


# ── SLIDE 6: Research Contributions ──────────────────────────
s = content_slide(prs, "Research Contributions", SEC_INTRO)
add_page_number(s, 8)
contribs = [
    ("C1: MutBench Framework",
     "First systematic benchmark for viral\nmutation hotspot detection with\n"
     "standardized pipeline & evaluation"),
    ("C2: 3-Layer Ground Truth",
     "Adaptive threshold + constrained +\nDMS-based ground truth design\n"
     "reduces single-reference bias"),
    ("C3: Cross-Pathogen Evidence",
     "9 pathogens, 2,544 evaluations prove\n"
     "no universal best method exists\n"
     "(interaction omega-sq = 0.285)"),
    ("C4: PAHD Proof of Concept",
     "Pathogen-Adaptive Hotspot Detection:\n"
     "profile-based method selection\n"
     "outperforms any single approach"),
]
for i, (ttl, desc) in enumerate(contribs):
    col = i % 2
    row = i // 2
    cx = Inches(0.600) + col * Inches(6.300)
    cy = Inches(1.200) + row * Inches(2.600)
    card = add_card(s, cx, cy, Inches(5.800), Inches(2.300), border=MID_GRAY)
    color = [SEC_FRAME, GREEN, SEC_STAGE2, PURPLE][i]
    add_accent_bar(s, cx, cy, Inches(0.080), Inches(2.300), color)
    add_textbox(s, cx + Inches(0.250), cy + Inches(0.200), Inches(5.300), Inches(0.400),
                ttl, size=16, bold=True, color=color)
    add_textbox(s, cx + Inches(0.250), cy + Inches(0.700), Inches(5.300), Inches(1.400),
                desc, size=12, color=DARK_GRAY)


# ── SLIDE 9: MutBench Pipeline ───────────────────────────────
s = content_slide(prs, "MutBench Pipeline Overview", SEC_FRAME)
add_page_number(s, 9)

# flow boxes
flow_items = ["Input\nSequences", "MSA\nAlignment", "Scoring\n(9 methods)", "Detection\n(39 methods)",
              "Ground\nTruth", "Evaluation\n(MCC)"]
box_w = Inches(1.700)
box_h = Inches(1.000)
start_x = Inches(0.600)
for i, item in enumerate(flow_items):
    bx = start_x + i * Inches(2.050)
    by = Inches(1.300)
    fill_c = SEC_FRAME if i < 2 else (GREEN if i < 4 else (ORANGE if i == 4 else DARK_NAVY))
    box = add_rounded_rect(s, bx, by, box_w, box_h, fill=fill_c)
    set_text(box, item, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    # arrow
    if i < len(flow_items) - 1:
        add_textbox(s, bx + box_w, by + Inches(0.200), Inches(0.350), Inches(0.500),
                    ">", size=20, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# scoring table
add_textbox(s, Inches(0.600), Inches(2.600), Inches(5.800), Inches(0.400),
            "9 Scoring Methods", size=14, bold=True, color=SEC_FRAME)
scoring_items = [
    "Shannon Entropy", "Jensen-Shannon Divergence", "Mutation Frequency",
    "Wavelet (CWT)", "Sliding Window", "Phylo-aware Score",
    "dN/dS ratio", "Kabat Variability", "Property Entropy"
]
for i, item in enumerate(scoring_items):
    col = i % 3
    row = i // 3
    add_textbox(s, Inches(0.800) + col * Inches(1.900),
                Inches(3.050) + row * Inches(0.400),
                Inches(1.800), Inches(0.350),
                f"  {item}", size=10, color=DARK_GRAY)

# detection table
add_textbox(s, Inches(6.900), Inches(2.600), Inches(5.800), Inches(0.400),
            "Detection Families (39 methods)", size=14, bold=True, color=GREEN)
detect_items = [
    "Z-score (4 variants)", "Percentile (4 variants)", "IQR outlier (2 variants)",
    "DBSCAN (6 variants)", "HDBSCAN (3 variants)", "Gaussian Mixture (4 variants)",
    "Kernel Density (4 variants)", "Isolation Forest (3 variants)", "Wavelet peaks (3 variants)",
    "LOF (2 variants)", "Threshold (2 variants)", "Bayesian CP (2 variants)",
]
for i, item in enumerate(detect_items):
    col = i % 2
    row = i // 2
    add_textbox(s, Inches(7.100) + col * Inches(2.900),
                Inches(3.050) + row * Inches(0.400),
                Inches(2.800), Inches(0.350),
                f"  {item}", size=10, color=DARK_GRAY)


# ── SLIDE 10: 3-Layer Ground Truth ────────────────────────────
s = content_slide(prs, "3-Layer Ground Truth Design", SEC_FRAME)
add_page_number(s, 10)
add_figure(s, os.path.join(FIG_DISS, "multi_ground_truth_figure.png"),
           Inches(0.400), Inches(1.100), Inches(7.500), Inches(5.600))

# 3 layer cards on right
layers = [
    ("Layer 1: Adaptive", "Entropy-based threshold auto-adjusted\nper pathogen. Captures high-variability\npositions relative to background."),
    ("Layer 2: Constrained", "Intersection of top-K positions across\nmultiple scoring methods. Reduces\nsingle-method bias."),
    ("Layer 3: DMS-based", "Deep mutational scanning data from\nexperimental fitness measurements.\nGold-standard biological validation."),
]
for i, (ttl, desc) in enumerate(layers):
    cy = Inches(1.200) + i * Inches(1.900)
    card = add_card(s, Inches(8.200), cy, Inches(4.500), Inches(1.650), border=MID_GRAY)
    add_textbox(s, Inches(8.400), cy + Inches(0.100), Inches(4.100), Inches(0.400),
                ttl, size=13, bold=True, color=SEC_FRAME)
    add_textbox(s, Inches(8.400), cy + Inches(0.550), Inches(4.100), Inches(1.000),
                desc, size=11, color=DARK_GRAY)


# ── SLIDE 11: Evaluation Metric: MCC ──────────────────────────
s = content_slide(prs, "Evaluation Metric: Matthews Correlation Coefficient", SEC_FRAME)
add_page_number(s, 11)

# MCC formula box
fbox = add_card(s, Inches(0.600), Inches(1.200), Inches(6.500), Inches(2.500),
                fill=RGBColor(0xF3, 0xE5, 0xF5), border=PURPLE)
add_textbox(s, Inches(0.800), Inches(1.300), Inches(6.100), Inches(0.400),
            "MCC Formula", size=16, bold=True, color=PURPLE)
add_textbox(s, Inches(0.800), Inches(1.750), Inches(6.100), Inches(0.500),
            "MCC = (TP*TN - FP*FN) / sqrt((TP+FP)(TP+FN)(TN+FP)(TN+FN))",
            size=13, bold=True, color=DARK_GRAY)
add_textbox(s, Inches(0.800), Inches(2.350), Inches(6.100), Inches(1.200),
            "Range: [-1, +1]   |   0 = random   |   1 = perfect\n\n"
            "Why MCC over F1?\n"
            "  - Balanced even with class imbalance (hotspots are rare)\n"
            "  - Uses all four confusion matrix quadrants\n"
            "  - Random baseline: mean=0.001, std=0.003",
            size=11, color=DARK_GRAY)

# figure
add_figure(s, os.path.join(FIG_MAIN, "mcc_vs_f1_fairness.png"),
           Inches(7.400), Inches(1.200), Inches(5.300), Inches(5.400))


# ── SLIDE 12: Two-Stage Design ────────────────────────────────
s = content_slide(prs, "Two-Stage Experimental Design", SEC_FRAME)
add_page_number(s, 12)

# Stage 1 card
s1 = add_card(s, Inches(0.600), Inches(1.200), Inches(5.800), Inches(5.500),
              fill=RGBColor(0xE8, 0xF5, 0xE9), border=GREEN)
add_rect(s, Inches(0.600), Inches(1.200), Inches(5.800), Inches(0.600), fill=GREEN)
add_textbox(s, Inches(0.800), Inches(1.220), Inches(5.400), Inches(0.560),
            "Stage 1: Controlled Comparison", size=18, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
add_textbox(s, Inches(0.800), Inches(2.000), Inches(5.400), Inches(4.500),
            "4 pathogens\n"
            "  SARS-CoV-2, Influenza, HIV, Dengue\n\n"
            "351 scoring-detection combinations\n"
            "  9 scoring x 39 detection methods\n\n"
            "Focus areas:\n"
            "  Method comparison & ranking\n"
            "  Parameter sensitivity analysis\n"
            "  Synthetic vs real data gap\n"
            "  Initial ranking reversal observation",
            size=13, color=DARK_GRAY)

# Stage 2 card
s2 = add_card(s, Inches(6.900), Inches(1.200), Inches(5.800), Inches(5.500),
              fill=RGBColor(0xFF, 0xF3, 0xE0), border=ORANGE)
add_rect(s, Inches(6.900), Inches(1.200), Inches(5.800), Inches(0.600), fill=ORANGE)
add_textbox(s, Inches(7.100), Inches(1.220), Inches(5.400), Inches(0.560),
            "Stage 2: Large-Scale Validation", size=18, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
add_textbox(s, Inches(7.100), Inches(2.000), Inches(5.400), Inches(4.500),
            "9 pathogens (+5 more)\n"
            "  + MERS, RSV, Norovirus, HCV, Ebola\n\n"
            "2,544 evaluations\n"
            "  9 scoring x 39 detection x 9 pathogens\n\n"
            "Focus areas:\n"
            "  Two-way ANOVA (interaction effects)\n"
            "  Per-pathogen best combination\n"
            "  Friedman test + LOPO CV\n"
            "  PAHD proof of concept",
            size=13, color=DARK_GRAY)


# ── SLIDE 13: Stage 1 Results ────────────────────────────────
s = content_slide(prs, "Stage 1: Method Comparison Results", SEC_STAGE1)
add_page_number(s, 13)
add_figure(s, os.path.join(FIG_DISS, "new_methodology_comparison.png"),
           Inches(0.400), Inches(1.100), Inches(8.000), Inches(5.600))

# findings
add_textbox(s, Inches(8.700), Inches(1.200), Inches(4.000), Inches(0.400),
            "Key Findings", size=16, bold=True, color=SEC_STAGE1)
findings = [
    "Wavelet scoring achieves highest\nmean MCC across 4 pathogens",
    "HDBSCAN-based detection outperforms\ntraditional threshold methods",
    "Large variance across pathogens\nsuggests generalization limits",
    "Top combination differs between\nSARS-CoV-2 and HIV",
]
for i, f in enumerate(findings):
    cy = Inches(1.800) + i * Inches(1.250)
    card = add_card(s, Inches(8.700), cy, Inches(4.000), Inches(1.050), border=MID_GRAY)
    add_textbox(s, Inches(8.900), cy + Inches(0.100), Inches(3.600), Inches(0.850),
                f, size=11, color=DARK_GRAY)


# ── SLIDE 14: Parameter Sensitivity ──────────────────────────
s = content_slide(prs, "Parameter Sensitivity Analysis", SEC_STAGE1)
add_page_number(s, 14)
add_figure(s, os.path.join(FIG_DISS, "sensitivity_heatmap.png"),
           Inches(0.400), Inches(1.100), Inches(7.500), Inches(5.600))
add_textbox(s, Inches(8.200), Inches(1.200), Inches(4.500), Inches(0.400),
            "Sensitivity Findings", size=16, bold=True, color=SEC_STAGE1)
sens_items = [
    "Window size (scoring): moderate\nimpact, optimal varies by pathogen",
    "Epsilon (DBSCAN): high sensitivity\nrequires per-pathogen tuning",
    "Min-cluster (HDBSCAN): robust\nacross wide parameter range",
    "Threshold percentile: simple but\nshows pathogen-dependent optima",
]
for i, item in enumerate(sens_items):
    cy = Inches(1.800) + i * Inches(1.250)
    card = add_card(s, Inches(8.200), cy, Inches(4.500), Inches(1.050), border=MID_GRAY)
    add_textbox(s, Inches(8.400), cy + Inches(0.100), Inches(4.100), Inches(0.850),
                item, size=11, color=DARK_GRAY)


# ── SLIDE 15: Synthetic vs Real Gap ──────────────────────────
s = content_slide(prs, "Synthetic vs Real Data Gap", SEC_STAGE1)
add_page_number(s, 15)

# table header
add_rect(s, Inches(0.600), Inches(1.300), Inches(12.100), Inches(0.500), fill=DARK_NAVY)
cols = ["Aspect", "Synthetic Data", "Real Data"]
col_ws = [Inches(3.000), Inches(4.500), Inches(4.600)]
cx = Inches(0.600)
for c, w in zip(cols, col_ws):
    add_textbox(s, cx, Inches(1.310), w, Inches(0.480),
                c, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    cx += w

# table rows
rows = [
    ("Mutation distribution", "Uniform / Poisson", "Highly skewed, founder effects"),
    ("Selection pressure", "None (neutral)", "Complex positive + purifying"),
    ("Ground truth", "Known by construction", "Approximate (DMS / literature)"),
    ("Phylo correlation", "Independent samples", "Strong phylogenetic non-independence"),
    ("Method ranking", "Threshold-based wins", "Clustering-based wins"),
]
for r, (asp, syn, real) in enumerate(rows):
    ry = Inches(1.800) + r * Inches(0.700)
    fill = LIGHT_GRAY if r % 2 == 0 else WHITE
    add_rect(s, Inches(0.600), ry, Inches(12.100), Inches(0.700), fill=fill)
    cx = Inches(0.600)
    for val, w in zip([asp, syn, real], col_ws):
        add_textbox(s, cx, ry, w, Inches(0.700), val, size=12, color=DARK_GRAY,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += w

# key message
msg = add_card(s, Inches(2.000), Inches(5.600), Inches(9.300), Inches(0.800),
               fill=RGBColor(0xFF, 0xEB, 0xEE), border=RED_ACCENT)
set_text(msg, "Ranking reversals between synthetic and real data demonstrate the need for real-pathogen benchmarking",
         size=13, bold=True, color=RED_ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ── SLIDE 16: Ranking Reversal ───────────────────────────────
s = content_slide(prs, "Cross-Pathogen Ranking Reversal", SEC_STAGE1)
add_page_number(s, 16)
add_figure(s, os.path.join(FIG_MAIN, "detection_ranking.png"),
           Inches(0.400), Inches(1.100), Inches(7.800), Inches(5.600))

# big number
big = add_card(s, Inches(8.700), Inches(1.300), Inches(3.800), Inches(2.200),
               fill=DARK_NAVY)
add_textbox(s, Inches(8.700), Inches(1.500), Inches(3.800), Inches(1.000),
            "9 / 9", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(8.700), Inches(2.500), Inches(3.800), Inches(0.800),
            "unique best\ncombinations", size=16, color=RGBColor(0xBB, 0xDE, 0xFB),
            align=PP_ALIGN.CENTER)

add_textbox(s, Inches(8.500), Inches(3.800), Inches(4.200), Inches(2.800),
            "Every pathogen has a different\noptimal scoring + detection pair.\n\n"
            "No single method can be\n\"recommended\" universally.\n\n"
            "This motivates the need for\npathogen-adaptive selection.",
            size=12, color=DARK_GRAY)


# ── SLIDE 17: 9-Pathogen Scale ───────────────────────────────
s = content_slide(prs, "Stage 2: Scaling to 9 Pathogens", SEC_STAGE2)
add_page_number(s, 17)

bignums = [
    ("9", "Pathogens", "SARS-CoV-2, Influenza, HIV,\nDengue, MERS, RSV,\nNorovirus, HCV, Ebola"),
    ("2,544", "Evaluations", "9 scoring methods\nx 39 detection methods\nx 9 pathogens (multilayer GT)"),
    ("4", "Statistical Tests", "Two-way ANOVA\nPermutation test\nFriedman test\nLOPO cross-validation"),
]
for i, (num, label, desc) in enumerate(bignums):
    cx = Inches(0.600) + i * Inches(4.200)
    cy = Inches(1.500)
    card = add_card(s, cx, cy, Inches(3.900), Inches(4.800), border=MID_GRAY)
    add_textbox(s, cx, cy + Inches(0.300), Inches(3.900), Inches(1.200),
                num, size=54, bold=True, color=SEC_STAGE2, align=PP_ALIGN.CENTER)
    add_textbox(s, cx, cy + Inches(1.500), Inches(3.900), Inches(0.500),
                label, size=20, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)
    add_textbox(s, cx + Inches(0.300), cy + Inches(2.200), Inches(3.300), Inches(2.400),
                desc, size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ── SLIDE 18: ANOVA ──────────────────────────────────────────
s = content_slide(prs, "Two-Way ANOVA: Variance Decomposition", SEC_STAGE2)
add_page_number(s, 18)
add_figure(s, os.path.join(FIG_MAIN, "variance_decomposition.png"),
           Inches(0.400), Inches(1.100), Inches(7.200), Inches(5.600))

# omega-squared cards
anova_cards = [
    ("Detection", "~35%", "Largest effect:\nchoice of detection\nmethod matters most"),
    ("Interaction", "28.5%", "Scoring x Pathogen\ninteraction is strong:\nno universal best"),
    ("Scoring", "~8.3%", "Scoring method\ncontributes less than\ndetection or interaction"),
]
for i, (ttl, val, desc) in enumerate(anova_cards):
    cy = Inches(1.200) + i * Inches(1.850)
    card = add_card(s, Inches(7.900), cy, Inches(4.800), Inches(1.600), border=MID_GRAY)
    add_textbox(s, Inches(8.100), cy + Inches(0.100), Inches(1.500), Inches(0.400),
                ttl, size=12, bold=True, color=SEC_STAGE2)
    add_textbox(s, Inches(8.100), cy + Inches(0.500), Inches(1.500), Inches(0.800),
                val, size=28, bold=True, color=SEC_STAGE2)
    add_textbox(s, Inches(9.800), cy + Inches(0.200), Inches(2.700), Inches(1.200),
                desc, size=11, color=DARK_GRAY)


# ── SLIDE 19: Per-Pathogen Best ───────────────────────────────
s = content_slide(prs, "Per-Pathogen Best Combinations", SEC_STAGE2)
add_page_number(s, 19)
add_figure(s, os.path.join(FIG_MAIN, "cross_pathogen_top5.png"),
           Inches(0.300), Inches(1.100), Inches(8.500), Inches(5.600))

# annotation
big = add_card(s, Inches(9.100), Inches(1.300), Inches(3.600), Inches(2.000),
               fill=SEC_STAGE2)
add_textbox(s, Inches(9.100), Inches(1.400), Inches(3.600), Inches(1.000),
            "9 / 9", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(9.100), Inches(2.300), Inches(3.600), Inches(0.800),
            "unique best\ncombinations", size=16, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(s, Inches(9.100), Inches(3.600), Inches(3.600), Inches(3.000),
            "Each pathogen requires a\ndifferent scoring + detection pair.\n\n"
            "High entropy-function divergence\n(e.g., Norovirus r=0.39 vs\nSARS-CoV-2 r=0.03)\n\n"
            "H-score founder bias\n(rho = -0.876)",
            size=11, color=DARK_GRAY)


# ── SLIDE 20: Friedman + LOPO ─────────────────────────────────
s = content_slide(prs, "Friedman Test & Leave-One-Pathogen-Out CV", SEC_STAGE2)
add_page_number(s, 20)

# Friedman card
fc = add_card(s, Inches(0.600), Inches(1.200), Inches(5.800), Inches(5.400),
              fill=WHITE, border=MID_GRAY)
add_rect(s, Inches(0.600), Inches(1.200), Inches(5.800), Inches(0.600), fill=SEC_STAGE2)
add_textbox(s, Inches(0.800), Inches(1.220), Inches(5.400), Inches(0.560),
            "Friedman Rank Test (Top-20)", size=18, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
add_textbox(s, Inches(0.800), Inches(2.100), Inches(5.400), Inches(4.200),
            "Test: Are top-20 method rankings\nconsistent across pathogens?\n\n"
            "Result:\n"
            "  chi-squared = 42.44\n"
            "  p-value = 0.0015\n\n"
            "Interpretation:\n"
            "  Rankings differ significantly\n"
            "  across pathogens (p < 0.01)\n\n"
            "  No single method maintains\n"
            "  top rank across all pathogens",
            size=13, color=DARK_GRAY)

# LOPO card
lc = add_card(s, Inches(6.900), Inches(1.200), Inches(5.800), Inches(5.400),
              fill=WHITE, border=MID_GRAY)
add_rect(s, Inches(6.900), Inches(1.200), Inches(5.800), Inches(0.600), fill=RED_ACCENT)
add_textbox(s, Inches(7.100), Inches(1.220), Inches(5.400), Inches(0.560),
            "LOPO Cross-Validation", size=18, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
add_textbox(s, Inches(7.100), Inches(2.100), Inches(5.400), Inches(4.200),
            "Protocol: Train on 8 pathogens,\ntest on held-out 1 pathogen\n\n"
            "Question: Can we predict the best\nmethod for an unseen pathogen?\n\n"
            "Result:\n"
            "  0 / 9 correct predictions\n\n"
            "Interpretation:\n"
            "  Simple majority voting fails\n"
            "  completely for new pathogens\n\n"
            "  Motivates PAHD: profile-based\n"
            "  adaptive selection needed",
            size=13, color=DARK_GRAY)


# ── SLIDE 21: ESM-2 Validation ───────────────────────────────
s = content_slide(prs, "ESM-2 Protein Language Model Validation", SEC_STAGE2)
add_page_number(s, 21)

# table
add_rect(s, Inches(1.500), Inches(1.400), Inches(10.300), Inches(0.500), fill=DARK_NAVY)
esm_cols = ["Metric", "Without ESM-2", "With ESM-2", "Change"]
esm_ws = [Inches(2.800), Inches(2.500), Inches(2.500), Inches(2.500)]
cx = Inches(1.500)
for c, w in zip(esm_cols, esm_ws):
    add_textbox(s, cx, Inches(1.410), w, Inches(0.480),
                c, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    cx += w

esm_rows = [
    ("Mean MCC (top-20)", "0.347", "0.352", "+0.005"),
    ("Median MCC (top-20)", "0.338", "0.341", "+0.003"),
    ("Best single combo", "0.390", "0.393", "+0.003"),
    ("Coverage (>0.3 MCC)", "12/20", "13/20", "+1"),
    ("Rank correlation", "--", "rho=0.94", "High agreement"),
]
for r, (m, wo, wi, ch) in enumerate(esm_rows):
    ry = Inches(1.900) + r * Inches(0.650)
    fill = LIGHT_GRAY if r % 2 == 0 else WHITE
    add_rect(s, Inches(1.500), ry, Inches(10.300), Inches(0.650), fill=fill)
    cx = Inches(1.500)
    for val, w in zip([m, wo, wi, ch], esm_ws):
        clr = GREEN if val.startswith("+") or val == "High agreement" else DARK_GRAY
        add_textbox(s, cx, ry, w, Inches(0.650), val, size=12, color=clr,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += w

# conclusion box
msg = add_card(s, Inches(1.500), Inches(5.400), Inches(10.300), Inches(1.000),
               fill=RGBColor(0xE8, 0xF5, 0xE9), border=GREEN)
set_text(msg, "ESM-2 embeddings provide marginal improvement, confirming that\n"
         "sequence-based scoring already captures most of the relevant signal.\n"
         "The benchmark rankings remain stable (rho=0.94).",
         size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ── SLIDE 22: KEY FINDING ─────────────────────────────────────
s = content_slide(prs, "KEY FINDING: No Universal Best Method", SEC_STAGE2)
add_page_number(s, 22)

# central message
central = add_card(s, Inches(2.000), Inches(1.200), Inches(9.300), Inches(1.200),
                   fill=DARK_NAVY)
set_text(central, "No single scoring + detection combination works best for all pathogens",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 4 evidence cards
evidence = [
    ("ANOVA Interaction", "omega-sq = 0.285", "Scoring x Pathogen interaction\nexplains 28.5% of variance"),
    ("Unique Best Combos", "9 / 9", "Every pathogen has a different\noptimal method pair"),
    ("Friedman Test", "p = 0.0015", "Rankings differ significantly\nacross pathogens"),
    ("LOPO CV", "0 / 9 correct", "Majority voting fails completely\nfor unseen pathogens"),
]
for i, (ttl, val, desc) in enumerate(evidence):
    col = i % 2
    row = i // 2
    cx = Inches(0.600) + col * Inches(6.300)
    cy = Inches(2.700) + row * Inches(2.200)
    card = add_card(s, cx, cy, Inches(5.800), Inches(1.950), border=MID_GRAY)
    add_textbox(s, cx + Inches(0.200), cy + Inches(0.100), Inches(5.400), Inches(0.350),
                ttl, size=13, bold=True, color=SEC_STAGE2)
    add_textbox(s, cx + Inches(0.200), cy + Inches(0.500), Inches(5.400), Inches(0.600),
                val, size=28, bold=True, color=DARK_NAVY)
    add_textbox(s, cx + Inches(0.200), cy + Inches(1.150), Inches(5.400), Inches(0.700),
                desc, size=11, color=DARK_GRAY)


# ── SLIDE 23: PAHD Proof of Concept ──────────────────────────
s = content_slide(prs, "PAHD: Pathogen-Adaptive Hotspot Detection", SEC_EXT)
add_page_number(s, 23)

# 3-step flow
steps = [
    ("Step 1: Profile", "Extract pathogen profile\n(entropy distribution, genome\nlength, mutation rate, diversity)"),
    ("Step 2: Match", "Compare profile against\nMutBench knowledge base\n(similarity-based retrieval)"),
    ("Step 3: Select", "Recommend optimal scoring +\ndetection combination based\non matched pathogen results"),
]
for i, (ttl, desc) in enumerate(steps):
    cx = Inches(0.600) + i * Inches(4.200)
    cy = Inches(1.300)
    card = add_card(s, cx, cy, Inches(3.900), Inches(2.500), border=MID_GRAY)
    add_rect(s, cx, cy, Inches(3.900), Inches(0.500), fill=PURPLE)
    add_textbox(s, cx, cy + Inches(0.050), Inches(3.900), Inches(0.400),
                ttl, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, cx + Inches(0.200), cy + Inches(0.650), Inches(3.500), Inches(1.700),
                desc, size=12, color=DARK_GRAY)
    if i < 2:
        add_textbox(s, cx + Inches(3.900), cy + Inches(0.700), Inches(0.300), Inches(0.500),
                    ">", size=24, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

# results
res = add_card(s, Inches(0.600), Inches(4.200), Inches(12.100), Inches(2.500),
               fill=RGBColor(0xF3, 0xE5, 0xF5), border=PURPLE)
add_textbox(s, Inches(0.800), Inches(4.300), Inches(11.700), Inches(0.400),
            "Proof of Concept Results", size=16, bold=True, color=PURPLE)
add_textbox(s, Inches(0.800), Inches(4.800), Inches(5.500), Inches(1.700),
            "PAHD approach:\n"
            "  Profile-based selection achieves\n"
            "  higher mean MCC than any fixed method\n\n"
            "  Outperforms MutBench top-1 on 6/9 pathogens\n"
            "  when using leave-one-out profile matching",
            size=12, color=DARK_GRAY)
add_textbox(s, Inches(6.500), Inches(4.800), Inches(5.500), Inches(1.700),
            "Baseline comparison:\n"
            "  MutBench best:     MCC = 0.390\n"
            "  FreqThresh:          MCC = 0.248\n"
            "  SWAN:                  MCC = 0.193\n"
            "  MutClust-Orig:      MCC = 0.138\n"
            "  Random baseline:  MCC = 0.001",
            size=12, color=DARK_GRAY)


# ── SLIDE 24: Phylo + Region-overlap ─────────────────────────
s = content_slide(prs, "Phylogenetic Correction & Region Overlap", SEC_EXT)
add_page_number(s, 24)

# Phylo card
pc = add_card(s, Inches(0.600), Inches(1.200), Inches(5.800), Inches(5.400), border=MID_GRAY)
add_rect(s, Inches(0.600), Inches(1.200), Inches(5.800), Inches(0.600), fill=PURPLE)
add_textbox(s, Inches(0.800), Inches(1.220), Inches(5.400), Inches(0.560),
            "Phylogenetic Non-Independence", size=16, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
add_textbox(s, Inches(0.800), Inches(2.000), Inches(5.400), Inches(4.300),
            "Problem:\n"
            "  Closely related sequences share mutations\n"
            "  by descent, not independent selection.\n"
            "  D614G spread by founder effect,\n"
            "  not necessarily positive selection.\n\n"
            "Approach:\n"
            "  TreeTime phylogenetic reconstruction\n"
            "  Weight mutations by branch independence\n"
            "  Down-weight redundant lineage mutations\n\n"
            "Status:\n"
            "  Acknowledged as key limitation\n"
            "  Partial correction in MutBench v1\n"
            "  Full integration planned for future work",
            size=12, color=DARK_GRAY)

# Region overlap card
rc = add_card(s, Inches(6.900), Inches(1.200), Inches(5.800), Inches(5.400), border=MID_GRAY)
add_rect(s, Inches(6.900), Inches(1.200), Inches(5.800), Inches(0.600), fill=TEAL_ACCENT)
add_textbox(s, Inches(7.100), Inches(1.220), Inches(5.400), Inches(0.560),
            "Region-Level Overlap Analysis", size=16, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
add_textbox(s, Inches(7.100), Inches(2.000), Inches(5.400), Inches(4.300),
            "Motivation:\n"
            "  Position-level MCC may be too strict.\n"
            "  Methods may detect the right region\n"
            "  but miss exact positions.\n\n"
            "Region-overlap metric:\n"
            "  Window-based overlap scoring\n"
            "  Relaxed evaluation with +/- N positions\n\n"
            "Results:\n"
            "  MCC improves from 0.289 to 0.712\n"
            "  when using region-level evaluation\n\n"
            "  Suggests methods capture hotspot regions\n"
            "  even when exact positions differ",
            size=12, color=DARK_GRAY)


# ── SLIDE 25: Baseline Comparison ─────────────────────────────
s = content_slide(prs, "Scoring-Detection Heatmap & Baseline Comparison", SEC_EXT)
add_page_number(s, 25)
add_figure(s, os.path.join(FIG_MAIN, "scoring_detection_heatmap.png"),
           Inches(0.300), Inches(1.100), Inches(7.800), Inches(5.600))

# MCC comparison bars (text-based)
add_textbox(s, Inches(8.500), Inches(1.200), Inches(4.200), Inches(0.400),
            "MCC Comparison", size=16, bold=True, color=SEC_EXT)
bars = [
    ("MutBench Best", "0.390", GREEN),
    ("FreqThresh", "0.248", BLUE),
    ("SWAN", "0.193", ORANGE),
    ("MutClust-Orig", "0.138", RED_ACCENT),
    ("Random", "0.001", GRAY),
]
for i, (name, val, color) in enumerate(bars):
    cy = Inches(1.800) + i * Inches(1.000)
    add_textbox(s, Inches(8.500), cy, Inches(2.500), Inches(0.350),
                name, size=12, bold=True, color=color)
    bar_w = float(val) / 0.4 * 4.0  # scale
    add_rect(s, Inches(8.500), cy + Inches(0.380), Inches(bar_w), Inches(0.350), fill=color)
    add_textbox(s, Inches(8.500) + Inches(bar_w) + Inches(0.100), cy + Inches(0.380),
                Inches(1.000), Inches(0.350), val, size=11, bold=True, color=color,
                anchor=MSO_ANCHOR.MIDDLE)


# ── SLIDE 26: Discussion ──────────────────────────────────────
s = content_slide(prs, "Discussion", SEC_CONCL)
add_page_number(s, 26)

disc = [
    ("Benchmarking Matters",
     "Without standardized evaluation, published methods appear to perform well in isolation "
     "but fail under cross-pathogen comparison. MutBench provides the missing infrastructure "
     "for fair, reproducible comparison of hotspot detection methods."),
    ("Method Selection is Pathogen-Dependent",
     "The strong interaction effect (omega-sq = 0.285) means that recommending a single method "
     "is fundamentally flawed. Pathogen-specific characteristics (genome structure, mutation rate, "
     "selection landscape) determine which approach works best."),
    ("From Benchmark to Algorithm",
     "MutBench is not just a benchmark but a knowledge base. The systematic evaluation across "
     "9 pathogens enables PAHD: data-driven, profile-based method selection that outperforms "
     "any fixed approach."),
]
for i, (ttl, desc) in enumerate(disc):
    cy = Inches(1.200) + i * Inches(1.850)
    card = add_card(s, Inches(0.600), cy, Inches(12.100), Inches(1.650), border=MID_GRAY)
    colors = [SEC_FRAME, SEC_STAGE2, PURPLE]
    add_accent_bar(s, Inches(0.600), cy, Inches(0.080), Inches(1.650), colors[i])
    add_textbox(s, Inches(0.900), cy + Inches(0.120), Inches(11.500), Inches(0.400),
                ttl, size=16, bold=True, color=colors[i])
    add_textbox(s, Inches(0.900), cy + Inches(0.580), Inches(11.500), Inches(0.950),
                desc, size=12, color=DARK_GRAY)


# ── SLIDE 27: Contributions + Limitations + Future ───────────
s = content_slide(prs, "Contributions, Limitations & Future Work", SEC_CONCL)
add_page_number(s, 27)

cols_data = [
    ("Contributions", GREEN, [
        "First systematic viral hotspot\ndetection benchmark (9 pathogens)",
        "3-layer ground truth design\n(adaptive + constrained + DMS)",
        "Statistical proof: no universal\nbest method (interaction eta-sq=0.285)",
        "PAHD proof of concept for\nadaptive method selection",
    ]),
    ("Limitations", RED_ACCENT, [
        "Ground truth approximation\n(no perfect gold standard)",
        "Phylogenetic non-independence\nnot fully corrected",
        "Limited to substitution mutations\n(no indels, recombination)",
        "PAHD is proof-of-concept,\nnot production-ready",
    ]),
    ("Future Work", BLUE, [
        "Full PAHD implementation with\nlearned pathogen profiles",
        "Integrate phylogenetic correction\n(TreeTime branch weighting)",
        "Extend to indels, recombination,\nand structural variants",
        "Real-time surveillance integration\nfor emerging pathogens",
    ]),
]
for ci, (ttl, color, items) in enumerate(cols_data):
    cx = Inches(0.400) + ci * Inches(4.250)
    # header
    add_rect(s, cx, Inches(1.200), Inches(4.000), Inches(0.500), fill=color)
    add_textbox(s, cx, Inches(1.210), Inches(4.000), Inches(0.480),
                ttl, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    for j, item in enumerate(items):
        iy = Inches(1.850) + j * Inches(1.250)
        card = add_card(s, cx, iy, Inches(4.000), Inches(1.050), border=MID_GRAY)
        add_textbox(s, cx + Inches(0.150), iy + Inches(0.100), Inches(3.700), Inches(0.850),
                    item, size=11, color=DARK_GRAY)


# ── SLIDE 28: Key Takeaway ────────────────────────────────────
s = content_slide(prs, "Key Takeaway", SEC_CONCL)
add_page_number(s, 28)

# central message
central = add_card(s, Inches(1.500), Inches(1.300), Inches(10.300), Inches(1.600),
                   fill=DARK_NAVY)
set_text(central,
         "MutBench transforms viral hotspot detection from\n"
         "ad-hoc single-virus evaluation to systematic cross-pathogen benchmarking",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 4 evidence boxes
evidences = [
    ("9 Pathogens", "Largest cross-pathogen\nbenchmark for hotspot\ndetection"),
    ("2,544 Evaluations", "Comprehensive coverage\nof scoring x detection\ncombinations"),
    ("omega-sq = 0.285", "Strong interaction proves\npathogen-adaptive\nselection is needed"),
    ("PAHD Concept", "From benchmark insight\nto actionable algorithm\nfor method selection"),
]
for i, (ttl, desc) in enumerate(evidences):
    cx = Inches(0.500) + i * Inches(3.200)
    cy = Inches(3.300)
    card = add_card(s, cx, cy, Inches(2.900), Inches(3.200), border=MID_GRAY)
    colors = [SEC_FRAME, SEC_STAGE2, PURPLE, GREEN]
    add_textbox(s, cx, cy + Inches(0.200), Inches(2.900), Inches(0.500),
                ttl, size=16, bold=True, color=colors[i], align=PP_ALIGN.CENTER)
    add_textbox(s, cx + Inches(0.200), cy + Inches(0.850), Inches(2.500), Inches(2.100),
                desc, size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ── SLIDE 29: Thank You ──────────────────────────────────────
s = add_blank(prs)
# dark top area
add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(3.500), fill=DARK_NAVY)
# KNU logo
s.shapes.add_picture(LOGO, Inches(0.600), Inches(0.300), Inches(0.650), Inches(0.650))
add_textbox(s, Inches(0.600), Inches(1.200), Inches(12.0), Inches(1.200),
            "Thank You", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.600), Inches(2.400), Inches(12.0), Inches(0.600),
            "Questions & Discussion", size=20, color=RGBColor(0xBB, 0xDE, 0xFB),
            align=PP_ALIGN.CENTER)

# contact info cards
add_textbox(s, Inches(3.000), Inches(4.200), Inches(7.300), Inches(0.400),
            "Hwijun Kwon", size=22, bold=True, color=DARK_NAVY, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(3.000), Inches(4.700), Inches(7.300), Inches(0.350),
            "Dept. of Computer Science, Kyungpook National University",
            size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(3.000), Inches(5.200), Inches(7.300), Inches(0.350),
            "Advisor: Prof. Inuk Jung", size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(3.000), Inches(5.700), Inches(7.300), Inches(0.350),
            "March 2026", size=14, color=GRAY, align=PP_ALIGN.CENTER)

# bottom bar
add_rect(s, Inches(0), Inches(7.120), SLIDE_W, Inches(0.380), fill=DARK_NAVY)
add_textbox(s, Inches(0.400), Inches(7.140), Inches(10.0), Inches(0.340),
            "Kyungpook National University  |  Dept. of Computer Science  |  Hwijun Kwon",
            size=10, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)


# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Total slides: {len(prs.slides)}")
