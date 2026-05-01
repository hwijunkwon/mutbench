from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation


ROOT = Path(".")
REPORT = Path("docs/final_adversarial_audit_2026-04-27.md")
PPTX = Path("paper/dissertation/presentation/mutbench_defense_pahd_r.pptx")

TEXT_SOURCES = [
    Path("paper/dissertation/chapters_en/ch4_results.tex"),
    Path("paper/dissertation/chapters_en/ch5_discussion.tex"),
    Path("paper/dissertation/chapters_en/ch6_conclusion.tex"),
    Path("paper/dissertation/dissertation_easy_guide_v2.md"),
    Path("docs/pahd_r_easy_guide.md"),
    Path("docs/pahd_r_thesis_easyguide_sync.md"),
    Path("scripts/build_defense_pptx.py"),
]


@dataclass
class Finding:
    status: str
    check: str
    detail: str


def read_text_sources() -> dict[str, str]:
    out: dict[str, str] = {}
    for path in TEXT_SOURCES:
        if path.exists():
            out[str(path)] = path.read_text(encoding="utf-8", errors="replace")
        else:
            out[str(path)] = ""
    return out


def pptx_text_and_stats(path: Path) -> tuple[str, int, int]:
    if not path.exists():
        return "", 0, 0
    prs = Presentation(path)
    texts: list[str] = []
    pictures = 0
    for i, slide in enumerate(prs.slides, start=1):
        texts.append(f"\n[PPTX slide {i}]\n")
        for shape in slide.shapes:
            if getattr(shape, "shape_type", None) == 13:
                pictures += 1
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n".join(texts), len(prs.slides), pictures


def snippets(pattern: str, text: str, window: int = 200) -> list[str]:
    hits = []
    for m in re.finditer(pattern, text, flags=re.I):
        start = max(0, m.start() - window)
        end = min(len(text), m.end() + window)
        snippet = re.sub(r"\s+", " ", text[start:end]).strip()
        hits.append(snippet)
    return hits


def safe_negated(snippet: str) -> bool:
    s = snippet.lower()
    guards = [
        "not ",
        "should not",
        "do not",
        "never",
        "candidate only",
        "candidate-only",
        "not final",
        "not adopted",
        "remain candidate",
        "remains candidate",
        "future-validation",
        "reject",
        "avoid",
        "what not to claim",
        "do not claim",
        "not a universal",
        "not primarily",
        "not an all-pathogen",
        "without claiming",
        "not as a completed",
        "rather than adopted",
        "crosses zero",
        "말하면 안",
        "안 되는 표현",
        "피해야",
        "금지",
    ]
    return any(g in s for g in guards)


def add_presence(findings: list[Finding], name: str, corpus: str, required: list[str]) -> None:
    missing = [x for x in required if x not in corpus]
    if missing:
        findings.append(Finding("WARN", name, "Missing exact strings: " + ", ".join(missing)))
    else:
        findings.append(Finding("PASS", name, "All required numeric/reporting strings are present."))


def add_risk_check(findings: list[Finding], name: str, pattern: str, corpus: str) -> None:
    hits = snippets(pattern, corpus)
    risky = [h for h in hits if not safe_negated(h)]
    if risky:
        findings.append(Finding("FAIL", name, "Potential overclaim: " + " | ".join(risky[:3])))
    elif hits:
        findings.append(Finding("PASS", name, f"{len(hits)} occurrence(s), all in guarded/candidate context."))
    else:
        findings.append(Finding("PASS", name, "No risky occurrence found."))


def main() -> None:
    source_text = read_text_sources()
    deck_text, slides, pictures = pptx_text_and_stats(PPTX)
    corpus = "\n".join(source_text.values()) + "\n" + deck_text

    findings: list[Finding] = []

    if slides == 14 and pictures >= 3:
        findings.append(Finding("PASS", "Designed PPTX integrity", f"{slides} slides; {pictures} embedded figure(s)."))
    else:
        findings.append(Finding("WARN", "Designed PPTX integrity", f"{slides} slides; {pictures} embedded figure(s)."))

    add_presence(
        findings,
        "Adopted PAHD-R modes",
        corpus,
        [
            "0.1689",
            "0.5866",
            "0.2111",
            "0.1108",
            "0.1680",
            "0.6079",
            "0.1944",
            "0.1128",
            "0.1666",
            "0.6366",
            "0.1833",
            "0.0696",
        ],
    )
    add_presence(
        findings,
        "Candidate variant metrics",
        corpus,
        ["0.1770", "0.6085", "0.0947", "0.2742", "0.7276", "0.3500", "0.0354", "4/9"],
    )
    add_presence(
        findings,
        "Repair-layer metrics",
        corpus,
        ["0.5257", "0.0000", "0.2304", "0.8761", "0.1000", "0.0065"],
    )

    add_risk_check(findings, "Universal-AI overclaim", r"universal\s+AI\s+predictor|universal\s+adaptive\s+(?:algorithm|predictor)", corpus)
    add_risk_check(findings, "AI-primary framing", r"primarily\s+AI[- ]based|AI[- ]based\s+predictor", corpus)
    add_risk_check(findings, "Virus-specific algorithm framing", r"virus[- ]specific\s+algorithm(?:s| collection)?", corpus)
    add_risk_check(findings, "SARS-CoV-2/MERS adoption overclaim", r"(?:SARS-CoV-2|MERS).{0,120}(?:adopted|finalized|default)", corpus)
    add_risk_check(findings, "Core-Calibrated default overclaim", r"Core[- ]Calibrated.{0,120}(?:default|adopted)", corpus)
    add_risk_check(findings, "Selective all-pathogen overclaim", r"selective.{0,120}all[- ]pathogen", corpus)
    add_risk_check(findings, "Region-only reporting overclaim", r"(?:\+/-10|±10|region MCC).{0,120}(?:alone|standalone|proves)", corpus)
    add_risk_check(findings, "Pooled-permutation overclaim", r"pooled permutation.{0,120}(?:every|all pathogen|full pathogen-level)", corpus)
    add_risk_check(findings, "Raw/time robustness overclaim", r"(?:raw|time[- ]forward|temporal).{0,120}(?:proven|validated|adopted default)", corpus)

    status_order = {"FAIL": 0, "WARN": 1, "PASS": 2}
    worst = min(findings, key=lambda f: status_order[f.status]).status if findings else "WARN"
    blocking = any(f.status == "FAIL" for f in findings)

    lines = [
        "# Final adversarial audit - 2026-04-27",
        "",
        f"Overall: {'BLOCKING ISSUES FOUND' if blocking else 'PASS with non-blocking caveats'}",
        "",
        "## Scope",
        "",
        "- English thesis result/discussion/conclusion chapters",
        "- Korean easy guide v2 and concise PAHD-R guide",
        "- PAHD-R thesis/easy-guide sync note",
        "- Designed PowerPoint source script and generated PPTX text",
        "",
        "## Findings",
        "",
    ]
    for f in findings:
        lines.append(f"- **{f.status}** `{f.check}`: {f.detail}")
    lines.extend(
        [
            "",
            "## Residual caveats",
            "",
            "- The current claim is internally audited, but stronger adaptive-learning claims still require additional pathogens or time-forward external labels.",
            "- SARS-CoV-2 and MERS repair layers must remain candidate-only until external/structural validation is completed.",
            "- Selective callability is a triage/reject-option result, not an all-pathogen benchmark result.",
            "- Pooled permutation is acceptable as a global sanity check, but not as proof that every pathogen-level null is solved.",
            "",
            f"Machine status: {worst}",
            "",
        ]
    )
    REPORT.parent.mkdir(parents=True, exist_ok=True)
    REPORT.write_text("\n".join(lines), encoding="utf-8")
    print(REPORT)
    print(f"overall={'FAIL' if blocking else 'PASS'} slides={slides} pictures={pictures}")


if __name__ == "__main__":
    main()
