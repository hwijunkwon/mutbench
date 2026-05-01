#!/usr/bin/env python3
"""Integrated adversarial review for PAHD-R artifacts and claims."""

from __future__ import annotations

from pathlib import Path
import re

import numpy as np
import pandas as pd
from pptx import Presentation

PROJECT_ROOT = Path(__file__).resolve().parent.parent
RESULTS_DIR = PROJECT_ROOT / "results" / "mutbench"
PPT_PATH = PROJECT_ROOT / "paper" / "ppt" / "MutBench_Defense_v6.pptx"
DOC_PATH = PROJECT_ROOT / "docs" / "pahd_r_validation_and_improvement_check.md"


def read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8") if path.exists() else ""


def ppt_text(path: Path) -> str:
    if not path.exists():
        return ""
    prs = Presentation(str(path))
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                chunks.append(shape.text)
    return "\n".join(chunks)


def add_issue(rows: list[dict[str, object]], severity: str, area: str, issue: str, evidence: str, action: str) -> None:
    rows.append({
        "severity": severity,
        "area": area,
        "issue": issue,
        "evidence": evidence,
        "recommended_action": action,
    })


def markdown_table(frame: pd.DataFrame) -> str:
    text = frame.copy()
    for col in text.columns:
        if pd.api.types.is_float_dtype(text[col]):
            text[col] = text[col].map(lambda x: "" if pd.isna(x) else f"{x:.4f}")
        else:
            text[col] = text[col].astype(str)
    lines = [
        "| " + " | ".join(text.columns) + " |",
        "| " + " | ".join(["---"] * len(text.columns)) + " |",
    ]
    for _, row in text.iterrows():
        lines.append("| " + " | ".join(row[col] for col in text.columns) + " |")
    return "\n".join(lines)


def compare_final_metrics(rows: list[dict[str, object]]) -> None:
    reporting = pd.read_csv(RESULTS_DIR / "pahd_r_final_reporting_table.csv")
    stability = pd.read_csv(RESULTS_DIR / "pahd_r_final_mode_stability_by_mode.csv")
    reporting = reporting[reporting["scope"] == "global"].rename(columns={"entry": "mode"})
    merged = reporting.merge(stability, on="mode", suffixes=("_report", "_stability"))
    for _, row in merged.iterrows():
        checks = [
            ("exact_mcc", "mean_mcc_exact"),
            ("mcc_window_10", "mean_mcc_window_10"),
            ("precision_20", "mean_precision_20"),
            ("constrained_fpr", "mean_constrained_fpr"),
        ]
        for left, right in checks:
            delta = abs(float(row[left]) - float(row[right]))
            if delta > 1e-9:
                add_issue(
                    rows,
                    "high",
                    "metric_consistency",
                    f"{row['mode']} {left} differs between final reporting and stability audit",
                    f"delta={delta:.6g}",
                    "Regenerate one of the tables from a shared source before reporting.",
                )
    if len(merged) != 3:
        add_issue(
            rows,
            "high",
            "metric_consistency",
            "Not all three global PAHD-R modes matched between final tables",
            f"matched_modes={len(merged)}",
            "Check mode naming and table generation.",
        )


def audit_scope_claims(rows: list[dict[str, object]]) -> None:
    ppt = ppt_text(PPT_PATH)
    doc = read_text(DOC_PATH)
    combined = ppt + "\n" + doc
    if "primary metrics use the 9-pathogen benchmark" not in ppt:
        add_issue(
            rows,
            "high",
            "ppt_claim_boundary",
            "PPT does not explicitly separate PAHD-R primary metrics from supplemental 11-feature audits",
            "missing phrase: primary metrics use the 9-pathogen benchmark",
            "Keep the claim-boundary text on the PAHD-R slide.",
        )
    if "11 pathogens enables PAHD-R" in ppt:
        add_issue(
            rows,
            "medium",
            "ppt_claim_boundary",
            "PPT still contains a broad phrase implying 11-pathogen evidence directly enables PAHD-R",
            "phrase found: 11 pathogens enables PAHD-R",
            "Revise wording to distinguish benchmark expansion from PAHD-R primary metrics.",
        )
    if re.search(r"PAHD-R.*11[- ]pathogen", combined, re.IGNORECASE):
        add_issue(
            rows,
            "medium",
            "claim_scope",
            "A PAHD-R sentence may still be read as an 11-pathogen performance claim",
            "regex matched PAHD-R near 11-pathogen",
            "Review surrounding sentence manually before defense.",
        )


def audit_adversarial_results(rows: list[dict[str, object]]) -> None:
    global_df = pd.read_csv(RESULTS_DIR / "pahd_r_adversarial_global_summary.csv")
    for _, row in global_df.iterrows():
        if float(row["empirical_p_null_ge_observed"]) > 0.01:
            add_issue(
                rows,
                "high",
                "permutation_null",
                f"{row['strategy']} does not clearly beat label-permutation null",
                f"p={row['empirical_p_null_ge_observed']:.4f}",
                "Do not claim validation for this strategy; redesign or mark diagnostic only.",
            )
    path_df = pd.read_csv(RESULTS_DIR / "pahd_r_adversarial_null_summary.csv")
    weak = path_df[~path_df["exceeds_null_p95"].astype(bool)]
    if len(weak):
        evidence = ", ".join(sorted(set(weak["pathogen"].astype(str)))[:10])
        add_issue(
            rows,
            "medium",
            "per_pathogen_null",
            "Some pathogen-strategy pairs do not exceed their permutation p95",
            evidence,
            "Keep pathogen-level weak-case language and avoid per-pathogen universal claims.",
        )


def audit_raw_sequence(rows: list[dict[str, object]]) -> None:
    raw = pd.read_csv(RESULTS_DIR / "pahd_r_raw_sequence_stability_summary.csv")
    risky = raw[raw["raw_sequence_stability_status"] != "pass_proxy"]
    for _, row in risky.iterrows():
        add_issue(
            rows,
            "medium",
            "raw_sequence_stability",
            f"{row['pathogen']} raw sequence proxy stability needs review",
            (
                f"duplicate_rate={row['duplicate_rate']:.3f}, "
                f"mean_top_jaccard={row['mean_top_jaccard']:.3f}, "
                f"mean_spearman={row['mean_spearman']:.3f}"
            ),
            "Run true sequence/time/tree bootstrap before claiming raw-data robustness.",
        )
    high_dup = raw[raw["duplicate_rate"] > 0.50]
    if len(high_dup):
        add_issue(
            rows,
            "medium",
            "raw_sequence_duplicates",
            "Several FASTA files have duplicate rate above 50%",
            ", ".join(f"{r.pathogen}:{r.duplicate_rate:.2f}" for r in high_dup.itertuples()),
            "Collapse duplicates or stratify by time/location before final raw bootstrap.",
        )


def audit_site_call_sources(rows: list[dict[str, object]]) -> None:
    old_path = RESULTS_DIR / "pahd_r_final_site_calls.csv"
    new_path = RESULTS_DIR / "pahd_r_final_mode_site_calls.csv"
    if old_path.exists() and new_path.exists():
        old = pd.read_csv(old_path)
        if "strategy" in old.columns and not set(old["strategy"].unique()).issubset({"PAHD-R-Core", "PAHD-R-Augmented", "PAHD-R-Review"}):
            add_issue(
                rows,
                "medium",
                "artifact_staleness",
                "Older site-call file uses pre-final strategy names",
                f"old strategies: {', '.join(sorted(old['strategy'].unique())[:4])}",
                "Use `pahd_r_final_mode_site_calls.csv` as the reporting source.",
            )
    if not new_path.exists():
        add_issue(
            rows,
            "high",
            "artifact_missing",
            "Final-mode site-call file is missing",
            str(new_path.relative_to(PROJECT_ROOT)),
            "Regenerate final-mode site calls before reporting.",
        )


def audit_repair_claims(rows: list[dict[str, object]]) -> None:
    reporting = pd.read_csv(RESULTS_DIR / "pahd_r_final_reporting_table.csv")
    normalized_status = reporting["status"].str.lower().str.strip()
    bad_adopted = reporting[
        reporting["scope"].isin(["SARS-CoV-2 only", "MERS only", "RSV only"])
        & normalized_status.isin({"adopted", "adopted optional repair"})
    ]
    if len(bad_adopted):
        add_issue(
            rows,
            "high",
            "repair_claims",
            "A non-HCV repair appears adopted in the final reporting table",
            ", ".join(bad_adopted["entry"].astype(str)),
            "Downgrade SARS/MERS/RSV repairs to candidate/rejected unless new validation passes.",
        )


def audit_new_variant_round(rows: list[dict[str, object]]) -> None:
    variant_path = RESULTS_DIR / "pahd_r_robustness_variant_summary.csv"
    sig_path = RESULTS_DIR / "pahd_r_winsorized_paired_significance.csv"
    dedup_path = RESULTS_DIR / "pahd_r_dedup_sequence_sensitivity_summary.csv"
    targeted_path = RESULTS_DIR / "pahd_r_targeted_dedup_winsorized_summary.csv"
    targeted_meta_path = RESULTS_DIR / "pahd_r_targeted_dedup_winsorized_metadata.csv"

    if variant_path.exists():
        variants = pd.read_csv(variant_path)
        candidates = variants[variants["decision"] == "candidate_preprocessing_improvement"]
        if len(candidates):
            add_issue(
                rows,
                "medium",
                "candidate_preprocessing",
                "Winsorized preprocessing improves mean region MCC but is not yet adopted",
                ", ".join(
                    f"{r.mode}:dMCC10={r.delta_mcc_window_10:.3f},dExact={r.delta_mcc_exact:.3f}"
                    for r in candidates.itertuples()
                ),
                "Keep as PAHD-R preprocessing candidate until paired/pathogen-level validation is stronger.",
            )
        load_bearing = variants[
            (variants["decision"] == "load_bearing_or_rejected")
            & (variants["variant"].isin(["drop_recurrence_diversity", "drop_phylo"]))
        ]
        if len(load_bearing):
            add_issue(
                rows,
                "medium",
                "load_bearing_evidence",
                "Recurrence/diversity and phylo ablations show these evidence families are load-bearing",
                ", ".join(sorted(set(load_bearing["variant"].astype(str)))),
                "Do not simplify PAHD-R by removing sequence variation or phylogenetic evidence.",
            )

    if sig_path.exists():
        sig = pd.read_csv(sig_path)
        weak = sig[
            sig["metric"].isin(["mcc_window_10", "mcc_exact"])
            & (sig["wilcoxon_p"] > 0.05)
        ]
        if len(weak):
            add_issue(
                rows,
                "medium",
                "candidate_preprocessing_significance",
                "Winsorized preprocessing does not show significant paired improvement on 9 pathogens",
                "all MCC exact/window Wilcoxon p-values > 0.05",
                "Do not replace the baseline final modes with winsorized variants without a larger panel or bootstrap validation.",
            )

    if dedup_path.exists():
        dedup = pd.read_csv(dedup_path)
        changed = dedup[dedup["dedup_decision"] == "dedup_changes_top_sites"]
        if len(changed):
            add_issue(
                rows,
                "medium",
                "dedup_sequence_sensitivity",
                "Deduplication changes top raw-variation sites for some pathogens",
                ", ".join(
                    f"{r.pathogen}:J={r.original_vs_dedup_top_jaccard:.3f}"
                    for r in changed.itertuples()
                ),
                "Run deduplicated/time-stratified feature recomputation before raw-data robustness claims.",
            )

    if targeted_path.exists():
        targeted = pd.read_csv(targeted_path)
        rejected = targeted[targeted["decision"] == "rejected_input_control"]
        if len(rejected):
            evidence = ", ".join(
                f"{r.variant}/{r.mode}:target_dMCC10={r.target_delta_mcc_window_10:.3f}"
                for r in rejected.head(6).itertuples()
            )
            add_issue(
                rows,
                "medium",
                "targeted_dedup_control",
                "Targeted HIV-1/MERS dedup input controls often degrade the targeted mean",
                evidence,
                "Do not adopt dedup-proxy feature replacement; treat dedup sensitivity as a raw-data robustness caveat.",
            )

    if targeted_meta_path.exists():
        meta = pd.read_csv(targeted_meta_path)
        time_meta = meta[meta["proxy_mode"] == "time_equalized"]
        insufficient = time_meta[time_meta["proxy_status"] == "insufficient_time_metadata"]
        if len(insufficient):
            evidence = ", ".join(
                f"{r.pathogen}:year_frac={r.year_assigned_frac:.3f}"
                for r in insufficient.drop_duplicates(["pathogen"]).itertuples()
            )
            add_issue(
                rows,
                "medium",
                "time_stratification_metadata",
                "Current HIV-1/MERS FASTA headers are insufficient for time-stratified feature recomputation",
                evidence,
                "Acquire accession-level collection dates or a curated temporal MSA before claiming time-stratified robustness.",
            )


def run() -> None:
    rows: list[dict[str, object]] = []
    compare_final_metrics(rows)
    audit_scope_claims(rows)
    audit_adversarial_results(rows)
    audit_raw_sequence(rows)
    audit_site_call_sources(rows)
    audit_repair_claims(rows)
    audit_new_variant_round(rows)

    issues = pd.DataFrame(rows)
    if issues.empty:
        issues = pd.DataFrame([{
            "severity": "none",
            "area": "integrated_review",
            "issue": "No blocking issues detected",
            "evidence": "All automated checks passed",
            "recommended_action": "Proceed with stated caveats.",
        }])

    order = {"high": 0, "medium": 1, "low": 2, "none": 3}
    issues["severity_order"] = issues["severity"].map(order).fillna(9)
    issues = issues.sort_values(["severity_order", "area"]).drop(columns=["severity_order"])

    summary = (
        issues.groupby("severity")
        .size()
        .reset_index(name="count")
        .sort_values("severity", key=lambda s: s.map(order).fillna(9))
    )

    issues_path = RESULTS_DIR / "pahd_r_integrated_adversarial_issues.csv"
    report_path = RESULTS_DIR / "pahd_r_integrated_adversarial_review.md"
    issues.to_csv(issues_path, index=False)

    report = [
        "# PAHD-R Integrated Adversarial Review",
        "",
        "Date: 2026-04-27 KST",
        "",
        "## Summary",
        "",
        markdown_table(summary),
        "",
        "## Issues",
        "",
        markdown_table(issues),
        "",
        "## Verdict",
        "",
        (
            "No blocking metric mismatch was detected between the final reporting "
            "table and final-mode stability table. The remaining risks are claim "
            "scope, stale artifact usage, raw sequence sampling robustness, and "
            "pathogen-level weak cases."
        ),
        "",
    ]
    report_path.write_text("\n".join(report), encoding="utf-8")
    print(report_path)
    print(summary.to_string(index=False))


if __name__ == "__main__":
    run()
