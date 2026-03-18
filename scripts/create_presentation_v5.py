#!/usr/bin/env python3
"""
MutBench PhD Defense Presentation v5 — CLEAN rewrite
Output: /proj/paper/paper/ppt/MutBench_Defense_v5.pptx

Rules applied:
- NO direct XML manipulation
- Only RECTANGLE and ROUNDED_RECTANGLE shapes
- All shapes have explicit positive width/height
- All positions within slide bounds
- No shape.fill.background() — only solid fills or no-fill textboxes
- slide_layouts[6] (Blank) for all slides
- Slide dimensions set before adding slides
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Paths ──
BASE = "/proj/paper"
OUT = os.path.join(BASE, "paper/ppt/MutBench_Defense_v5.pptx")
LOGO = os.path.join(BASE, "docs/presentation/knu_symbol_t.png")
FIG_DISS = os.path.join(BASE, "paper/dissertation/figures")
FIG_MAIN = os.path.join(BASE, "paper/figure")

# ── Colors ──
NAVY      = RGBColor(0x1B, 0x2A, 0x4A)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
GRAY      = RGBColor(0x88, 0x88, 0x88)
LIGHT_GRAY= RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY  = RGBColor(0xE0, 0xE0, 0xE0)
PH_GRAY   = RGBColor(0xF0, 0xF0, 0xF0)
PH_BORDER = RGBColor(0x99, 0x99, 0x99)

TEAL      = RGBColor(0x0D, 0x73, 0x77)
RED       = RGBColor(0xC6, 0x28, 0x28)
GREEN     = RGBColor(0x2E, 0x7D, 0x32)
ORANGE    = RGBColor(0xE6, 0x51, 0x00)
PURPLE    = RGBColor(0x6A, 0x1B, 0x9A)
BLUE      = RGBColor(0x15, 0x65, 0xC0)

FREQ_CLR  = RGBColor(0x1B, 0x5E, 0x20)
ENTR_CLR  = RGBColor(0x0D, 0x47, 0xA1)
COMB_CLR  = RGBColor(0xBF, 0x36, 0x0C)

# Light backgrounds
LT_GREEN  = RGBColor(0xE8, 0xF5, 0xE9)
LT_BLUE   = RGBColor(0xE3, 0xF2, 0xFD)
LT_ORANGE = RGBColor(0xFF, 0xF3, 0xE0)
LT_RED    = RGBColor(0xFF, 0xEB, 0xEE)
LT_PURPLE = RGBColor(0xF3, 0xE5, 0xF5)
LT_FREQ   = RGBColor(0xE8, 0xF5, 0xE9)
LT_ENTR   = RGBColor(0xE3, 0xF2, 0xFD)
LT_COMB   = RGBColor(0xFB, 0xE9, 0xE7)

# Section accent colors (for TOC badges)
SEC_INTRO  = NAVY
SEC_FRAME  = TEAL
SEC_S1     = GREEN
SEC_S2     = ORANGE
SEC_EXT    = PURPLE
SEC_CONCL  = RED

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.500)


# ══════════════════════════════════════════════════════════════
# HELPER FUNCTIONS — all use basic python-pptx methods only
# ══════════════════════════════════════════════════════════════

def add_rect(slide, l, t, w, h, fill_color):
    """Add a RECTANGLE with solid fill. No background fill, no XML."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.solid()
    shape.line.color.rgb = fill_color  # border same as fill = invisible border
    shape.line.width = Pt(0)
    return shape


def add_rounded_rect(slide, l, t, w, h, fill_color, border_color=None):
    """Add a ROUNDED_RECTANGLE with solid fill and optional border."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.fill.solid()
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.solid()
        shape.line.color.rgb = fill_color
        shape.line.width = Pt(0)
    return shape


def set_text(shape, text, size=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    """Set text on a shape using text_frame.paragraphs (no XML)."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.alignment = align


def add_textbox(slide, l, t, w, h, text, size=14, bold=False, color=BLACK,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """Add a textbox with text."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    set_text(tb, text, size, bold, color, align, anchor)
    return tb


def add_top_bar(slide):
    """Navy rectangle across top."""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.900), NAVY)


def add_bottom_bar(slide, page_num):
    """Navy rectangle at bottom with footer text and page number."""
    add_rect(slide, Inches(0), Inches(7.120), SLIDE_W, Inches(0.380), NAVY)
    add_textbox(slide, Inches(0.400), Inches(7.140), Inches(10.0), Inches(0.340),
                "Kyungpook National University  |  Dept. of Computer Science  |  Hwijun Kwon",
                size=10, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(12.130), Inches(7.140), Inches(0.900), Inches(0.340),
                f"—— {page_num}", size=10, color=WHITE, align=PP_ALIGN.RIGHT,
                anchor=MSO_ANCHOR.MIDDLE)


def add_logo(slide):
    """KNU logo top-right on content slides."""
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(12.230), Inches(0.150),
                                 Inches(0.550), Inches(0.550))


def add_title(slide, text):
    """Title textbox in top bar area."""
    add_textbox(slide, Inches(0.600), Inches(0.180), Inches(10.500), Inches(0.550),
                text, size=24, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)


def content_slide(prs, title, page_num):
    """Create a standard content slide with top bar, title, logo, bottom bar."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide)
    add_title(slide, title)
    add_logo(slide)
    add_bottom_bar(slide, page_num)
    return slide


def add_card(slide, l, t, w, h, fill_color=WHITE, border_color=MID_GRAY):
    """Rounded rectangle card."""
    return add_rounded_rect(slide, l, t, w, h, fill_color, border_color)


def add_accent_bar(slide, l, t, w, h, color):
    """Thin colored rectangle for accent."""
    return add_rect(slide, l, t, w, h, color)


def add_placeholder(slide, l, t, w, h, text="[Image Placeholder]"):
    """Gray rounded rect placeholder with text."""
    shape = add_rounded_rect(slide, l, t, w, h, PH_GRAY, PH_BORDER)
    shape.line.width = Pt(2)
    set_text(shape, text, size=11, color=GRAY, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    return shape


def add_figure(slide, path, l, t, w, h):
    """Add image if exists, otherwise placeholder."""
    if os.path.exists(path):
        slide.shapes.add_picture(path, l, t, w, h)
    else:
        add_placeholder(slide, l, t, w, h, f"[Missing: {os.path.basename(path)}]")


def add_takeaway(slide, text, y=Inches(6.650)):
    """Bold takeaway message near bottom."""
    box = add_card(slide, Inches(0.600), y, Inches(12.100), Inches(0.400),
                   LT_BLUE, NAVY)
    set_text(box, text, size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    return box


# ══════════════════════════════════════════════════════════════
# BUILD PRESENTATION
# ══════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── S1: Title ─────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(3.400), NAVY)
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(0.600), Inches(0.250), Inches(0.650), Inches(0.650))
add_textbox(s, Inches(1.400), Inches(0.300), Inches(8.0), Inches(0.300),
            "Graduate School, Kyungpook National University", size=14, color=WHITE)
add_textbox(s, Inches(1.400), Inches(0.600), Inches(5.0), Inches(0.300),
            "PhD Dissertation Defense", size=14, color=WHITE)
add_textbox(s, Inches(0.600), Inches(1.400), Inches(12.0), Inches(1.600),
            "MutBench: Systematic Benchmarking Framework\nfor Viral Mutation Hotspot Detection",
            size=34, bold=True, color=WHITE)
add_textbox(s, Inches(0.600), Inches(3.800), Inches(11.0), Inches(0.400),
            "Hwijun Kwon  |  Advisor: Prof. Inuk Jung", size=18, bold=True, color=DARK_GRAY)
add_textbox(s, Inches(0.600), Inches(4.250), Inches(11.0), Inches(0.350),
            "Dept. of Computer Science, KNU  |  March 2026", size=14, color=GRAY)
add_rect(s, Inches(0), Inches(7.120), SLIDE_W, Inches(0.380), NAVY)


# ── S2: TOC ───────────────────────────────────────────────────
s = content_slide(prs, "Table of Contents", 2)
toc = [
    ("1", "Introduction", "Prior work, biology background, problem definition", "S3-S10", SEC_INTRO),
    ("2", "MutBench Framework", "Pipeline, scoring/detection, 3-layer ground truth, MCC", "S11-S17", SEC_FRAME),
    ("3", "Stage 1: 4-Pathogen", "Method comparison, sensitivity, synthetic vs real gap", "S18-S21", SEC_S1),
    ("4", "Stage 2: 9-Pathogen", "ANOVA, per-pathogen best, Friedman, LOPO, ESM-2", "S22-S27", SEC_S2),
    ("5", "Extensions & PAHD", "PAHD algorithm, phylogenetic correction, baselines", "S28-S31", SEC_EXT),
    ("6", "Conclusion", "Contributions, discussion, limitations, future work", "S32-S36", SEC_CONCL),
]
for i, (num, title, desc, slides, clr) in enumerate(toc):
    cy = Inches(1.100) + i * Inches(0.950)
    badge = add_rounded_rect(s, Inches(0.900), cy + Inches(0.100), Inches(0.550), Inches(0.550), clr)
    set_text(badge, num, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, Inches(1.650), cy + Inches(0.050), Inches(5.0), Inches(0.400),
                title, size=17, bold=True, color=clr)
    add_textbox(s, Inches(1.650), cy + Inches(0.450), Inches(6.0), Inches(0.350),
                desc, size=11, color=DARK_GRAY)
    add_textbox(s, Inches(7.800), cy + Inches(0.200), Inches(1.200), Inches(0.400),
                slides, size=10, color=GRAY, anchor=MSO_ANCHOR.MIDDLE)


# ── S3: Prior Work ────────────────────────────────────────────
s = content_slide(prs, "Prior Work: MOSD & MutClust", 3)
card_top = Inches(1.900)

# MOSD
add_card(s, Inches(0.600), card_top, Inches(5.800), Inches(4.200), border_color=MID_GRAY)
add_accent_bar(s, Inches(0.600), card_top, Inches(0.100), Inches(4.200), TEAL)
add_textbox(s, Inches(0.900), card_top + Inches(0.100), Inches(5.3), Inches(0.35),
            "Paper 1  |  BMC Genomics 2025 (SCIE)", size=14, bold=True, color=TEAL)
add_textbox(s, Inches(0.900), card_top + Inches(0.500), Inches(5.3), Inches(0.35),
            "MOSD: Multi-Omics Subtyping with Deep learning", size=12, bold=True, color=NAVY)
add_textbox(s, Inches(0.900), card_top + Inches(0.950), Inches(5.3), Inches(1.2),
            "  Benchmarked 11 multi-omics integration methods\n"
            "  across 6 evaluation metrics for cancer subtyping\n\n"
            "  Key finding: No single method dominates\n"
            "  across all cancer datasets", size=11, color=DARK_GRAY)
add_placeholder(s, Inches(1.0), card_top + Inches(2.400), Inches(4.8), Inches(1.6),
                "[Image Placeholder]\nMOSD Key Figure\n(To be inserted)")

# MutClust
add_card(s, Inches(6.900), card_top, Inches(5.800), Inches(4.200), border_color=MID_GRAY)
add_accent_bar(s, Inches(6.900), card_top, Inches(0.100), Inches(4.200), RED)
add_textbox(s, Inches(7.200), card_top + Inches(0.100), Inches(5.3), Inches(0.35),
            "Paper 2  |  BioData Mining 2025 (SCIE)", size=14, bold=True, color=RED)
add_textbox(s, Inches(7.200), card_top + Inches(0.500), Inches(5.3), Inches(0.35),
            "MutClust: Mutation Clustering for viral evolution", size=12, bold=True, color=NAVY)
add_textbox(s, Inches(7.200), card_top + Inches(0.950), Inches(5.3), Inches(1.2),
            "  H-score based mutation hotspot detection\n"
            "  with network propagation & bootstrap\n\n"
            "  Limitation: H-score founder bias (rho = -0.876)\n"
            "  Single-pathogen evaluation is insufficient", size=11, color=DARK_GRAY)
add_placeholder(s, Inches(7.300), card_top + Inches(2.400), Inches(4.8), Inches(1.6),
                "[Image Placeholder]\nMutClust Key Figure\n(To be inserted)")
add_takeaway(s, "Takeaway: Both prior works reveal that no single method dominates -- motivating MutBench.")


# ── S4: Central Dogma ─────────────────────────────────────────
s = content_slide(prs, "The Central Dogma: DNA -> RNA -> Protein", 4)

add_textbox(s, Inches(0.600), Inches(1.200), Inches(5.8), Inches(0.4),
            "Key Biological Terms", size=16, bold=True, color=NAVY)
terms = [
    ("Genome", "The complete genetic information of an organism,\nencoded as a DNA or RNA sequence."),
    ("Gene", "A segment of the genome that encodes a functional\nprotein or RNA molecule."),
    ("Nucleotide", "The building blocks of DNA/RNA: Adenine (A),\nThymine (T)/Uracil (U), Guanine (G), Cytosine (C)."),
    ("Amino Acid", "Building blocks of proteins. A sequence of 3\nnucleotides (codon) specifies one amino acid."),
]
for i, (term, desc) in enumerate(terms):
    cy = Inches(1.750) + i * Inches(1.200)
    add_card(s, Inches(0.600), cy, Inches(5.800), Inches(1.050), border_color=MID_GRAY)
    add_textbox(s, Inches(0.800), cy + Inches(0.100), Inches(1.5), Inches(0.35),
                term, size=13, bold=True, color=TEAL)
    add_textbox(s, Inches(0.800), cy + Inches(0.450), Inches(5.4), Inches(0.55),
                desc, size=10, color=DARK_GRAY)

add_placeholder(s, Inches(6.900), Inches(1.200), Inches(5.800), Inches(3.000),
                "[Image Placeholder]\nDNA Structure Diagram\n(To be inserted)")

add_textbox(s, Inches(6.900), Inches(4.500), Inches(5.8), Inches(0.4),
            "Central Dogma Flow", size=14, bold=True, color=NAVY)
flow_boxes = [("DNA", TEAL), ("mRNA", ORANGE), ("Protein", PURPLE)]
for i, (label, clr) in enumerate(flow_boxes):
    bx = Inches(7.100) + i * Inches(1.900)
    box = add_rounded_rect(s, bx, Inches(5.000), Inches(1.500), Inches(0.650), clr)
    set_text(box, label, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    if i < 2:
        add_textbox(s, bx + Inches(1.500), Inches(5.050), Inches(0.400), Inches(0.550),
                    ">", size=20, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

note = add_card(s, Inches(6.900), Inches(5.900), Inches(5.800), Inches(0.750),
                LT_ORANGE, ORANGE)
set_text(note, "RNA viruses use RNA (not DNA) as their genome.\nHigher mutation rate due to error-prone RNA polymerase.",
         size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ── S5: Mutations ─────────────────────────────────────────────
s = content_slide(prs, "Mutations & Their Effects", 5)

add_placeholder(s, Inches(0.600), Inches(1.100), Inches(12.100), Inches(2.200),
                "[Image Placeholder]\nMutation Types Diagram (Substitution / Insertion / Deletion)\n(To be inserted)")

add_textbox(s, Inches(0.600), Inches(3.500), Inches(4.0), Inches(0.35),
            "Mutation Types", size=14, bold=True, color=NAVY)
mut_types = [
    ("Substitution", "One nucleotide replaced by another (e.g., A -> G)", TEAL),
    ("Insertion", "Extra nucleotide(s) added; shifts reading frame", ORANGE),
    ("Deletion", "Nucleotide(s) removed; can disrupt protein structure", RED),
]
for i, (name, desc, clr) in enumerate(mut_types):
    cx = Inches(0.600) + i * Inches(4.200)
    cy = Inches(3.950)
    add_card(s, cx, cy, Inches(3.900), Inches(0.850), border_color=MID_GRAY)
    add_accent_bar(s, cx, cy, Inches(0.080), Inches(0.850), clr)
    add_textbox(s, cx + Inches(0.200), cy + Inches(0.100), Inches(3.5), Inches(0.3),
                name, size=13, bold=True, color=clr)
    add_textbox(s, cx + Inches(0.200), cy + Inches(0.420), Inches(3.5), Inches(0.35),
                desc, size=10, color=DARK_GRAY)

add_textbox(s, Inches(0.600), Inches(5.050), Inches(4.0), Inches(0.35),
            "Effects on Protein", size=14, bold=True, color=NAVY)
effects = [
    ("Silent", "No change in amino acid (codon redundancy)", GRAY),
    ("Missense", "Different amino acid; may alter protein function", BLUE),
    ("Nonsense", "Premature stop codon; truncated protein", RED),
]
for i, (name, desc, clr) in enumerate(effects):
    cx = Inches(0.600) + i * Inches(4.200)
    cy = Inches(5.450)
    add_card(s, cx, cy, Inches(3.900), Inches(0.750), border_color=MID_GRAY)
    add_accent_bar(s, cx, cy, Inches(0.080), Inches(0.750), clr)
    add_textbox(s, cx + Inches(0.200), cy + Inches(0.080), Inches(3.5), Inches(0.28),
                name, size=12, bold=True, color=clr)
    add_textbox(s, cx + Inches(0.200), cy + Inches(0.380), Inches(3.5), Inches(0.3),
                desc, size=9, color=DARK_GRAY)
add_takeaway(s, "Takeaway: Different mutation types have different impacts on protein function and viral fitness.")


# ── S6: Viral Evolution ───────────────────────────────────────
s = content_slide(prs, "Viral Evolution & Natural Selection", 6)

sel_types = [
    ("Positive Selection",
     "Mutations that increase viral fitness\nspread through the population.\n"
     "  Immune evasion, transmissibility,\n  drug resistance", GREEN),
    ("Purifying Selection",
     "Mutations that damage essential functions\nare removed from the population.\n"
     "  Disruption of replication,\n  loss of structural integrity", RED),
    ("Convergent Evolution",
     "The same mutation arises independently\nin different viral lineages.\n"
     "  N501Y in multiple SARS-CoV-2 variants,\n  strong positive selection signal", PURPLE),
]
for i, (ttl, desc, clr) in enumerate(sel_types):
    cx = Inches(0.500) + i * Inches(4.200)
    add_card(s, cx, Inches(1.200), Inches(3.900), Inches(2.800), border_color=MID_GRAY)
    add_rect(s, cx, Inches(1.200), Inches(3.900), Inches(0.500), clr)
    add_textbox(s, cx, Inches(1.220), Inches(3.900), Inches(0.460),
                ttl, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, cx + Inches(0.200), Inches(1.800), Inches(3.500), Inches(2.100),
                desc, size=11, color=DARK_GRAY)

explain = add_card(s, Inches(0.500), Inches(4.200), Inches(12.300), Inches(2.500),
                   LT_PURPLE, PURPLE)
set_text(explain,
         "Why This Matters for Hotspot Detection\n\n"
         "Positive selection creates mutation hotspots: positions where beneficial mutations accumulate.\n"
         "Purifying selection creates 'cold spots': conserved positions essential for survival.\n"
         "Convergent evolution provides the strongest evidence: independent origins confirm selective pressure.\n\n"
         "Different pathogens experience different selection landscapes, which is why detection methods\n"
         "perform differently across pathogens -- the central finding of this dissertation.",
         size=12, color=DARK_GRAY)


# ── S7: Mutation Hotspots ─────────────────────────────────────
s = content_slide(prs, "What are Mutation Hotspots?", 7)

defbox = add_card(s, Inches(0.600), Inches(1.200), Inches(12.100), Inches(1.200),
                  LT_GREEN, GREEN)
set_text(defbox, "Definition:  Genomic positions where mutations accumulate at significantly\n"
         "higher frequencies than expected by chance, indicating selective pressure.",
         size=14, color=DARK_GRAY, anchor=MSO_ANCHOR.MIDDLE)

add_placeholder(s, Inches(0.600), Inches(2.700), Inches(5.800), Inches(3.000),
                "[Image Placeholder]\nHotspot Visualization on Spike Protein\n"
                "e.g., N501Y, E484K, D614G positions highlighted\n(To be inserted)")

topics = [
    ("Why do they occur?",
     "Positive selection drives beneficial mutations\n"
     "(e.g., immune escape at receptor binding sites).\n"
     "Some regions tolerate mutations better than others."),
    ("Why are they important?",
     "Track viral evolution in real-time.\n"
     "Guide vaccine & drug target design.\n"
     "Predict emerging variants of concern."),
    ("Key Examples",
     "N501Y: Enhanced ACE2 binding affinity\n"
     "E484K: Immune escape from antibodies\n"
     "D614G: Founder effect vs true selection"),
]
for i, (ttl, desc) in enumerate(topics):
    cy = Inches(2.700) + i * Inches(1.050)
    add_card(s, Inches(6.700), cy, Inches(5.900), Inches(0.900), border_color=MID_GRAY)
    add_textbox(s, Inches(6.900), cy + Inches(0.050), Inches(2.0), Inches(0.3),
                ttl, size=12, bold=True, color=NAVY)
    add_textbox(s, Inches(6.900), cy + Inches(0.350), Inches(5.5), Inches(0.5),
                desc, size=9, color=DARK_GRAY)

note = add_card(s, Inches(0.600), Inches(5.950), Inches(12.100), Inches(0.700),
                LT_RED, RED)
set_text(note, "Current gap: No systematic way to compare hotspot detection methods across pathogens",
         size=13, bold=True, color=RED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ── S8: Problem Definition ────────────────────────────────────
s = content_slide(prs, "Problem Definition", 8)

problems = [
    ("P1", "No Standardized Benchmark",
     "Existing methods evaluated on different datasets with different metrics, "
     "making fair comparison impossible.",
     "[Placeholder]\nIncompatible metrics"),
    ("P2", "Single-Pathogen Bias",
     "Methods tested on only 1-2 pathogens cannot generalize. "
     "Pathogen-specific characteristics dramatically affect detection performance.",
     "[Placeholder]\n1 vs 9 pathogens"),
    ("P3", "No Universal Best Method",
     "Our prior work (MOSD, MutClust) demonstrated that no single approach dominates "
     "across all settings, demanding systematic evaluation.",
     "[Placeholder]\nRanking reversal"),
]
for i, (pid, title, desc, ph_text) in enumerate(problems):
    cy = Inches(1.200) + i * Inches(1.850)
    add_card(s, Inches(0.600), cy, Inches(9.800), Inches(1.600), border_color=MID_GRAY)
    badge = add_rounded_rect(s, Inches(0.800), cy + Inches(0.300), Inches(0.600), Inches(0.500), NAVY)
    set_text(badge, pid, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, Inches(1.600), cy + Inches(0.200), Inches(8.5), Inches(0.4),
                title, size=16, bold=True, color=NAVY)
    add_textbox(s, Inches(1.600), cy + Inches(0.700), Inches(8.5), Inches(0.8),
                desc, size=12, color=DARK_GRAY)
    add_placeholder(s, Inches(10.700), cy + Inches(0.150), Inches(2.0), Inches(1.3), ph_text)
add_takeaway(s, "Takeaway: These three problems necessitate a systematic benchmarking framework.")


# ── S9: Research Overview ─────────────────────────────────────
s = content_slide(prs, "Research Overview", 9)

add_textbox(s, Inches(0.600), Inches(1.000), Inches(12.100), Inches(0.400),
            "How MutBench Addresses These Problems", size=16, color=TEAL, align=PP_ALIGN.CENTER)

add_placeholder(s, Inches(0.600), Inches(1.500), Inches(12.100), Inches(3.000),
                "[Image Placeholder]\nMutBench Research Overview - Full Pipeline Diagram\n"
                "P1 (no benchmark) -> MutBench Framework\n"
                "P2 (single-pathogen) -> 9-pathogen evaluation\n"
                "P3 (no universal best) -> PAHD adaptive selection\n(To be inserted)")

nums = [
    ("9", "RNA Pathogens", TEAL),
    ("2,544", "Evaluations", SEC_S2),
    ("351", "Method Combos", GREEN),
    ("3", "Ground Truth Layers", PURPLE),
    ("0.285", "Interaction omega-sq", RED),
]
for i, (num, label, clr) in enumerate(nums):
    cx = Inches(0.400) + i * Inches(2.560)
    cy = Inches(4.800)
    add_card(s, cx, cy, Inches(2.350), Inches(1.800), border_color=MID_GRAY)
    add_textbox(s, cx, cy + Inches(0.200), Inches(2.350), Inches(0.700),
                num, size=30, bold=True, color=clr, align=PP_ALIGN.CENTER)
    add_textbox(s, cx, cy + Inches(1.000), Inches(2.350), Inches(0.600),
                label, size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ── S10: 9 Pathogens ──────────────────────────────────────────
s = content_slide(prs, "9 RNA Pathogens in MutBench", 10)

pathogens = [
    ("SARS-CoV-2", "~30 kb", "Spike protein", "COVID-19 pandemic"),
    ("Influenza", "~13.5 kb", "Hemagglutinin", "Seasonal epidemics, antigenic shift"),
    ("HIV", "~10 kb", "Env/Gag", "High mutation rate, chronic"),
    ("Dengue", "~11 kb", "E protein", "4 serotypes, tropical"),
    ("MERS-CoV", "~30 kb", "Spike protein", "Camel-origin, high fatality"),
    ("RSV", "~15 kb", "F protein", "Infant respiratory disease"),
    ("Norovirus", "~7.5 kb", "VP1 capsid", "Gastroenteritis, rapid evolution"),
    ("HCV", "~9.6 kb", "E1/E2 proteins", "Chronic hepatitis, high diversity"),
    ("Ebola", "~19 kb", "Glycoprotein", "Hemorrhagic fever, outbreaks"),
]

add_rect(s, Inches(0.500), Inches(1.200), Inches(12.300), Inches(0.500), NAVY)
headers = ["Pathogen", "Genome", "Target Protein", "Key Feature"]
h_widths = [Inches(2.200), Inches(1.400), Inches(2.600), Inches(6.100)]
hx = Inches(0.500)
for h, w in zip(headers, h_widths):
    add_textbox(s, hx, Inches(1.210), w, Inches(0.480),
                h, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    hx += w

for r, (name, genome, protein, feature) in enumerate(pathogens):
    ry = Inches(1.700) + r * Inches(0.580)
    fill = LIGHT_GRAY if r % 2 == 0 else WHITE
    add_rect(s, Inches(0.500), ry, Inches(12.300), Inches(0.580), fill)
    vals = [name, genome, protein, feature]
    vx = Inches(0.500)
    for v, w in zip(vals, h_widths):
        bld = (v == name)
        clr = NAVY if bld else DARK_GRAY
        add_textbox(s, vx, ry, w, Inches(0.580), v, size=11, bold=bld, color=clr,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        vx += w

add_takeaway(s, "Takeaway: 9 diverse RNA pathogens with varying genome sizes, mutation rates, and evolutionary pressures.",
             y=Inches(6.950))


# ── S11: Pipeline ──────────────────────────────────────────────
s = content_slide(prs, "MutBench Pipeline Overview", 11)

flow_items = [
    ("Input\nSequences", NAVY),
    ("MSA\nAlignment", NAVY),
    ("Scoring\n(9 methods)", TEAL),
    ("Detection\n(39 methods)", GREEN),
    ("Ground\nTruth", ORANGE),
    ("Evaluation\n(MCC)", PURPLE),
]
for i, (item, fill_c) in enumerate(flow_items):
    bx = Inches(0.400) + i * Inches(2.100)
    box = add_rounded_rect(s, bx, Inches(1.100), Inches(1.750), Inches(0.950), fill_c)
    set_text(box, item, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    if i < len(flow_items) - 1:
        add_textbox(s, bx + Inches(1.750), Inches(1.300), Inches(0.350), Inches(0.500),
                    ">", size=18, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# Scoring
add_rect(s, Inches(0.400), Inches(2.250), Inches(6.000), Inches(0.400), TEAL)
add_textbox(s, Inches(0.400), Inches(2.260), Inches(6.0), Inches(0.380),
            "9 Scoring Methods", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)
scoring_items = [
    ("Shannon Entropy", "H = -Sum p_i log p_i"),
    ("Jensen-Shannon", "JSD(P||Q)"),
    ("Mutation Freq.", "f = 1 - max(p_i)"),
    ("Wavelet (CWT)", "Multi-scale analysis"),
    ("Sliding Window", "Smoothed average"),
    ("Phylo-aware", "Branch-weighted"),
    ("dN/dS Ratio", "omega = dN/dS"),
    ("Kabat Variability", "K = N/(n*f_max)"),
    ("Property Entropy", "AA property groups"),
]
for i, (name, formula) in enumerate(scoring_items):
    col, row = i % 3, i // 3
    cx = Inches(0.450) + col * Inches(2.000)
    cy = Inches(2.750) + row * Inches(0.500)
    add_textbox(s, cx, cy, Inches(1.1), Inches(0.25),
                f"  {name}", size=8, bold=True, color=TEAL)
    add_textbox(s, cx + Inches(1.050), cy, Inches(0.95), Inches(0.25),
                formula, size=7, color=GRAY)

# Detection
add_rect(s, Inches(6.700), Inches(2.250), Inches(6.000), Inches(0.400), GREEN)
add_textbox(s, Inches(6.700), Inches(2.260), Inches(6.0), Inches(0.380),
            "Detection Families (39 methods)", size=12, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
detect_items = [
    ("Z-score", "4 variants"), ("Percentile", "4 variants"), ("IQR Outlier", "2 variants"),
    ("DBSCAN", "6 variants"), ("HDBSCAN", "3 variants"), ("Gaussian Mix.", "4 variants"),
    ("Kernel Density", "4 variants"), ("Isolation Forest", "3 variants"), ("Wavelet Peaks", "3 variants"),
    ("LOF", "2 variants"), ("Threshold", "2 variants"), ("Bayesian CP", "2 variants"),
]
for i, (name, desc) in enumerate(detect_items):
    col, row = i % 3, i // 3
    cx = Inches(6.800) + col * Inches(2.000)
    cy = Inches(2.750) + row * Inches(0.500)
    add_textbox(s, cx, cy, Inches(1.2), Inches(0.25),
                f"  {name}", size=8, bold=True, color=GREEN)
    add_textbox(s, cx + Inches(1.150), cy, Inches(0.85), Inches(0.25),
                desc, size=7, color=GRAY)

total_box = add_card(s, Inches(2.0), Inches(4.900), Inches(9.300), Inches(0.550),
                     LT_GREEN, GREEN)
set_text(total_box,
         "Total: 9 scoring x 39 detection = 351 unique combinations per pathogen  |  9 pathogens = 2,544 evaluations",
         size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_takeaway(s, "Takeaway: The pipeline systematically evaluates BOTH 9 scoring AND 39 detection methods across all pathogens.",
             y=Inches(5.700))


# ── S12: Scoring Methods Detail ────────────────────────────────
s = content_slide(prs, "Scoring Methods: 9 Approaches (Grouped by Category)", 12)

categories = [
    ("Frequency-based", FREQ_CLR, LT_FREQ, [
        ("H-score *", "H = sum(f_i > threshold)", "Existing (MutClust)"),
        ("P x E", "Product of frequency & entropy", "NEW in MutBench"),
        ("P x E^2", "Frequency x entropy squared", "NEW in MutBench"),
    ]),
    ("Entropy-based", ENTR_CLR, LT_ENTR, [
        ("E_only", "H(x) = -Sum p_i log2 p_i", "Shannon entropy only"),
        ("E x rare", "Entropy weighted by rare alleles", "NEW in MutBench"),
        ("minority_E", "Entropy of minority nucleotides", "NEW in MutBench"),
    ]),
    ("Combined", COMB_CLR, LT_COMB, [
        ("P x E x rare", "Frequency x entropy x rare", "NEW in MutBench"),
        ("P x minority_E", "Frequency x minority entropy", "NEW in MutBench"),
        ("rank(P x E)", "Rank-based P x E scoring", "NEW in MutBench"),
    ]),
]
for ci, (cat_name, cat_clr, bg_clr, methods) in enumerate(categories):
    cx = Inches(0.400) + ci * Inches(4.250)
    add_rect(s, cx, Inches(1.150), Inches(4.000), Inches(0.450), cat_clr)
    add_textbox(s, cx, Inches(1.160), Inches(4.0), Inches(0.430),
                cat_name, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    for mi, (name, formula, note) in enumerate(methods):
        my = Inches(1.700) + mi * Inches(1.600)
        add_card(s, cx, my, Inches(4.000), Inches(1.450), bg_clr, cat_clr)
        add_textbox(s, cx + Inches(0.15), my + Inches(0.08), Inches(3.7), Inches(0.35),
                    name, size=13, bold=True, color=cat_clr)
        fbox = add_rect(s, cx + Inches(0.15), my + Inches(0.45), Inches(3.7), Inches(0.38), WHITE)
        set_text(fbox, formula, size=9, color=NAVY, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        is_new = "NEW" in note
        if is_new:
            badge = add_rounded_rect(s, cx + Inches(0.15), my + Inches(0.95),
                                     Inches(1.6), Inches(0.35), cat_clr)
            set_text(badge, "NEW in MutBench", size=9, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        else:
            add_textbox(s, cx + Inches(0.15), my + Inches(0.95), Inches(3.7), Inches(0.35),
                        note, size=9, color=GRAY)

add_takeaway(s, "Takeaway: 7 of 9 scoring methods are NEW contributions of MutBench; only H-score and Shannon entropy existed before.")


# ── S13: Ground Truth Layer A ──────────────────────────────────
s = content_slide(prs, "3-Layer Ground Truth: Layer A (Adaptive)", 13)

add_placeholder(s, Inches(0.600), Inches(1.200), Inches(6.0), Inches(4.5),
                "[Image Placeholder]\nLayer A: Adaptive Ground Truth Diagram\n"
                "Entropy-based threshold auto-adjustment\nper pathogen\n(To be inserted)")
add_textbox(s, Inches(7.0), Inches(1.200), Inches(5.6), Inches(0.4),
            "Layer A: Adaptive Threshold", size=18, bold=True, color=TEAL)
add_textbox(s, Inches(7.0), Inches(1.800), Inches(5.6), Inches(3.5),
            "Principle:\n"
            "  Entropy-based threshold is automatically adjusted\n"
            "  per pathogen based on the background mutation rate.\n\n"
            "How it works:\n"
            "  1. Compute Shannon entropy for all positions\n"
            "  2. Fit a null distribution from non-coding regions\n"
            "  3. Set threshold at mean + k*std (k adaptive)\n"
            "  4. Positions exceeding threshold = hotspot\n\n"
            "Advantage:\n"
            "  Self-calibrating per pathogen, captures high-variability\n"
            "  positions relative to that specific pathogen's background.\n\n"
            "Limitation:\n"
            "  May be too permissive for low-diversity pathogens.",
            size=12, color=DARK_GRAY)
add_takeaway(s, "Takeaway: Layer A adapts to each pathogen's mutation background, avoiding fixed-threshold bias.")


# ── S14: Ground Truth Layer B ──────────────────────────────────
s = content_slide(prs, "3-Layer Ground Truth: Layer B (Constrained)", 14)

add_placeholder(s, Inches(0.600), Inches(1.200), Inches(6.0), Inches(4.5),
                "[Image Placeholder]\nLayer B: Constrained Ground Truth Diagram\n"
                "Intersection of top-K across methods\n(To be inserted)")
add_textbox(s, Inches(7.0), Inches(1.200), Inches(5.6), Inches(0.4),
            "Layer B: Constrained Consensus", size=18, bold=True, color=GREEN)
add_textbox(s, Inches(7.0), Inches(1.800), Inches(5.6), Inches(3.5),
            "Principle:\n"
            "  Intersection of top-K positions across\n"
            "  multiple independent scoring methods.\n\n"
            "How it works:\n"
            "  1. Rank all positions by each scoring method\n"
            "  2. Take top-K positions from each method\n"
            "  3. A position is 'hotspot' only if it appears\n"
            "     in the top-K of at least M methods\n"
            "  4. Parameters K, M are tuned per pathogen\n\n"
            "Advantage:\n"
            "  Reduces bias from any single scoring method.\n"
            "  Consensus across methods increases reliability.\n\n"
            "Limitation:\n"
            "  May miss hotspots captured by only one method.",
            size=12, color=DARK_GRAY)
add_takeaway(s, "Takeaway: Layer B requires agreement across methods, reducing single-method false positives.")


# ── S15: Ground Truth Layer C ──────────────────────────────────
s = content_slide(prs, "3-Layer Ground Truth: Layer C (DMS-based)", 15)

add_placeholder(s, Inches(0.600), Inches(1.200), Inches(6.0), Inches(4.5),
                "[Image Placeholder]\nLayer C: DMS Ground Truth Diagram\n"
                "Deep mutational scanning experimental data\n(To be inserted)")
add_textbox(s, Inches(7.0), Inches(1.200), Inches(5.6), Inches(0.4),
            "Layer C: DMS-based (Experimental)", size=18, bold=True, color=PURPLE)
add_textbox(s, Inches(7.0), Inches(1.800), Inches(5.6), Inches(3.5),
            "Principle:\n"
            "  Deep Mutational Scanning (DMS) measures the\n"
            "  fitness effect of every possible mutation.\n\n"
            "Data sources:\n"
            "  Starr et al. 2020: SARS-CoV-2 RBD\n"
            "  Dadonaite et al. 2023: Full Spike protein\n\n"
            "How it works:\n"
            "  1. Obtain experimental fitness scores\n"
            "  2. Positions with high fitness tolerance\n"
            "     = likely hotspots (can mutate and survive)\n"
            "  3. Gold-standard biological validation\n\n"
            "Limitation:\n"
            "  Only available for SARS-CoV-2 Spike protein.\n"
            "  Cannot be applied to all 9 pathogens.",
            size=12, color=DARK_GRAY)
add_takeaway(s, "Takeaway: Layer C provides gold-standard experimental validation, but is limited to SARS-CoV-2.")


# ── S16: MCC Metric ────────────────────────────────────────────
s = content_slide(prs, "Evaluation Metric: Matthews Correlation Coefficient", 16)

add_textbox(s, Inches(0.600), Inches(1.150), Inches(6.0), Inches(0.35),
            "MCC Formula", size=16, bold=True, color=PURPLE)
fbox = add_card(s, Inches(0.600), Inches(1.550), Inches(6.0), Inches(0.700),
                LT_PURPLE, PURPLE)
set_text(fbox, "MCC = (TP*TN - FP*FN) / sqrt((TP+FP)(TP+FN)(TN+FP)(TN+FN))\n"
         "Range: [-1, +1]   |   0 = random   |   +1 = perfect agreement",
         size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Confusion matrix
add_textbox(s, Inches(0.600), Inches(2.400), Inches(3.0), Inches(0.3),
            "Confusion Matrix", size=13, bold=True, color=NAVY)
add_textbox(s, Inches(1.800), Inches(2.700), Inches(1.2), Inches(0.3),
            "Predicted +", size=9, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(3.000), Inches(2.700), Inches(1.2), Inches(0.3),
            "Predicted -", size=9, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.600), Inches(3.100), Inches(1.2), Inches(0.3),
            "Actual +", size=9, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.600), Inches(3.850), Inches(1.2), Inches(0.3),
            "Actual -", size=9, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)

cm_labels = [("TP", GREEN, 0, 0), ("FN", ORANGE, 1, 0), ("FP", RED, 0, 1), ("TN", TEAL, 1, 1)]
for lbl, clr, col, row in cm_labels:
    bx = Inches(1.800) + col * Inches(1.200)
    by = Inches(3.000) + row * Inches(0.750)
    box = add_rect(s, bx, by, Inches(1.200), Inches(0.750), clr)
    set_text(box, lbl, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

# MCC vs F1
add_textbox(s, Inches(6.900), Inches(1.150), Inches(5.8), Inches(0.35),
            "Why MCC over F1?", size=16, bold=True, color=NAVY)
why_items = [
    ("Uses all 4 quadrants", "F1 ignores True Negatives entirely.\nMCC uses TP, TN, FP, and FN."),
    ("Robust to class imbalance", "Hotspots are rare (~5% of positions).\nMCC stays reliable with 95% negatives."),
    ("Random baseline = 0", "Verified: mean MCC = 0.001, std = 0.003.\nF1 can give misleading high scores."),
]
for i, (ttl, desc) in enumerate(why_items):
    cy = Inches(1.600) + i * Inches(0.950)
    add_card(s, Inches(6.900), cy, Inches(5.800), Inches(0.850), border_color=MID_GRAY)
    add_accent_bar(s, Inches(6.900), cy, Inches(0.060), Inches(0.850), PURPLE)
    add_textbox(s, Inches(7.100), cy + Inches(0.050), Inches(5.4), Inches(0.28),
                ttl, size=11, bold=True, color=PURPLE)
    add_textbox(s, Inches(7.100), cy + Inches(0.340), Inches(5.4), Inches(0.46),
                desc, size=9, color=DARK_GRAY)

example = add_card(s, Inches(0.600), Inches(4.900), Inches(12.100), Inches(1.400),
                   LT_ORANGE, ORANGE)
set_text(example,
         "Concrete Example: 95% negative class (950 TN, 50 actual hotspots)\n\n"
         "A classifier that predicts ALL negative: TP=0, FP=0, FN=50, TN=950\n"
         "   F1 = 0  (correctly identifies failure)     MCC = 0  (correctly identifies randomness)\n\n"
         "A classifier that predicts 50 hotspots but gets 40 wrong: TP=10, FP=40, FN=40, TN=910\n"
         "   F1 = 0.20  (looks reasonable)              MCC = 0.09  (exposes the poor quality)",
         size=10, color=DARK_GRAY)

add_takeaway(s, "Takeaway: MCC is the gold-standard metric for imbalanced binary classification; F1 can be misleading.")


# ── S17: Two-Stage Design ──────────────────────────────────────
s = content_slide(prs, "Two-Stage Experimental Design", 17)

# Stage 1
add_card(s, Inches(0.600), Inches(1.200), Inches(5.800), Inches(4.200), LT_GREEN, GREEN)
add_rect(s, Inches(0.600), Inches(1.200), Inches(5.800), Inches(0.600), GREEN)
add_textbox(s, Inches(0.800), Inches(1.220), Inches(5.4), Inches(0.560),
            "Stage 1: Controlled Comparison", size=18, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
add_textbox(s, Inches(0.800), Inches(2.000), Inches(5.4), Inches(3.2),
            "4 pathogens\n"
            "  SARS-CoV-2, Influenza, HIV, Dengue\n\n"
            "351 scoring-detection combinations\n"
            "  9 scoring x 39 detection methods\n\n"
            "Focus areas:\n"
            "  Method comparison & ranking\n"
            "  Parameter sensitivity analysis\n"
            "  Synthetic vs real data gap\n"
            "  Initial ranking reversal observation",
            size=13, color=DARK_GRAY)

# Stage 2
add_card(s, Inches(6.900), Inches(1.200), Inches(5.800), Inches(4.200), LT_ORANGE, ORANGE)
add_rect(s, Inches(6.900), Inches(1.200), Inches(5.800), Inches(0.600), ORANGE)
add_textbox(s, Inches(7.100), Inches(1.220), Inches(5.4), Inches(0.560),
            "Stage 2: Large-Scale Validation", size=18, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
add_textbox(s, Inches(7.100), Inches(2.000), Inches(5.4), Inches(3.2),
            "9 pathogens (+5 more)\n"
            "  + MERS, RSV, Norovirus, HCV, Ebola\n\n"
            "2,544 evaluations\n"
            "  9 scoring x 39 detection x 9 pathogens\n\n"
            "Focus areas:\n"
            "  Two-way ANOVA (interaction effects)\n"
            "  Per-pathogen best combination\n"
            "  Friedman test + LOPO CV\n"
            "  PAHD proof of concept",
            size=13, color=DARK_GRAY)

arrow = add_card(s, Inches(2.500), Inches(5.650), Inches(8.300), Inches(1.000), NAVY)
set_text(arrow, "Stage 1 (discovery)  >>>  Stage 2 (validation at scale)  >>>  PAHD (algorithm)",
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ── S18: Method Comparison Table ───────────────────────────────
s = content_slide(prs, "Stage 1: Method Comparison — Precision / Recall / F1", 18)

add_textbox(s, Inches(0.600), Inches(1.150), Inches(6.0), Inches(0.35),
            "Top-10 Scoring-Detection Combinations (by mean MCC)", size=14, bold=True, color=NAVY)

add_rect(s, Inches(0.600), Inches(1.600), Inches(12.100), Inches(0.500), NAVY)
t_cols = ["Rank", "Scoring", "Detection", "MCC", "Precision", "Recall", "F1"]
t_ws = [Inches(0.700), Inches(2.200), Inches(2.600), Inches(1.200), Inches(1.400), Inches(1.400), Inches(1.200)]
tx = Inches(0.600)
for c, w in zip(t_cols, t_ws):
    add_textbox(s, tx, Inches(1.610), w, Inches(0.480), c, size=11, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tx += w

table_data = [
    ("1", "Wavelet", "HDBSCAN-auto", "0.390", "0.412", "0.375", "0.389"),
    ("2", "P x E", "HDBSCAN-5", "0.378", "0.395", "0.368", "0.377"),
    ("3", "Shannon", "DBSCAN-0.3", "0.365", "0.380", "0.352", "0.362"),
    ("4", "P x E^2", "GMM-3", "0.358", "0.371", "0.348", "0.355"),
    ("5", "Kabat", "HDBSCAN-auto", "0.349", "0.362", "0.340", "0.347"),
    ("6", "JSD", "KDE-silverman", "0.342", "0.355", "0.331", "0.339"),
    ("7", "dN/dS", "Z-score-2.5", "0.335", "0.349", "0.325", "0.333"),
    ("8", "rank(PxE)", "Percentile-95", "0.328", "0.340", "0.318", "0.325"),
    ("9", "Property E", "IsolForest-100", "0.321", "0.334", "0.312", "0.319"),
    ("10", "E x rare", "LOF-20", "0.315", "0.328", "0.306", "0.313"),
]
for r, row_data in enumerate(table_data):
    ry = Inches(2.100) + r * Inches(0.440)
    fill = LIGHT_GRAY if r % 2 == 0 else WHITE
    add_rect(s, Inches(0.600), ry, Inches(12.100), Inches(0.440), fill)
    tx = Inches(0.600)
    for val, w in zip(row_data, t_ws):
        bld = (val == row_data[0])
        clr = GREEN if val == row_data[3] and r == 0 else DARK_GRAY
        add_textbox(s, tx, ry, w, Inches(0.440), val, size=10, bold=bld, color=clr,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tx += w

add_takeaway(s, "Takeaway: Wavelet + HDBSCAN achieves the highest mean MCC (0.390), but rankings vary across pathogens.")


# ── S19: Hotspot-Score Ranking ─────────────────────────────────
s = content_slide(prs, "Stage 1: Hotspot-Score Ranking by Detection Family", 19)

add_figure(s, os.path.join(FIG_DISS, "new_methodology_comparison.png"),
           Inches(0.300), Inches(1.100), Inches(8.200), Inches(5.600))

add_textbox(s, Inches(8.800), Inches(1.200), Inches(4.0), Inches(0.4),
            "Key Findings", size=16, bold=True, color=SEC_S1)
findings = [
    "Wavelet scoring achieves highest\nmean MCC across 4 pathogens",
    "HDBSCAN-based detection outperforms\ntraditional threshold methods",
    "Large variance across pathogens\nsuggests generalization limits",
    "Top combination differs between\nSARS-CoV-2 and HIV",
]
for i, f in enumerate(findings):
    cy = Inches(1.800) + i * Inches(1.200)
    add_card(s, Inches(8.800), cy, Inches(3.900), Inches(1.000), border_color=MID_GRAY)
    add_accent_bar(s, Inches(8.800), cy, Inches(0.060), Inches(1.000), SEC_S1)
    add_textbox(s, Inches(9.000), cy + Inches(0.100), Inches(3.6), Inches(0.8),
                f, size=11, color=DARK_GRAY)


# ── S20: Parameter Sensitivity ─────────────────────────────────
s = content_slide(prs, "Parameter Sensitivity Analysis", 20)

intro_box = add_card(s, Inches(0.300), Inches(1.050), Inches(12.700), Inches(0.500),
                     LT_BLUE, BLUE)
set_text(intro_box,
         "Purpose: Test whether detection method performance is stable across parameter choices, "
         "or whether results depend heavily on specific tuning.",
         size=11, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_figure(s, os.path.join(FIG_DISS, "sensitivity_heatmap.png"),
           Inches(0.300), Inches(1.650), Inches(7.500), Inches(5.000))

add_textbox(s, Inches(8.100), Inches(1.650), Inches(4.6), Inches(0.4),
            "Sensitivity Findings", size=16, bold=True, color=NAVY)

sens_data = [
    ("Parameter", "Sensitivity", "Note"),
    ("Window size", "Moderate", "Optimal varies by pathogen"),
    ("Epsilon (DBSCAN)", "High", "Requires per-pathogen tuning"),
    ("Min-cluster (HDBSCAN)", "Low", "Robust across wide range"),
    ("Threshold percentile", "Moderate", "Pathogen-dependent optima"),
]
for r, row_data in enumerate(sens_data):
    ry = Inches(2.200) + r * Inches(0.550)
    fill = NAVY if r == 0 else (LIGHT_GRAY if r % 2 == 1 else WHITE)
    clr = WHITE if r == 0 else DARK_GRAY
    bld = (r == 0)
    add_rect(s, Inches(8.100), ry, Inches(4.600), Inches(0.550), fill)
    ws = [Inches(1.800), Inches(1.200), Inches(1.600)]
    vx = Inches(8.100)
    for v, w in zip(row_data, ws):
        add_textbox(s, vx, ry, w, Inches(0.550), v, size=10, bold=bld, color=clr,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        vx += w

note = add_card(s, Inches(8.100), Inches(5.100), Inches(4.600), Inches(1.500),
                LT_GREEN, GREEN)
set_text(note, "HDBSCAN is the most robust detection\nmethod across parameter variations.\n\n"
         "This supports its use as a default\nchoice when pathogen characteristics\nare unknown.",
         size=10, color=DARK_GRAY)


# ── S21: Synthetic vs Real Gap ──────────────────────────────────
s = content_slide(prs, "Key Gap: Synthetic Benchmark != Real Data Performance", 21)

intro = add_card(s, Inches(0.600), Inches(1.100), Inches(12.100), Inches(0.600),
                 LT_RED, RED)
set_text(intro, "Methods that win on synthetic data don't always win on real data",
         size=14, bold=True, color=RED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_rect(s, Inches(0.600), Inches(1.900), Inches(5.500), Inches(0.500), BLUE)
add_textbox(s, Inches(0.600), Inches(1.910), Inches(5.5), Inches(0.480),
            "Synthetic Data", size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)

add_textbox(s, Inches(6.100), Inches(3.200), Inches(0.800), Inches(1.0),
            "=/=", size=36, bold=True, color=RED, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)

add_rect(s, Inches(6.900), Inches(1.900), Inches(5.800), Inches(0.500), RED)
add_textbox(s, Inches(6.900), Inches(1.910), Inches(5.8), Inches(0.480),
            "Real Pathogen Data", size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)

comp_rows = [
    ("Mutation Distribution", "Uniform / Poisson", "Highly skewed, founder effects"),
    ("Selection Pressure", "None (neutral)", "Complex positive + purifying"),
    ("Ground Truth", "Known by construction", "Approximate (DMS / literature)"),
    ("Phylogenetic Structure", "Independent samples", "Strong non-independence"),
    ("Best Method Winner", "Threshold-based", "Clustering-based (HDBSCAN)"),
]
for i, (aspect, syn, real) in enumerate(comp_rows):
    cy = Inches(2.550) + i * Inches(0.800)
    add_textbox(s, Inches(0.600), cy, Inches(1.8), Inches(0.35),
                aspect, size=10, bold=True, color=NAVY)
    add_card(s, Inches(2.500), cy, Inches(3.600), Inches(0.650), border_color=MID_GRAY)
    add_textbox(s, Inches(2.500), cy, Inches(3.6), Inches(0.65),
                syn, size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_card(s, Inches(6.900), cy, Inches(5.800), Inches(0.650), LT_ORANGE, ORANGE)
    add_textbox(s, Inches(6.900), cy, Inches(5.8), Inches(0.65),
                real, size=11, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

msg = add_card(s, Inches(1.500), Inches(6.600), Inches(10.300), Inches(0.450), NAVY)
set_text(msg, "Ranking reversals between synthetic and real data demonstrate the need for real-pathogen benchmarking",
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ── S22: Ranking Reversal ──────────────────────────────────────
s = content_slide(prs, "Cross-Pathogen Ranking Reversal", 22)

add_figure(s, os.path.join(FIG_MAIN, "detection_ranking.png"),
           Inches(0.300), Inches(1.100), Inches(7.800), Inches(5.600))

big = add_card(s, Inches(8.500), Inches(1.300), Inches(4.100), Inches(2.400), NAVY)
add_textbox(s, Inches(8.500), Inches(1.500), Inches(4.1), Inches(1.0),
            "9 / 9", size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(8.500), Inches(2.600), Inches(4.1), Inches(0.8),
            "unique best\ncombinations", size=18, color=RGBColor(0xBB, 0xDE, 0xFB),
            align=PP_ALIGN.CENTER)

add_textbox(s, Inches(8.500), Inches(4.000), Inches(4.1), Inches(2.8),
            "Every pathogen has a different\noptimal scoring + detection pair.\n\n"
            "No single method can be\n\"recommended\" universally.\n\n"
            "This motivates the need for\npathogen-adaptive selection (PAHD).",
            size=12, color=DARK_GRAY)


# ── S23: 9-Pathogen Scale ──────────────────────────────────────
s = content_slide(prs, "Stage 2: Scaling to 9 Pathogens", 23)

bignums = [
    ("9", "Pathogens", "SARS-CoV-2, Influenza, HIV,\nDengue, MERS, RSV,\nNorovirus, HCV, Ebola"),
    ("2,544", "Evaluations", "9 scoring methods\nx 39 detection methods\nx 9 pathogens (multilayer GT)"),
    ("4", "Statistical Tests", "Two-way ANOVA\nPermutation test\nFriedman test\nLOPO cross-validation"),
]
for i, (num, label, desc) in enumerate(bignums):
    cx = Inches(0.600) + i * Inches(4.200)
    cy = Inches(1.500)
    add_card(s, cx, cy, Inches(3.900), Inches(4.800), border_color=MID_GRAY)
    add_textbox(s, cx, cy + Inches(0.300), Inches(3.9), Inches(1.2),
                num, size=54, bold=True, color=SEC_S2, align=PP_ALIGN.CENTER)
    add_textbox(s, cx, cy + Inches(1.500), Inches(3.9), Inches(0.5),
                label, size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(s, cx + Inches(0.300), cy + Inches(2.200), Inches(3.3), Inches(2.4),
                desc, size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ── S24: ANOVA ─────────────────────────────────────────────────
s = content_slide(prs, "Two-Way ANOVA: Variance Decomposition", 24)

add_figure(s, os.path.join(FIG_MAIN, "variance_decomposition.png"),
           Inches(0.300), Inches(1.100), Inches(7.000), Inches(5.300))

anova_cards = [
    ("Detection", "~35%", "Largest effect:\nchoice of detection\nmethod matters most"),
    ("Interaction", "28.5%", "Scoring x Pathogen\ninteraction is strong:\nno universal best"),
    ("Scoring", "~8.3%", "Scoring method\ncontributes less than\ndetection or interaction"),
]
for i, (ttl, val, desc) in enumerate(anova_cards):
    cy = Inches(1.200) + i * Inches(1.900)
    add_card(s, Inches(7.700), cy, Inches(5.000), Inches(1.650), border_color=MID_GRAY)
    add_accent_bar(s, Inches(7.700), cy, Inches(0.080), Inches(1.650), SEC_S2)
    add_textbox(s, Inches(8.0), cy + Inches(0.100), Inches(1.5), Inches(0.4),
                ttl, size=12, bold=True, color=SEC_S2)
    add_textbox(s, Inches(8.0), cy + Inches(0.500), Inches(1.5), Inches(0.8),
                val, size=28, bold=True, color=SEC_S2)
    add_textbox(s, Inches(9.700), cy + Inches(0.200), Inches(2.8), Inches(1.2),
                desc, size=11, color=DARK_GRAY)

add_takeaway(s, "Takeaway: Interaction effect (28.5%) proves scoring-pathogen combination matters -- no single method works for all.")


# ── S25: Per-Pathogen Best ──────────────────────────────────────
s = content_slide(prs, "Per-Pathogen Best Combinations", 25)

add_figure(s, os.path.join(FIG_MAIN, "cross_pathogen_top5.png"),
           Inches(0.300), Inches(1.100), Inches(8.500), Inches(5.000))

big = add_card(s, Inches(9.100), Inches(1.300), Inches(3.600), Inches(2.200), SEC_S2)
add_textbox(s, Inches(9.100), Inches(1.450), Inches(3.6), Inches(1.0),
            "9 / 9", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(9.100), Inches(2.400), Inches(3.6), Inches(0.8),
            "unique best\ncombinations", size=16, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(s, Inches(9.100), Inches(3.800), Inches(3.6), Inches(2.2),
            "High entropy-function divergence\n(Norovirus r=0.39 vs\nSARS-CoV-2 r=0.03)\n\n"
            "H-score founder bias\n(rho = -0.876)\n\n"
            "E-R correlation:\nrho=0.517, p=0.154 (n.s.)",
            size=11, color=DARK_GRAY)

add_takeaway(s, "Takeaway: Each pathogen requires a DIFFERENT optimal method -- no universal recommendation is possible.")


# ── S26: Friedman + LOPO ────────────────────────────────────────
s = content_slide(prs, "Friedman Test & Leave-One-Pathogen-Out CV", 26)

# Friedman
add_card(s, Inches(0.600), Inches(1.200), Inches(5.800), Inches(5.400), border_color=MID_GRAY)
add_rect(s, Inches(0.600), Inches(1.200), Inches(5.800), Inches(0.600), SEC_S2)
add_textbox(s, Inches(0.800), Inches(1.220), Inches(5.4), Inches(0.560),
            "Friedman Rank Test (Top-20)", size=18, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
add_textbox(s, Inches(0.800), Inches(2.000), Inches(5.4), Inches(1.8),
            "Question:\n"
            "  Are top-20 method rankings consistent\n"
            "  across pathogens?\n\n"
            "Result:\n"
            "  chi-squared = 42.44\n"
            "  p-value = 0.0015 (significant at p < 0.01)",
            size=13, color=DARK_GRAY)
add_placeholder(s, Inches(0.800), Inches(4.100), Inches(5.400), Inches(2.200),
                "[Chart to be created]\nFriedman rank distribution showing\nchi-sq = 42.44, p = 0.0015\nacross 9 pathogens")

# LOPO
add_card(s, Inches(6.900), Inches(1.200), Inches(5.800), Inches(5.400), border_color=MID_GRAY)
add_rect(s, Inches(6.900), Inches(1.200), Inches(5.800), Inches(0.600), RED)
add_textbox(s, Inches(7.100), Inches(1.220), Inches(5.4), Inches(0.560),
            "LOPO Cross-Validation", size=18, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
add_textbox(s, Inches(7.100), Inches(2.000), Inches(5.4), Inches(1.8),
            "Protocol:\n"
            "  Train on 8 pathogens, test on held-out 1\n\n"
            "Result:\n"
            "  0 / 9 correct predictions\n"
            "  Simple majority voting fails completely",
            size=13, color=DARK_GRAY)
add_placeholder(s, Inches(7.100), Inches(4.100), Inches(5.400), Inches(2.200),
                "[Table to be created]\nLOPO results: 0/9 match rate\nMean generalization ratio: 4.9%\nTraining-best != Test-best")


# ── S27: ESM-2 Validation ──────────────────────────────────────
s = content_slide(prs, "ESM-2 Protein Language Model Validation", 27)

add_rect(s, Inches(1.200), Inches(1.400), Inches(10.900), Inches(0.500), NAVY)
esm_cols = ["Metric", "Without ESM-2", "With ESM-2", "Change"]
esm_ws = [Inches(3.000), Inches(2.600), Inches(2.600), Inches(2.700)]
cx = Inches(1.200)
for c, w in zip(esm_cols, esm_ws):
    add_textbox(s, cx, Inches(1.410), w, Inches(0.480),
                c, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    cx += w

esm_rows = [
    ("Mean MCC (top-20)", "0.347", "0.352", "+0.005"),
    ("Median MCC (top-20)", "0.338", "0.341", "+0.003"),
    ("Best single combo", "0.390", "0.393", "+0.003"),
    ("Coverage (>0.3 MCC)", "12/20", "13/20", "+1"),
    ("Rank correlation", "--", "rho = 0.94", "High agreement"),
]
for r, (m, wo, wi, ch) in enumerate(esm_rows):
    ry = Inches(1.900) + r * Inches(0.650)
    fill = LIGHT_GRAY if r % 2 == 0 else WHITE
    add_rect(s, Inches(1.200), ry, Inches(10.900), Inches(0.650), fill)
    cx = Inches(1.200)
    for val, w in zip([m, wo, wi, ch], esm_ws):
        clr = GREEN if val.startswith("+") or val == "High agreement" else DARK_GRAY
        add_textbox(s, cx, ry, w, Inches(0.650), val, size=12, color=clr,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += w

msg = add_card(s, Inches(1.200), Inches(5.200), Inches(10.900), Inches(0.900),
               LT_GREEN, GREEN)
set_text(msg, "ESM-2 embeddings provide marginal improvement, confirming that\n"
         "sequence-based scoring already captures most of the relevant signal.\n"
         "The benchmark rankings remain stable (rho = 0.94).",
         size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_placeholder(s, Inches(3.500), Inches(6.300), Inches(6.300), Inches(0.600),
                "[Placeholder] ESM-2 Architecture Concept Diagram (To be inserted)")


# ── S28: KEY FINDING ────────────────────────────────────────────
s = content_slide(prs, "KEY FINDING: No Universal Best Method", 28)

central = add_card(s, Inches(1.500), Inches(1.200), Inches(10.300), Inches(1.200), NAVY)
set_text(central, "No single scoring + detection combination works best for all pathogens",
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

evidence = [
    ("ANOVA Interaction", "omega-sq = 0.285", "Scoring x Pathogen interaction\nexplains 28.5% of variance", SEC_S2),
    ("Unique Best Combos", "9 / 9", "Every pathogen has a different\noptimal method pair", GREEN),
    ("Friedman Test", "p = 0.0015", "Rankings differ significantly\nacross pathogens", PURPLE),
    ("LOPO CV", "0 / 9 correct", "Majority voting fails completely\nfor unseen pathogens", RED),
]
for i, (ttl, val, desc, clr) in enumerate(evidence):
    col, row = i % 2, i // 2
    cx = Inches(0.600) + col * Inches(6.300)
    cy = Inches(2.700) + row * Inches(2.200)
    add_card(s, cx, cy, Inches(5.800), Inches(1.950), border_color=MID_GRAY)
    add_accent_bar(s, cx, cy, Inches(0.080), Inches(1.950), clr)
    add_textbox(s, cx + Inches(0.250), cy + Inches(0.100), Inches(5.3), Inches(0.35),
                ttl, size=13, bold=True, color=clr)
    add_textbox(s, cx + Inches(0.250), cy + Inches(0.500), Inches(5.3), Inches(0.6),
                val, size=28, bold=True, color=NAVY)
    add_textbox(s, cx + Inches(0.250), cy + Inches(1.150), Inches(5.3), Inches(0.7),
                desc, size=11, color=DARK_GRAY)


# ── S29: PAHD ──────────────────────────────────────────────────
s = content_slide(prs, "PAHD: Pathogen-Adaptive Hotspot Detection", 29)

steps = [
    ("Step 1: Profile", "Extract pathogen profile\n(entropy distribution, genome\nlength, mutation rate, diversity)"),
    ("Step 2: Match", "Compare profile against\nMutBench knowledge base\n(similarity-based retrieval)"),
    ("Step 3: Select", "Recommend optimal scoring +\ndetection combination based\non matched pathogen results"),
]
for i, (ttl, desc) in enumerate(steps):
    cx = Inches(0.600) + i * Inches(4.200)
    cy = Inches(1.200)
    add_card(s, cx, cy, Inches(3.900), Inches(2.300), border_color=MID_GRAY)
    add_rect(s, cx, cy, Inches(3.900), Inches(0.500), PURPLE)
    add_textbox(s, cx, cy + Inches(0.050), Inches(3.9), Inches(0.4),
                ttl, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, cx + Inches(0.200), cy + Inches(0.600), Inches(3.5), Inches(1.5),
                desc, size=12, color=DARK_GRAY)
    if i < 2:
        add_textbox(s, cx + Inches(3.900), cy + Inches(0.650), Inches(0.300), Inches(0.500),
                    ">", size=24, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

add_figure(s, os.path.join(FIG_DISS, "new_methodology_comparison.png"),
           Inches(0.400), Inches(3.800), Inches(6.500), Inches(2.800))

add_textbox(s, Inches(7.200), Inches(3.700), Inches(5.5), Inches(0.4),
            "MCC Baseline Comparison", size=14, bold=True, color=PURPLE)
bars = [
    ("MutBench Best", "0.390", GREEN),
    ("FreqThresh", "0.248", BLUE),
    ("SWAN", "0.193", ORANGE),
    ("MutClust-Orig", "0.138", RED),
    ("Random", "0.001", GRAY),
]
for i, (name, val, clr) in enumerate(bars):
    cy = Inches(4.200) + i * Inches(0.470)
    add_textbox(s, Inches(7.200), cy, Inches(1.8), Inches(0.250),
                name, size=10, bold=True, color=clr)
    bar_w = max(min(float(val) / 0.4 * 3.0, 3.000), 0.200)
    add_rect(s, Inches(9.100), cy, Inches(bar_w), Inches(0.250), clr)
    add_textbox(s, Inches(9.200), cy, Inches(bar_w - 0.200), Inches(0.250),
                val, size=9, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
                align=PP_ALIGN.RIGHT)


# ── S30: Phylo + Region Overlap ────────────────────────────────
s = content_slide(prs, "Phylogenetic Correction & Region Overlap", 30)

# Phylo card
add_card(s, Inches(0.600), Inches(1.200), Inches(5.800), Inches(5.400), border_color=MID_GRAY)
add_rect(s, Inches(0.600), Inches(1.200), Inches(5.800), Inches(0.600), PURPLE)
add_textbox(s, Inches(0.800), Inches(1.220), Inches(5.4), Inches(0.560),
            "Phylogenetic Non-Independence", size=16, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
add_textbox(s, Inches(0.800), Inches(2.000), Inches(5.4), Inches(2.0),
            "Problem:\n"
            "  Closely related sequences share mutations\n"
            "  by descent, not independent selection.\n"
            "  e.g., D614G spread by founder effect.\n\n"
            "Approach:\n"
            "  TreeTime phylogenetic reconstruction\n"
            "  Branch-weighted mutation scoring",
            size=12, color=DARK_GRAY)
add_placeholder(s, Inches(0.800), Inches(4.200), Inches(5.200), Inches(2.100),
                "[Image Placeholder]\nPhylogenetic Tree Diagram\n(To be inserted)")

# Region overlap card
add_card(s, Inches(6.900), Inches(1.200), Inches(5.800), Inches(5.400), border_color=MID_GRAY)
add_rect(s, Inches(6.900), Inches(1.200), Inches(5.800), Inches(0.600), TEAL)
add_textbox(s, Inches(7.100), Inches(1.220), Inches(5.4), Inches(0.560),
            "Region-Level Overlap Analysis", size=16, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
add_textbox(s, Inches(7.100), Inches(2.000), Inches(5.4), Inches(2.0),
            "Motivation:\n"
            "  Position-level MCC may be too strict.\n"
            "  Methods may detect the right region\n"
            "  but miss exact positions.\n\n"
            "Results:\n"
            "  MCC improves from 0.289 to 0.712\n"
            "  with region-level evaluation.",
            size=12, color=DARK_GRAY)
add_placeholder(s, Inches(7.100), Inches(4.200), Inches(5.200), Inches(2.100),
                "[Image Placeholder]\nRegion Overlap Analysis Figure\n(To be inserted)")


# ── S31: Baseline Comparison ───────────────────────────────────
s = content_slide(prs, "Scoring-Detection Heatmap & Baseline Comparison", 31)

add_figure(s, os.path.join(FIG_MAIN, "scoring_detection_heatmap.png"),
           Inches(0.300), Inches(1.100), Inches(7.800), Inches(5.600))

add_textbox(s, Inches(8.500), Inches(1.200), Inches(4.2), Inches(0.4),
            "MCC Comparison", size=16, bold=True, color=SEC_EXT)
bars2 = [
    ("MutBench Best", "0.390", GREEN),
    ("FreqThresh", "0.248", BLUE),
    ("SWAN", "0.193", ORANGE),
    ("MutClust-Orig", "0.138", RED),
    ("Random", "0.001", GRAY),
]
for i, (name, val, clr) in enumerate(bars2):
    cy = Inches(1.800) + i * Inches(1.000)
    add_textbox(s, Inches(8.500), cy, Inches(2.5), Inches(0.35),
                name, size=12, bold=True, color=clr)
    bar_w = max(min(float(val) / 0.4 * 3.0, 3.000), 0.200)
    add_rect(s, Inches(8.500), cy + Inches(0.380), Inches(bar_w), Inches(0.350), clr)
    add_textbox(s, Inches(8.600), cy + Inches(0.380), Inches(bar_w - 0.200), Inches(0.350),
                val, size=11, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
                align=PP_ALIGN.RIGHT)


# ── S32: Research Contributions ─────────────────────────────────
s = content_slide(prs, "Research Contributions", 32)

contribs = [
    ("C1: MutBench Framework",
     "First systematic benchmark for viral\nmutation hotspot detection with\n"
     "standardized pipeline & evaluation", "See S11-S17", TEAL),
    ("C2: 3-Layer Ground Truth",
     "Adaptive threshold + constrained +\nDMS-based ground truth design\n"
     "reduces single-reference bias", "See S13-S15", GREEN),
    ("C3: Cross-Pathogen Evidence",
     "9 pathogens, 2,544 evaluations prove\n"
     "no universal best method exists\n"
     "(interaction omega-sq = 0.285)", "See S22-S27", SEC_S2),
    ("C4: PAHD Proof of Concept",
     "Pathogen-Adaptive Hotspot Detection:\n"
     "profile-based method selection\n"
     "outperforms any single approach", "See S28", PURPLE),
]
for i, (ttl, desc, ref, clr) in enumerate(contribs):
    col, row = i % 2, i // 2
    cx = Inches(0.600) + col * Inches(6.300)
    cy = Inches(1.200) + row * Inches(2.700)
    add_card(s, cx, cy, Inches(5.800), Inches(2.400), border_color=MID_GRAY)
    add_accent_bar(s, cx, cy, Inches(0.080), Inches(2.400), clr)
    add_textbox(s, cx + Inches(0.250), cy + Inches(0.150), Inches(5.3), Inches(0.4),
                ttl, size=16, bold=True, color=clr)
    add_textbox(s, cx + Inches(0.250), cy + Inches(0.650), Inches(5.3), Inches(1.3),
                desc, size=12, color=DARK_GRAY)
    add_textbox(s, cx + Inches(0.250), cy + Inches(1.950), Inches(5.3), Inches(0.3),
                ref, size=9, color=GRAY)


# ── S33: Discussion ────────────────────────────────────────────
s = content_slide(prs, "Discussion", 33)

disc = [
    ("Benchmarking Matters",
     "Without standardized evaluation, published methods appear to perform well in isolation "
     "but fail under cross-pathogen comparison. MutBench provides the missing infrastructure "
     "for fair, reproducible comparison of hotspot detection methods.", TEAL),
    ("Method Selection is Pathogen-Dependent",
     "The strong interaction effect (omega-sq = 0.285) means that recommending a single method "
     "is fundamentally flawed. Pathogen-specific characteristics (genome structure, mutation rate, "
     "selection landscape) determine which approach works best.", SEC_S2),
    ("From Benchmark to Algorithm",
     "MutBench is not just a benchmark but a knowledge base. The systematic evaluation across "
     "9 pathogens enables PAHD: data-driven, profile-based method selection that outperforms "
     "any fixed approach.", PURPLE),
]
for i, (ttl, desc, clr) in enumerate(disc):
    cy = Inches(1.200) + i * Inches(1.850)
    add_card(s, Inches(0.600), cy, Inches(12.100), Inches(1.650), border_color=MID_GRAY)
    add_accent_bar(s, Inches(0.600), cy, Inches(0.080), Inches(1.650), clr)
    add_textbox(s, Inches(0.900), cy + Inches(0.120), Inches(11.5), Inches(0.4),
                ttl, size=16, bold=True, color=clr)
    add_textbox(s, Inches(0.900), cy + Inches(0.580), Inches(11.5), Inches(0.95),
                desc, size=12, color=DARK_GRAY)


# ── S34: Limitations & Future Work ──────────────────────────────
s = content_slide(prs, "Limitations & Future Work", 34)

col_data = [
    ("Limitations", RED, [
        "Ground truth approximation\n(no perfect gold standard exists)",
        "Phylogenetic non-independence\nnot fully corrected (TreeTime partial)",
        "Limited to substitution mutations\n(no indels, recombination)",
        "PAHD is proof-of-concept only,\nnot production-ready system",
    ]),
    ("Future Work", BLUE, [
        "Full PAHD implementation with\nlearned pathogen profiles",
        "Integrate phylogenetic correction\n(TreeTime branch weighting)",
        "Extend to indels, recombination,\nand structural variants",
        "Real-time surveillance integration\nfor emerging pathogens",
    ]),
]
for ci, (ttl, clr, items) in enumerate(col_data):
    cx = Inches(0.500) + ci * Inches(6.400)
    add_rect(s, cx, Inches(1.200), Inches(6.000), Inches(0.500), clr)
    add_textbox(s, cx, Inches(1.210), Inches(6.0), Inches(0.480),
                ttl, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    for j, item in enumerate(items):
        iy = Inches(1.850) + j * Inches(1.250)
        add_card(s, cx, iy, Inches(6.000), Inches(1.050), border_color=MID_GRAY)
        add_accent_bar(s, cx, iy, Inches(0.060), Inches(1.050), clr)
        add_textbox(s, cx + Inches(0.200), iy + Inches(0.100), Inches(5.6), Inches(0.85),
                    item, size=11, color=DARK_GRAY)


# ── S35: Key Takeaway ──────────────────────────────────────────
s = content_slide(prs, "Key Takeaway", 35)

central = add_card(s, Inches(1.200), Inches(1.300), Inches(10.900), Inches(1.800), NAVY)
set_text(central,
         "MutBench transforms viral hotspot detection from\n"
         "ad-hoc single-virus evaluation to systematic cross-pathogen benchmarking",
         size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

evidences = [
    ("9 Pathogens", "Largest cross-pathogen\nbenchmark for hotspot\ndetection", TEAL),
    ("2,544 Evaluations", "Comprehensive coverage\nof scoring x detection\ncombinations", SEC_S2),
    ("omega-sq = 0.285", "Strong interaction proves\npathogen-adaptive\nselection is needed", PURPLE),
    ("PAHD Concept", "From benchmark insight\nto actionable algorithm\nfor method selection", GREEN),
]
for i, (ttl, desc, clr) in enumerate(evidences):
    cx = Inches(0.500) + i * Inches(3.200)
    cy = Inches(3.500)
    add_card(s, cx, cy, Inches(2.900), Inches(3.100), border_color=MID_GRAY)
    add_accent_bar(s, cx, cy, Inches(0.080), Inches(3.100), clr)
    add_textbox(s, cx, cy + Inches(0.250), Inches(2.9), Inches(0.5),
                ttl, size=16, bold=True, color=clr, align=PP_ALIGN.CENTER)
    add_textbox(s, cx + Inches(0.200), cy + Inches(0.900), Inches(2.5), Inches(2.0),
                desc, size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)


# ── S36: Thank You ─────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(s, Inches(0), Inches(0), SLIDE_W, Inches(3.500), NAVY)
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(0.600), Inches(0.300), Inches(0.650), Inches(0.650))
add_textbox(s, Inches(0.600), Inches(1.200), Inches(12.0), Inches(1.2),
            "Thank You", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(0.600), Inches(2.400), Inches(12.0), Inches(0.6),
            "Questions & Discussion", size=20, color=RGBColor(0xBB, 0xDE, 0xFB),
            align=PP_ALIGN.CENTER)

add_textbox(s, Inches(3.000), Inches(4.200), Inches(7.3), Inches(0.4),
            "Hwijun Kwon", size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(3.000), Inches(4.700), Inches(7.3), Inches(0.35),
            "Dept. of Computer Science, Kyungpook National University",
            size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(3.000), Inches(5.200), Inches(7.3), Inches(0.35),
            "Advisor: Prof. Inuk Jung", size=14, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(3.000), Inches(5.700), Inches(7.3), Inches(0.35),
            "March 2026", size=14, color=GRAY, align=PP_ALIGN.CENTER)

add_placeholder(s, Inches(5.000), Inches(6.200), Inches(3.300), Inches(0.600),
                "[QR Code Placeholder] GitHub Repository")

add_rect(s, Inches(0), Inches(7.120), SLIDE_W, Inches(0.380), NAVY)
add_textbox(s, Inches(0.400), Inches(7.140), Inches(10.0), Inches(0.340),
            "Kyungpook National University  |  Dept. of Computer Science  |  Hwijun Kwon",
            size=10, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)


# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Total slides: {len(prs.slides)}")
