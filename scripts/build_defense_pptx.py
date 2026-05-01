from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path("paper/dissertation/presentation/mutbench_defense_pahd_r.pptx")
FIG_DIR = Path("paper/dissertation/figures")
GUIDE_FIG_DIR = Path("paper/dissertation/guide_figures")

NAVY = RGBColor(10, 28, 52)
BLUE = RGBColor(28, 100, 170)
CYAN = RGBColor(42, 157, 244)
RED = RGBColor(179, 36, 36)
GREEN = RGBColor(38, 130, 73)
ORANGE = RGBColor(201, 112, 20)
PURPLE = RGBColor(94, 84, 166)
INK = RGBColor(24, 31, 42)
MUTED = RGBColor(91, 104, 124)
LINE = RGBColor(218, 225, 234)
PAPER = RGBColor(246, 248, 251)
SOFT = RGBColor(237, 243, 249)
WHITE = RGBColor(255, 255, 255)

W, H = 13.333, 7.5


def rgb(hex_code):
    hex_code = hex_code.strip("#")
    return RGBColor(int(hex_code[:2], 16), int(hex_code[2:4], 16), int(hex_code[4:], 16))


def font(run, size=18, bold=False, color=INK):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def set_bg(slide, color=PAPER):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.color.rgb = fill
        shape.line.width = Pt(0)
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.8)
    return shape


def text(slide, value, x, y, w, h, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    font(run, size=size, bold=bold, color=color)
    return box


def header(slide, title, kicker, idx):
    set_bg(slide)
    rect(slide, 0, 0, W, 0.68, NAVY)
    rect(slide, 0, 0.68, W, 0.03, CYAN)
    text(slide, kicker.upper(), 0.55, 0.19, 2.1, 0.22, 7.5, True, CYAN)
    text(slide, title, 1.75, 0.12, 9.9, 0.42, 17.5, True, WHITE)
    text(slide, f"{idx:02d}", 12.15, 0.19, 0.65, 0.22, 8, True, WHITE, PP_ALIGN.RIGHT)
    rect(slide, 0.55, 7.08, 12.2, 0.02, LINE)
    text(slide, "MutBench / PAHD-R defense deck", 0.58, 7.17, 4.0, 0.18, 7.5, False, MUTED)


def bullets(slide, items, x, y, w, h, size=15.5, color=INK, gap=5):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(gap)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def card(slide, title, body, x, y, w, h, accent=BLUE, title_size=13, body_size=11.5):
    rect(slide, x, y, w, h, WHITE, LINE)
    rect(slide, x, y, 0.09, h, accent)
    text(slide, title, x + 0.22, y + 0.15, w - 0.35, 0.25, title_size, True, accent)
    bullets(slide, body if isinstance(body, list) else [body], x + 0.22, y + 0.52, w - 0.35, h - 0.63, body_size, INK, 3)


def metric(slide, label, value, note, x, y, w, h, accent=BLUE):
    rect(slide, x, y, w, h, WHITE, LINE)
    rect(slide, x, y, w, 0.08, accent)
    text(slide, label, x + 0.14, y + 0.18, w - 0.28, 0.24, 9, True, MUTED)
    text(slide, value, x + 0.14, y + 0.46, w - 0.28, 0.42, 21, True, accent)
    text(slide, note, x + 0.14, y + 0.96, w - 0.28, 0.35, 8.5, False, MUTED)


def table(slide, rows, x, y, w, h, font_size=10.5, header_color=NAVY):
    shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = shape.table
    for c in range(len(rows[0])):
        tbl.columns[c].width = Inches(w / len(rows[0]))
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = str(val)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_color if r == 0 else (WHITE if r % 2 else rgb("F2F6FA"))
            for p in cell.text_frame.paragraphs:
                p.font.name = "Aptos"
                p.font.size = Pt(font_size)
                p.font.bold = r == 0
                p.font.color.rgb = WHITE if r == 0 else INK
                p.alignment = PP_ALIGN.CENTER if c > 1 else PP_ALIGN.LEFT
    return tbl


def picture(slide, path, x, y, max_w, max_h):
    path = Path(path)
    if not path.exists():
        return None
    with Image.open(path) as img:
        iw, ih = img.size
    scale = min(max_w / iw, max_h / ih)
    w = iw * scale
    h = ih * scale
    return slide.shapes.add_picture(str(path), Inches(x + (max_w - w) / 2), Inches(y + (max_h - h) / 2), width=Inches(w), height=Inches(h))


def slide(prs, title, kicker, idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    header(s, title, kicker, idx)
    return s


def section_slide(prs, idx, title, subtitle, tags):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, NAVY)
    rect(s, 0, 0, 0.18, H, CYAN)
    text(s, f"{idx:02d}", 0.65, 0.85, 0.7, 0.35, 13, True, CYAN)
    text(s, title, 0.65, 1.45, 10.8, 0.8, 34, True, WHITE)
    text(s, subtitle, 0.68, 2.45, 10.6, 0.55, 17, False, rgb("D7E7F7"))
    x = 0.68
    for tag in tags:
        rect(s, x, 5.65, 2.3, 0.38, rgb("183A5C"), CYAN)
        text(s, tag, x + 0.15, 5.75, 2.0, 0.16, 8.5, True, WHITE, PP_ALIGN.CENTER)
        x += 2.55
    return s


def build():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    prs.core_properties.title = "MutBench PAHD-R defense deck"
    prs.core_properties.subject = "Designed PowerPoint deck for dissertation defense"
    prs.core_properties.author = "Hwijun Kwon"

    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, NAVY)
    rect(s, 0, 0, W, H, NAVY)
    rect(s, 0, 5.95, W, 1.55, rgb("08182D"))
    rect(s, 0.65, 0.72, 0.12, 4.95, CYAN)
    text(s, "MUTBENCH", 0.95, 0.95, 5.7, 0.45, 14, True, CYAN)
    text(s, "Pathogen-dependent hotspot detection", 0.95, 1.48, 9.8, 1.25, 36, True, WHITE)
    text(s, "PAHD-R confirms adaptive evidence-fusion feasibility under audit", 0.98, 3.02, 8.6, 0.48, 19, False, rgb("D7E7F7"))
    for i, item in enumerate(["11 viruses", "20 scoring types", "8,580 evaluations"]):
        rect(s, 0.98 + i * 2.45, 4.45, 2.1, 0.46, rgb("183A5C"), CYAN)
        text(s, item, 1.08 + i * 2.45, 4.58, 1.9, 0.16, 9, True, WHITE, PP_ALIGN.CENTER)
    text(s, "Hwijun Kwon · Doctoral Dissertation Defense · 2026", 0.95, 6.45, 9.5, 0.25, 10.5, False, rgb("BFD3E7"))

    s = slide(prs, "Problem Framing", "motivation", 2)
    card(s, "What is missing?", ["No standardized biological benchmark for viral mutation hotspots.", "Single-virus results transfer poorly across pathogens."], 0.75, 1.15, 3.6, 2.05, CYAN)
    card(s, "What is usually fixed?", ["Many studies compare detectors while holding the information source fixed.", "That hides information-pathogen fit."], 4.85, 1.15, 3.6, 2.05, ORANGE)
    card(s, "Operational target", ["Rank positions for DMS, epitope mapping, and surveillance review.", "Report useful prioritization without overclaiming universality."], 8.95, 1.15, 3.6, 2.05, GREEN)
    rect(s, 0.78, 4.25, 11.75, 1.25, WHITE, LINE)
    text(s, "Thesis claim", 1.05, 4.52, 1.5, 0.22, 10, True, MUTED)
    text(s, "The best biological information source is pathogen-dependent; PAHD-R converts that finding into a bounded review workflow.", 2.45, 4.38, 9.4, 0.48, 20, True, NAVY)

    s = slide(prs, "MutBench Design", "benchmark", 3)
    metric(s, "Pathogens", "11", "RNA viruses / surface proteins", 0.75, 1.1, 2.2, 1.45, CYAN)
    metric(s, "Scoring types", "20", "frequency, MSA, phylo, structure, PLM, composite", 3.2, 1.1, 2.6, 1.45, ORANGE)
    metric(s, "Detector variants", "39", "14 families", 6.05, 1.1, 2.3, 1.45, GREEN)
    metric(s, "Evaluations", "8,580", "standardized comparisons", 8.6, 1.1, 2.5, 1.45, PURPLE)
    card(s, "Ground truth layers", ["Layer A: adaptive positives.", "Layer B: constrained negatives.", "Layer C: DMS/external validation when available."], 0.75, 3.15, 4.05, 2.25, BLUE)
    picture(s, GUIDE_FIG_DIR / "fig_multi_gt.png", 5.15, 3.0, 6.85, 2.75)

    s = slide(prs, "Information × Pathogen Interaction", "main result 1", 4)
    picture(s, FIG_DIR / "stage3_anova_decomposition.png", 0.7, 1.08, 5.7, 4.8)
    table(
        s,
        [
            ["Factor", "Meaning", "Omega²"],
            ["Scoring × pathogen", "information fit by virus", "0.296"],
            ["Pathogen", "baseline difficulty", "0.117"],
            ["Family × pathogen", "detector behavior", "0.082"],
            ["Scoring", "global information effect", "0.063"],
            ["Family", "global detector effect", "0.013"],
        ],
        6.65,
        1.22,
        5.85,
        2.7,
        10.2,
    )
    card(s, "Interpretation", "Detector choice matters, but the dominant modeled factor is whether the biological information source matches the pathogen.", 6.65, 4.25, 5.85, 1.25, CYAN, 12.5, 12.8)

    s = slide(prs, "No Universal Best Method", "main result 2", 5)
    picture(s, FIG_DIR / "stage3_scoring_pathogen_heatmap.png", 0.7, 1.05, 6.05, 4.95)
    card(s, "Observed pattern", ["11/11 pathogens have unique best scoring-detector combinations.", "Best combinations span 9 scoring types.", "Friedman p = 0.990 for top combinations."], 7.05, 1.15, 5.2, 2.1, ORANGE)
    card(s, "Implication", ["LOPO selection fails to identify a transferable best method.", "This is evidence of non-transfer, not just weak model selection."], 7.05, 3.55, 5.2, 1.75, BLUE)

    s = slide(prs, "External Validation", "evidence check", 6)
    table(
        s,
        [
            ["Pathogen", "Role", "Enrichment", "Interpretation"],
            ["HIV-1", "primary external anchor", "7.19x", "82% outside Layer A"],
            ["H3N2", "self-consistency", "9.36x", "69% Layer A overlap"],
            ["SARS-CoV-2", "exploratory", "7.63x", "Bonferroni not passed"],
            ["EqualWeight H3N2", "search reduction", "4.012x", "top-5% enrichment"],
        ],
        0.75,
        1.05,
        11.85,
        2.35,
        10.3,
    )
    picture(s, GUIDE_FIG_DIR / "search_space_reduction.png", 0.85, 3.85, 5.6, 2.35)
    card(s, "Conservative reading", "HIV-1 is the cleanest external validation anchor because most escape positions are disjoint from the benchmark positive labels.", 6.95, 4.0, 5.25, 1.55, GREEN, 12, 13)

    s = slide(prs, "Compact 4-Feature Core", "deployment floor", 7)
    picture(s, FIG_DIR / "feature_ablation_bars.png", 0.75, 1.05, 5.7, 4.95)
    table(
        s,
        [
            ["Feature set", "Mean MCC", "Operational meaning"],
            ["Homoplasy", "0.062", "independent recurrence"],
            ["+ pLDDT", "0.072", "structural flexibility"],
            ["+ entropy", "0.077", "allelic diversity"],
            ["+ frequency", "0.081", "4-feature CPU core"],
            ["Full 10-feature", "0.083", "GPU/PLM-expanded ensemble"],
        ],
        6.7,
        1.15,
        5.7,
        2.85,
        10.2,
    )
    card(s, "Decision", "The 4-feature core captures about 98% of the full ensemble and is retained as the cold-start floor.", 6.7, 4.45, 5.7, 1.15, CYAN, 12, 13)

    s = slide(prs, "Adaptive Weighting: Why Not Default?", "negative result", 8)
    card(s, "Rejected as default", ["kNN-weighted.", "Top-3 AUC.", "Correlation-aware weighting."], 0.8, 1.15, 3.45, 2.0, RED)
    card(s, "Failure mode", ["Only 11 reference pathogens.", "13-dimensional pathogen-profile space is undersampled.", "Model selection risk outweighs current gain."], 4.75, 1.15, 3.85, 2.0, ORANGE)
    card(s, "Still useful", ["Oracle result shows adaptive potential.", "Candidate variants confirm plausible design routes.", "Future adoption needs more pathogens or external labels."], 9.1, 1.15, 3.25, 2.0, BLUE)
    rect(s, 0.95, 4.35, 11.35, 0.92, WHITE, LINE)
    text(s, "Bottom line", 1.2, 4.62, 1.55, 0.22, 10, True, MUTED)
    text(s, "Current evidence supports bounded evidence fusion, not a learned universal adaptive predictor.", 2.55, 4.49, 8.9, 0.35, 18.5, True, NAVY)

    s = slide(prs, "PAHD-R: Possibility Under Audit", "algorithm", 9)
    table(
        s,
        [
            ["Mode", "Role", "Exact", "+/-10", "P@20", "FPR", "Status"],
            ["Core", "no-AI exact/top-k baseline", "0.1689", "0.5866", "0.2111", "0.1108", "adopted"],
            ["Augmented", "balanced region prioritization", "0.1680", "0.6079", "0.1944", "0.1128", "adopted"],
            ["Review", "conservative low-FPR review", "0.1666", "0.6366", "0.1833", "0.0696", "adopted"],
        ],
        0.55,
        1.18,
        12.2,
        2.25,
        9.8,
    )
    card(s, "Definition boundary", "PAHD-R confirms the feasibility of adaptive evidence fusion under audit constraints. It is not a completed universal adaptive predictor.", 0.75, 4.05, 5.75, 1.45, BLUE, 12, 13.5)
    card(s, "Reporting rule", "Do not report +/-10 MCC alone. Pair it with exact MCC, P@20, constrained FPR, and coverage.", 6.85, 4.05, 5.5, 1.45, GREEN, 12, 13.5)

    s = slide(prs, "Candidate Variants Kept for Future Validation", "algorithm audit", 10)
    table(
        s,
        [
            ["Candidate", "Improves", "Why not default"],
            ["Core-Calibrated", "Exact 0.1770, FPR 0.0947", "per-pathogen null weaknesses"],
            ["ShrinkCompact", "FPR 0.0782", "coverage trade-off"],
            ["ShrinkRegion", "+/-10 0.6232, FPR 0.0782", "weak decoy/window support"],
            ["Virus-aware hybrid", "Exact 0.1765, FPR 0.0940", "near-tie, model-selection risk"],
            ["Selective callable-only", "Exact 0.2742, FPR 0.0354", "accepts only 4/9 pathogens"],
        ],
        0.55,
        1.05,
        12.2,
        3.05,
        9.5,
    )
    card(s, "Final decision", "Adopt Core / Augmented / Review. Keep calibrated, shrinkage, virus-aware, and selective variants as possibility-confirming future candidates.", 1.0, 4.72, 11.2, 0.95, ORANGE, 12, 13.5)

    s = slide(prs, "Repair-Layer Decisions", "adversarial boundary", 11)
    table(
        s,
        [
            ["Repair", "Status", "Reason"],
            ["HCV E2", "adopted optional", "coordinate reconciliation; exact 0.5257, FPR 0.0000"],
            ["SARS-CoV-2 RBD/ACE2", "candidate", "exact 0.2304, +/-10 0.8761, P@20 0.1000, FPR 0.0065"],
            ["MERS DPP4/contact", "candidate", "sparse positives; structural-coordinate review needed"],
            ["RSV", "unrepaired", "repair ideas fail required controls"],
        ],
        0.75,
        1.15,
        11.85,
        2.65,
        10.2,
    )
    card(s, "Gatekeeping rule", "Repair layers are allowed only when coordinate validity, biological prior, decoy controls, and audits support them.", 1.05, 4.45, 5.25, 1.15, RED, 12, 13)
    card(s, "Practical consequence", "SARS-CoV-2 and MERS remain redesign targets, not final adopted repairs.", 6.85, 4.45, 5.05, 1.15, ORANGE, 12, 13)

    s = slide(prs, "Claim Boundaries", "final audit", 12)
    bullets(
        s,
        [
            "PAHD-R should not be described as a universal AI predictor.",
            "SARS-CoV-2 and MERS repair results are candidate-only.",
            "Selective callability is not all-pathogen performance; it accepts only 4/9 pathogens.",
            "Pooled permutation supports global sanity checking, not a full pathogen-level null for every virus.",
            "Further adoption requires external, time-forward, or tree/bootstrap validation.",
        ],
        0.9,
        1.12,
        6.0,
        4.6,
        15.3,
    )
    card(s, "Audit stance", "Prefer underclaiming: report exact-site utility, neighborhood utility, precision, FPR, and coverage together.", 7.35, 1.42, 4.95, 1.4, GREEN, 12, 13.5)
    card(s, "Risk still open", "Additional pathogens and time-forward labels are needed before moving from possibility to adoption.", 7.35, 3.32, 4.95, 1.4, BLUE, 12, 13.5)

    s = slide(prs, "Practical Workflow", "how to use", 13)
    steps = [
        ("1", "Collect", ">=200 sequences and build MSA."),
        ("2", "Compute core", "homoplasy, pLDDT, entropy, frequency."),
        ("3", "Prioritize", "use family-level recommendations as priors."),
        ("4", "Review", "run Core / Augmented / Review together."),
        ("5", "Validate", "DMS, epitope mapping, or surveillance follow-up."),
    ]
    for i, (num, title, body) in enumerate(steps):
        x = 0.75 + i * 2.48
        rect(s, x, 1.45, 1.68, 1.68, WHITE, LINE)
        text(s, num, x + 0.53, 1.75, 0.6, 0.38, 22, True, CYAN, PP_ALIGN.CENTER)
        text(s, title, x - 0.15, 3.42, 2.0, 0.28, 12, True, NAVY, PP_ALIGN.CENTER)
        text(s, body, x - 0.18, 3.82, 2.05, 0.58, 9.2, False, MUTED, PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            rect(s, x + 1.78, 2.26, 0.44, 0.06, CYAN)
    card(s, "Use rule", "PAHD-R is a prioritization and review tool. It narrows search space; it does not replace experimental validation.", 1.0, 5.1, 11.1, 0.82, ORANGE, 11.5, 12.5)

    s = section_slide(
        prs,
        14,
        "Final Take-Home",
        "Pathogen-specific information fit is the core discovery; PAHD-R confirms a bounded adaptive direction under audit.",
        ["Adopted: Core / Augmented / Review", "Candidate: calibrated / shrinkage / virus-aware", "Optional repair: HCV E2"],
    )
    rect(s, 0.68, 4.0, 10.9, 0.04, CYAN)
    text(s, "Conservative final wording: adaptive possibility confirmed, explicit audit boundaries retained.", 0.7, 4.38, 10.6, 0.42, 15.5, False, rgb("D7E7F7"))

    for shape in prs.slides[0].shapes:
        if hasattr(shape, "text_frame"):
            shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
